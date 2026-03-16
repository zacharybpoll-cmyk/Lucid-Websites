#!/usr/bin/env python3
"""Generate Lucid Persona Analysis PowerPoint deck."""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

# ─── Brand Colors ───
STEEL_BLUE = RGBColor(0x5B, 0x8D, 0xB8)
DARK_TEXT = RGBColor(0x1A, 0x1D, 0x21)
BODY_TEXT = RGBColor(0x5A, 0x62, 0x70)
BG_COLOR = RGBColor(0xF8, 0xF9, 0xFA)
MID_GRAY = RGBColor(0xF0, 0xF2, 0xF4)
DETAIL_GRAY = RGBColor(0xE4, 0xE8, 0xEC)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
DARK_BG = RGBColor(0x1A, 0x3A, 0x4F)
ACCENT_LIGHT = RGBColor(0xD4, 0xE4, 0xF0)

# Persona accent colors
PERSONA_COLORS = {
    "tech": RGBColor(0x5B, 0x8D, 0xB8),       # Steel blue
    "longevity": RGBColor(0x6B, 0xA3, 0x7E),   # Sage green
    "skeptic": RGBColor(0x9B, 0x7E, 0xB8),     # Muted purple
    "athlete": RGBColor(0xC4, 0x7A, 0x5A),      # Warm terracotta
    "disconnected": RGBColor(0x57, 0xF2, 0x87),   # Electric green
}

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)


def set_slide_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_shape(slide, left, top, width, height, fill_color, corner_radius=None):
    if corner_radius:
        shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
        # Adjust corner radius
        shape.adjustments[0] = corner_radius
    else:
        shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    return shape


def add_text_box(slide, left, top, width, height, text, font_size=14, bold=False,
                 color=BODY_TEXT, font_name="Inter", alignment=PP_ALIGN.LEFT,
                 italic=False, anchor=MSO_ANCHOR.TOP):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    tf.auto_size = None
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = font_name
    p.alignment = alignment
    if italic:
        p.font.italic = True
    return txBox


def add_multi_text(slide, left, top, width, height, lines, anchor=MSO_ANCHOR.TOP):
    """Add textbox with multiple formatted lines. Each line is (text, size, bold, color, font, alignment, italic, space_after)"""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    tf.auto_size = None

    for i, line_data in enumerate(lines):
        text = line_data[0]
        size = line_data[1] if len(line_data) > 1 else 14
        bold = line_data[2] if len(line_data) > 2 else False
        color = line_data[3] if len(line_data) > 3 else BODY_TEXT
        font = line_data[4] if len(line_data) > 4 else "Inter"
        align = line_data[5] if len(line_data) > 5 else PP_ALIGN.LEFT
        italic = line_data[6] if len(line_data) > 6 else False
        space_after = line_data[7] if len(line_data) > 7 else Pt(6)

        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()

        p.text = text
        p.font.size = Pt(size)
        p.font.bold = bold
        p.font.color.rgb = color
        p.font.name = font
        p.alignment = align
        p.space_after = space_after
        if italic:
            p.font.italic = True

    return txBox


def add_accent_bar(slide, left, top, width, height, color):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape


def add_card(slide, left, top, width, height, title, items, accent_color=STEEL_BLUE, title_size=16):
    """Add a card with title and bullet items."""
    card = add_shape(slide, left, top, width, height, WHITE, corner_radius=0.02)

    # Accent bar at top
    add_accent_bar(slide, left, top, width, Pt(4), accent_color)

    # Title
    add_text_box(slide, left + Inches(0.3), top + Inches(0.2), width - Inches(0.6), Inches(0.5),
                 title, font_size=title_size, bold=True, color=DARK_TEXT, font_name="Playfair Display")

    # Items
    lines = []
    for item in items:
        lines.append((item, 12, False, BODY_TEXT, "Inter", PP_ALIGN.LEFT, False, Pt(8)))

    add_multi_text(slide, left + Inches(0.3), top + Inches(0.65), width - Inches(0.6),
                   height - Inches(0.85), lines)

    return card


# ═══════════════════════════════════════════════════════════
# SLIDE 1: Title
# ═══════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank
set_slide_bg(slide, DARK_BG)

# Left accent bar
add_accent_bar(slide, Inches(0), Inches(0), Inches(0.15), SLIDE_H, STEEL_BLUE)

# Title
add_text_box(slide, Inches(1.2), Inches(1.8), Inches(8), Inches(1.2),
             "Lucid", font_size=60, bold=True, color=WHITE, font_name="Playfair Display")

add_text_box(slide, Inches(1.2), Inches(3.0), Inches(9), Inches(1.0),
             "User Persona Analysis", font_size=36, bold=False, color=ACCENT_LIGHT, font_name="Playfair Display")

add_text_box(slide, Inches(1.2), Inches(4.0), Inches(9), Inches(0.6),
             "Understanding who we serve and how to reach them.", font_size=18, color=RGBColor(0xA0, 0xB8, 0xCC), font_name="Inter", italic=True)

# Horizontal divider
add_accent_bar(slide, Inches(1.2), Inches(5.0), Inches(2), Pt(2), STEEL_BLUE)

add_text_box(slide, Inches(1.2), Inches(5.3), Inches(6), Inches(0.5),
             "March 2026  |  Confidential", font_size=14, color=RGBColor(0x80, 0x99, 0xAA), font_name="Inter")

# Right side stats
stats = [
    ("5", "Core Personas"),
    ("20+", "Vocal Biomarkers"),
    ("$14.99/mo", "Price Point"),
    ("100%", "On-Device Privacy"),
]
for i, (num, label) in enumerate(stats):
    y = Inches(1.8) + Inches(i * 1.2)
    add_text_box(slide, Inches(10.5), y, Inches(2.5), Inches(0.5),
                 num, font_size=32, bold=True, color=STEEL_BLUE, font_name="Playfair Display", alignment=PP_ALIGN.RIGHT)
    add_text_box(slide, Inches(10.5), y + Inches(0.45), Inches(2.5), Inches(0.4),
                 label, font_size=13, color=RGBColor(0x80, 0x99, 0xAA), font_name="Inter", alignment=PP_ALIGN.RIGHT)


# ═══════════════════════════════════════════════════════════
# SLIDE 2: Executive Summary
# ═══════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BG_COLOR)

add_text_box(slide, Inches(0.8), Inches(0.5), Inches(8), Inches(0.6),
             "Why Personas Matter", font_size=32, bold=True, color=DARK_TEXT, font_name="Playfair Display")

add_accent_bar(slide, Inches(0.8), Inches(1.15), Inches(1.5), Pt(3), STEEL_BLUE)

add_text_box(slide, Inches(0.8), Inches(1.5), Inches(7.5), Inches(1.2),
             "Lucid is a passive voice wellness monitor. But different users come for different reasons, stay for different features, and respond to different messaging. Understanding these five personas lets us craft targeted acquisition funnels, prioritize features by impact, and speak each user's language.",
             font_size=15, color=BODY_TEXT)

# Four persona preview cards
personas_preview = [
    ("The Burned-Out Builder", "Tech workers drowning in\nmeetings and screen time", PERSONA_COLORS["tech"], "34% of the TAM"),
    ("The Optimization Architect", "Longevity-focused biohackers\nwho track everything", PERSONA_COLORS["longevity"], "28% of the TAM"),
    ("The System Skeptic", "Privacy-first individuals who\ndon't trust the healthcare system", PERSONA_COLORS["skeptic"], "22% of the TAM"),
    ("The Performance Athlete", "Weekend warriors & competitors\nwhose mental edge defines results", PERSONA_COLORS["athlete"], "16% of the TAM"),
]

for i, (name, desc, color, tam) in enumerate(personas_preview):
    x = Inches(0.8) + Inches(i * 3.1)
    y = Inches(3.0)
    w = Inches(2.9)
    h = Inches(3.8)

    card = add_shape(slide, x, y, w, h, WHITE, corner_radius=0.02)
    add_accent_bar(slide, x, y, w, Pt(5), color)

    # Persona number circle
    circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, x + Inches(0.2), y + Inches(0.35), Inches(0.5), Inches(0.5))
    circle.fill.solid()
    circle.fill.fore_color.rgb = color
    circle.line.fill.background()
    tf = circle.text_frame
    tf.word_wrap = False
    p = tf.paragraphs[0]
    p.text = str(i + 1)
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.font.name = "Inter"
    p.alignment = PP_ALIGN.CENTER

    add_text_box(slide, x + Inches(0.85), y + Inches(0.35), w - Inches(1.1), Inches(0.5),
                 name, font_size=15, bold=True, color=DARK_TEXT, font_name="Playfair Display")

    add_text_box(slide, x + Inches(0.2), y + Inches(1.1), w - Inches(0.4), Inches(1.2),
                 desc, font_size=12, color=BODY_TEXT)

    # TAM estimate
    add_shape(slide, x + Inches(0.2), y + Inches(2.8), w - Inches(0.4), Inches(0.4), MID_GRAY, corner_radius=0.05)
    add_text_box(slide, x + Inches(0.2), y + Inches(2.85), w - Inches(0.4), Inches(0.3),
                 tam, font_size=11, bold=True, color=color, alignment=PP_ALIGN.CENTER)


# ═══════════════════════════════════════════════════════════
# SLIDE 3: Persona 1 — The Burned-Out Builder (Profile)
# ═══════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BG_COLOR)

# Header bar
add_shape(slide, Inches(0), Inches(0), SLIDE_W, Inches(1.2), DARK_BG)
add_accent_bar(slide, Inches(0), Inches(1.2), SLIDE_W, Pt(4), PERSONA_COLORS["tech"])

add_text_box(slide, Inches(0.8), Inches(0.15), Inches(4), Inches(0.4),
             "PERSONA 1", font_size=12, bold=True, color=STEEL_BLUE, font_name="Inter")
add_text_box(slide, Inches(0.8), Inches(0.45), Inches(8), Inches(0.6),
             "The Burned-Out Builder", font_size=30, bold=True, color=WHITE, font_name="Playfair Display")

# Left column — Profile
add_text_box(slide, Inches(0.8), Inches(1.6), Inches(5.5), Inches(0.4),
             "Who They Are", font_size=20, bold=True, color=DARK_TEXT, font_name="Playfair Display")

profile_items = [
    ("Demographics", "25-40 years old  |  $80K-$200K income  |  Urban/suburban\nSoftware engineers, product managers, designers, tech founders"),
    ("Psychographics", "Achievement-oriented but increasingly aware of burnout cost.\nThey've seen colleagues flame out. They want data, not therapy.\nIdentity: 'I optimize everything — except my mental health.'"),
    ("Current Tools", "Slack, Zoom, Google Calendar (11+ meetings/week avg)\nOura/Apple Watch for physical tracking\nCalm/Headspace (tried, abandoned — 3.3% 30-day retention)"),
    ("Spending", "$5,321/year avg on wellness (Fortune/GWI 2024)\nWilling to pay $6.50+/month for health apps (58.9% WTP)"),
]

y_pos = Inches(2.1)
for label, desc in profile_items:
    add_text_box(slide, Inches(0.8), y_pos, Inches(1.8), Inches(0.3),
                 label, font_size=11, bold=True, color=STEEL_BLUE, font_name="Inter")
    add_text_box(slide, Inches(2.8), y_pos, Inches(4.0), Inches(1.0),
                 desc, font_size=11, color=BODY_TEXT)
    y_pos += Inches(1.2)

# Right column — Key Stats
stats_card = add_shape(slide, Inches(7.5), Inches(1.6), Inches(5.0), Inches(5.5), WHITE, corner_radius=0.02)
add_accent_bar(slide, Inches(7.5), Inches(1.6), Inches(5.0), Pt(4), PERSONA_COLORS["tech"])

add_text_box(slide, Inches(7.8), Inches(1.8), Inches(4.5), Inches(0.4),
             "The Pain by the Numbers", font_size=18, bold=True, color=DARK_TEXT, font_name="Playfair Display")

pain_stats = [
    ("82%", "of tech employees feel close to burnout"),
    ("52%", "of tech workers report depression/anxiety"),
    ("65%", "of engineers experienced burnout in past year"),
    ("91%", "of 18-24 year-olds report fatigue and brain fog"),
    ("250+", "digital actions/day outside office hours = stress spike"),
    ("11+", "meetings/week significantly increases employee stress"),
    ("$300B", "annual US cost of stress-related productivity loss"),
]

y_pos = Inches(2.35)
for stat, desc in pain_stats:
    add_text_box(slide, Inches(7.8), y_pos, Inches(1.2), Inches(0.35),
                 stat, font_size=18, bold=True, color=PERSONA_COLORS["tech"], font_name="Inter")
    add_text_box(slide, Inches(9.1), y_pos + Inches(0.03), Inches(3.2), Inches(0.35),
                 desc, font_size=11, color=BODY_TEXT)
    y_pos += Inches(0.55)


# ═══════════════════════════════════════════════════════════
# SLIDE 4: Persona 1 — Pain Points & Messaging
# ═══════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BG_COLOR)

# Mini header
add_accent_bar(slide, Inches(0), Inches(0), Pt(6), SLIDE_H, PERSONA_COLORS["tech"])

add_text_box(slide, Inches(0.8), Inches(0.4), Inches(4), Inches(0.3),
             "PERSONA 1  |  THE BURNED-OUT BUILDER", font_size=11, bold=True, color=STEEL_BLUE, font_name="Inter")
add_text_box(slide, Inches(0.8), Inches(0.7), Inches(8), Inches(0.5),
             "Pain Points, Hooks & Feature Value", font_size=26, bold=True, color=DARK_TEXT, font_name="Playfair Display")

# Pain Points card
add_card(slide, Inches(0.8), Inches(1.5), Inches(5.8), Inches(2.8),
         "Core Pain Points", [
             "  'I don't realize I'm burned out until it's too late.'",
             "  23+ meetings/week but no data on which ones drain them.",
             "  Physical health is tracked — mental health is a black box.",
             "  Therapy feels too heavy; they want monitoring, not treatment.",
             "  Always-on culture: 250+ digital actions/day outside office hours.",
             "  Information overload accounts for 39% of exhaustion variance.",
         ], PERSONA_COLORS["tech"])

# Messaging That Resonates card
add_card(slide, Inches(7.0), Inches(1.5), Inches(5.8), Inches(2.8),
         "Messaging That Resonates", [
             "  'You Had 23 Meetings Last Week. Which Ones Are Burning You Out?'",
             "  'Know Your Stress Before You Feel It.'",
             "  'Your voice changes 30 minutes before you feel stressed. We catch it.'",
             "  'Most productivity tools tell you to do more. Lucid tells you the truth.'",
             "  'AI detects stress 6.7 minutes before you feel it.'",
             "  Data-driven language, not emotional/therapeutic language.",
         ], PERSONA_COLORS["tech"])

# Feature Value Ranking
add_text_box(slide, Inches(0.8), Inches(4.5), Inches(6), Inches(0.4),
             "Feature Value Ranking (Highest to Lowest)", font_size=16, bold=True, color=DARK_TEXT, font_name="Playfair Display")

features = [
    ("1. Meeting Stress Correlation", "Which calls spike stress — the killer feature for this persona", "10/10"),
    ("2. Stress & Calm Scores", "Real-time 0-100 score during the workday — passive, no effort", "9/10"),
    ("3. Morning Briefing", "Quick summary: 'Yesterday your stress peaked at 3pm after the all-hands'", "9/10"),
    ("4. Weekly Wrapped", "Shareable weekly recap — social proof spreads to teammates", "8/10"),
    ("5. Burnout Trajectory", "Trend line showing whether they're recovering or declining week-over-week", "8/10"),
    ("6. Echoes (Pattern Discovery)", "'Your stress is 40% higher after back-to-back meetings' — discovery delight", "7/10"),
]

y_pos = Inches(5.0)
for feat, desc, score in features:
    add_text_box(slide, Inches(0.8), y_pos, Inches(3.5), Inches(0.3),
                 feat, font_size=12, bold=True, color=DARK_TEXT)
    add_text_box(slide, Inches(4.5), y_pos, Inches(6.5), Inches(0.3),
                 desc, font_size=11, color=BODY_TEXT)
    add_text_box(slide, Inches(11.8), y_pos, Inches(1.0), Inches(0.3),
                 score, font_size=12, bold=True, color=PERSONA_COLORS["tech"], alignment=PP_ALIGN.RIGHT)
    y_pos += Inches(0.38)


# ═══════════════════════════════════════════════════════════
# SLIDE 5: Persona 1 — Acquisition Strategy
# ═══════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BG_COLOR)

add_accent_bar(slide, Inches(0), Inches(0), Pt(6), SLIDE_H, PERSONA_COLORS["tech"])

add_text_box(slide, Inches(0.8), Inches(0.4), Inches(4), Inches(0.3),
             "PERSONA 1  |  THE BURNED-OUT BUILDER", font_size=11, bold=True, color=STEEL_BLUE, font_name="Inter")
add_text_box(slide, Inches(0.8), Inches(0.7), Inches(8), Inches(0.5),
             "Acquisition & Conversion Strategy", font_size=26, bold=True, color=DARK_TEXT, font_name="Playfair Display")

# Channels card
add_card(slide, Inches(0.8), Inches(1.5), Inches(3.7), Inches(2.8),
         "Top Acquisition Channels", [
             "  1. LinkedIn (founder posts about burnout/productivity)",
             "  2. TikTok ('my meeting stress scores' screen recordings)",
             "  3. Reddit (r/cscareerquestions, r/ExperiencedDevs)",
             "  4. Twitter/X (tech Twitter loves data tools)",
             "  5. Hacker News / Product Hunt launch",
         ], PERSONA_COLORS["tech"])

# Hook Messages
add_card(slide, Inches(4.8), Inches(1.5), Inches(3.7), Inches(2.8),
         "Hook Messages That Convert", [
             "  'I tracked my meeting stress for 30 days. Here's what I found.'",
             "  'The app that told me my 1:1s were my best recovery.'",
             "  'You track sleep. You track steps. Why not your stress?'",
             "  'Your standup is costing you more than you think.'",
             "  Screen recordings of stress scores after different meetings.",
         ], PERSONA_COLORS["tech"])

# Conversion Strategy
add_card(slide, Inches(8.8), Inches(1.5), Inches(3.7), Inches(2.8),
         "Conversion Strategy", [
             "  Free tier: Basic daily wellness score (hook them)",
             "  Upgrade trigger: Meeting correlation (paid feature)",
             "  Social proof: 'One person starts tracking, then everyone wants their number'",
             "  Virality: Weekly Wrapped shareable cards",
             "  Team spread: First user shows coworkers their stress data",
         ], PERSONA_COLORS["tech"])

# Quote
add_shape(slide, Inches(0.8), Inches(4.6), Inches(11.7), Inches(2.5), WHITE, corner_radius=0.02)
add_accent_bar(slide, Inches(0.8), Inches(4.6), Pt(5), Inches(2.5), PERSONA_COLORS["tech"])

add_text_box(slide, Inches(1.3), Inches(4.8), Inches(10.7), Inches(0.8),
             '"I thought my Mondays were fine. Lucid showed me they weren\'t."',
             font_size=22, bold=False, color=DARK_TEXT, font_name="Playfair Display", italic=True)

add_text_box(slide, Inches(1.3), Inches(5.5), Inches(10.7), Inches(1.2),
             "This persona discovers value through revelation — showing them data about themselves they didn't know. The 'aha moment' is seeing their first meeting correlation: 'Your stress is 40% higher after back-to-back meetings.' That single insight converts trial to paid.\n\nKey insight: This persona doesn't want therapy. They want a dashboard for their mind — the same way they have dashboards for their code, their infrastructure, and their body.",
             font_size=13, color=BODY_TEXT)


# ═══════════════════════════════════════════════════════════
# SLIDE 6: Persona 2 — The Optimization Architect (Profile)
# ═══════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BG_COLOR)

add_shape(slide, Inches(0), Inches(0), SLIDE_W, Inches(1.2), DARK_BG)
add_accent_bar(slide, Inches(0), Inches(1.2), SLIDE_W, Pt(4), PERSONA_COLORS["longevity"])

add_text_box(slide, Inches(0.8), Inches(0.15), Inches(4), Inches(0.4),
             "PERSONA 2", font_size=12, bold=True, color=PERSONA_COLORS["longevity"], font_name="Inter")
add_text_box(slide, Inches(0.8), Inches(0.45), Inches(8), Inches(0.6),
             "The Optimization Architect", font_size=30, bold=True, color=WHITE, font_name="Playfair Display")

# Left column — Profile
add_text_box(slide, Inches(0.8), Inches(1.6), Inches(5.5), Inches(0.4),
             "Who They Are", font_size=20, bold=True, color=DARK_TEXT, font_name="Playfair Display")

profile_items = [
    ("Demographics", "28-50 years old  |  $100K-$300K+ income  |  Urban\nFounders, executives, high-performers, early tech adopters"),
    ("Psychographics", "Identity = 'self-knowledge through numbers.' They track glucose,\nHRV, sleep stages, blood biomarkers. Always the first to try new\nhealth tech. Healthspan > lifespan mindset (energy, cognition NOW)."),
    ("Current Stack", "Oura Ring ($400+), WHOOP ($199-$399/yr), Apple Watch, CGM\nBlood panels, sauna, cold plunge, supplements protocol\nSpend $5,321+/yr on wellness (often much more)"),
    ("The Gap", "'I have 50 data streams for my body and ZERO for my mind.'\nNo objective mental/emotional data in their health stack.\nVoice is the missing biomarker layer."),
]

y_pos = Inches(2.1)
for label, desc in profile_items:
    add_text_box(slide, Inches(0.8), y_pos, Inches(1.8), Inches(0.3),
                 label, font_size=11, bold=True, color=PERSONA_COLORS["longevity"], font_name="Inter")
    add_text_box(slide, Inches(2.8), y_pos, Inches(4.0), Inches(1.0),
                 desc, font_size=11, color=BODY_TEXT)
    y_pos += Inches(1.2)

# Right column — Market Stats
stats_card = add_shape(slide, Inches(7.5), Inches(1.6), Inches(5.0), Inches(5.5), WHITE, corner_radius=0.02)
add_accent_bar(slide, Inches(7.5), Inches(1.6), Inches(5.0), Pt(4), PERSONA_COLORS["longevity"])

add_text_box(slide, Inches(7.8), Inches(1.8), Inches(4.5), Inches(0.4),
             "Market Opportunity", font_size=18, bold=True, color=DARK_TEXT, font_name="Playfair Display")

market_stats = [
    ("$38B", "global biohacking market (2025), growing to $217B by 2035"),
    ("19%", "CAGR for biohacking market — explosive growth"),
    ("$500M+", "Oura revenue in 2024, on track for $1B in 2025"),
    ("155K", "members in r/quantifiedself — highly engaged community"),
    ("350K+", "members in r/Biohackers subreddit"),
    ("2-3x", "higher WTP for health subscriptions vs general population"),
    ("52%", "YoY spending growth among Oura users in 2024"),
]

y_pos = Inches(2.35)
for stat, desc in market_stats:
    add_text_box(slide, Inches(7.8), y_pos, Inches(1.2), Inches(0.35),
                 stat, font_size=18, bold=True, color=PERSONA_COLORS["longevity"], font_name="Inter")
    add_text_box(slide, Inches(9.1), y_pos + Inches(0.03), Inches(3.2), Inches(0.35),
                 desc, font_size=11, color=BODY_TEXT)
    y_pos += Inches(0.55)


# ═══════════════════════════════════════════════════════════
# SLIDE 7: Persona 2 — Pain Points & Messaging
# ═══════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BG_COLOR)

add_accent_bar(slide, Inches(0), Inches(0), Pt(6), SLIDE_H, PERSONA_COLORS["longevity"])

add_text_box(slide, Inches(0.8), Inches(0.4), Inches(4), Inches(0.3),
             "PERSONA 2  |  THE OPTIMIZATION ARCHITECT", font_size=11, bold=True, color=PERSONA_COLORS["longevity"], font_name="Inter")
add_text_box(slide, Inches(0.8), Inches(0.7), Inches(8), Inches(0.5),
             "Pain Points, Hooks & Feature Value", font_size=26, bold=True, color=DARK_TEXT, font_name="Playfair Display")

add_card(slide, Inches(0.8), Inches(1.5), Inches(5.8), Inches(2.8),
         "Core Pain Points", [
             "  'My body has a dashboard. My mind doesn't.'",
             "  They track HRV, sleep, glucose — but have ZERO emotional data.",
             "  No mental health biomarker exists in their stack (PMC 2023).",
             "  They want the missing data stream, not advice or treatment.",
             "  Frustrated by subjective mood journals — want objective signals.",
             "  Looking for cross-correlations between mental and physical data.",
         ], PERSONA_COLORS["longevity"])

add_card(slide, Inches(7.0), Inches(1.5), Inches(5.8), Inches(2.8),
         "Messaging That Resonates", [
             "  'Your Oura Tracks Your Body. Nothing Tracks Your Mind. Until Now.'",
             "  'Complete Your Health Stack.'",
             "  'You Compare Sleep Scores. Imagine Comparing Stress Resilience.'",
             "  'Your Mind Is 9.8 Years Behind Your Body.'",
             "  'The quantified self has a qualitative gap. Lucid fills it.'",
             "  Scientific credibility language: '244M parameters, 35K patients.'",
         ], PERSONA_COLORS["longevity"])

# Feature Value Ranking
add_text_box(slide, Inches(0.8), Inches(4.5), Inches(6), Inches(0.4),
             "Feature Value Ranking (Highest to Lowest)", font_size=16, bold=True, color=DARK_TEXT, font_name="Playfair Display")

features = [
    ("1. Correlation Explorer", "Cross-ref voice data with Oura, WHOOP, Apple Health — the 'killer link'", "10/10"),
    ("2. Raw Acoustic Data Export", "8 features, CSV/JSON/API — they want to build their own analyses", "10/10"),
    ("3. Daily Wellness Score", "One number for mental state — like Oura Readiness Score for the mind", "9/10"),
    ("4. Trend Lines & Trajectories", "Multi-week mood/energy/stress patterns over time", "9/10"),
    ("5. Echoes (Pattern Discovery)", "'Stress 40% higher on days you sleep <6 hours' — cross-modal insights", "8/10"),
    ("6. API Access", "Plug into Notion, Obsidian, Google Sheets — full data ownership", "8/10"),
]

y_pos = Inches(5.0)
for feat, desc, score in features:
    add_text_box(slide, Inches(0.8), y_pos, Inches(3.5), Inches(0.3),
                 feat, font_size=12, bold=True, color=DARK_TEXT)
    add_text_box(slide, Inches(4.5), y_pos, Inches(6.5), Inches(0.3),
                 desc, font_size=11, color=BODY_TEXT)
    add_text_box(slide, Inches(11.8), y_pos, Inches(1.0), Inches(0.3),
                 score, font_size=12, bold=True, color=PERSONA_COLORS["longevity"], alignment=PP_ALIGN.RIGHT)
    y_pos += Inches(0.38)


# ═══════════════════════════════════════════════════════════
# SLIDE 8: Persona 2 — Acquisition Strategy
# ═══════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BG_COLOR)

add_accent_bar(slide, Inches(0), Inches(0), Pt(6), SLIDE_H, PERSONA_COLORS["longevity"])

add_text_box(slide, Inches(0.8), Inches(0.4), Inches(4), Inches(0.3),
             "PERSONA 2  |  THE OPTIMIZATION ARCHITECT", font_size=11, bold=True, color=PERSONA_COLORS["longevity"], font_name="Inter")
add_text_box(slide, Inches(0.8), Inches(0.7), Inches(8), Inches(0.5),
             "Acquisition & Conversion Strategy", font_size=26, bold=True, color=DARK_TEXT, font_name="Playfair Display")

add_card(slide, Inches(0.8), Inches(1.5), Inches(3.7), Inches(2.8),
         "Top Acquisition Channels", [
             "  1. Reddit (r/quantifiedself, r/Biohackers, r/longevity)",
             "  2. YouTube (long-form 'I tracked my mind for 30 days')",
             "  3. Podcasts (Huberman-adjacent health optimization)",
             "  4. Twitter/X (health tech, biohacker community)",
             "  5. Product Hunt (early adopter discovery)",
         ], PERSONA_COLORS["longevity"])

add_card(slide, Inches(4.8), Inches(1.5), Inches(3.7), Inches(2.8),
         "Hook Messages That Convert", [
             "  'I added the mental layer to my health stack. Game changer.'",
             "  'My Oura told me I slept great. Lucid told me I was stressed.'",
             "  'Voice biomarkers: the data stream biohackers are sleeping on.'",
             "  'I found the correlation between my HRV and stress. Finally.'",
             "  Stack comparison visuals: physical stack + Lucid = complete.",
         ], PERSONA_COLORS["longevity"])

add_card(slide, Inches(8.8), Inches(1.5), Inches(3.7), Inches(2.8),
         "Conversion Strategy", [
             "  Free tier: Basic wellness score (shows the gap in their stack)",
             "  Upgrade trigger: Correlation Explorer + data export (paid)",
             "  'Complete your stack' CTA on quantified-self landing page",
             "  API/data export as premium differentiator",
             "  Community seeding in QS forums with genuine usage data",
         ], PERSONA_COLORS["longevity"])

# Insight box
add_shape(slide, Inches(0.8), Inches(4.6), Inches(11.7), Inches(2.5), WHITE, corner_radius=0.02)
add_accent_bar(slide, Inches(0.8), Inches(4.6), Pt(5), Inches(2.5), PERSONA_COLORS["longevity"])

add_text_box(slide, Inches(1.3), Inches(4.8), Inches(10.7), Inches(0.8),
             '"Your body has a dashboard. Your mind deserves one too."',
             font_size=22, bold=False, color=DARK_TEXT, font_name="Playfair Display", italic=True)

add_text_box(slide, Inches(1.3), Inches(5.5), Inches(10.7), Inches(1.2),
             "This persona already believes in data-driven health. They don't need to be convinced that monitoring works — they need to be convinced that VOICE is a valid biomarker. Lead with the science: '244M-parameter model, 35,000 patients, clinically validated.' Then show them the data gap their $400 Oura Ring can't fill.\n\nKey insight: This is your highest-LTV persona. They pay $400+ for rings, $200+/yr for subscriptions. $14.99/month is noise in their wellness budget. They churn the least because data compounds.",
             font_size=13, color=BODY_TEXT)


# ═══════════════════════════════════════════════════════════
# SLIDE 9: Persona 3 — The System Skeptic (Profile)
# ═══════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BG_COLOR)

add_shape(slide, Inches(0), Inches(0), SLIDE_W, Inches(1.2), DARK_BG)
add_accent_bar(slide, Inches(0), Inches(1.2), SLIDE_W, Pt(4), PERSONA_COLORS["skeptic"])

add_text_box(slide, Inches(0.8), Inches(0.15), Inches(4), Inches(0.4),
             "PERSONA 3", font_size=12, bold=True, color=PERSONA_COLORS["skeptic"], font_name="Inter")
add_text_box(slide, Inches(0.8), Inches(0.45), Inches(8), Inches(0.6),
             "The System Skeptic", font_size=30, bold=True, color=WHITE, font_name="Playfair Display")

# Left column — Profile
add_text_box(slide, Inches(0.8), Inches(1.6), Inches(5.5), Inches(0.4),
             "Who They Are", font_size=20, bold=True, color=DARK_TEXT, font_name="Playfair Display")

profile_items = [
    ("Demographics", "22-45 years old  |  All income levels  |  Urban/suburban\nFreelancers, creatives, privacy advocates, independent professionals"),
    ("Psychographics", "Deeply skeptical of institutions — insurance, pharma, therapy industry.\nPrivacy is non-negotiable. They believe mental health is personal, not\nclinical. They've had bad therapy experiences or know someone who has."),
    ("Barriers", "Therapy: $100-$250/session, 34.8% dropout rate, 9.8-year avg to diagnosis\nStigma: Don't want it 'on their record' or in insurance databases\nDistrust: Wary of AI companies harvesting data for profit"),
    ("What They Want", "Self-knowledge WITHOUT a gatekeeper. No referral needed, no\nwaiting room, no diagnosis. Just objective data they control.\nThey'd rather monitor themselves than be monitored by a system."),
]

y_pos = Inches(2.1)
for label, desc in profile_items:
    add_text_box(slide, Inches(0.8), y_pos, Inches(1.8), Inches(0.3),
                 label, font_size=11, bold=True, color=PERSONA_COLORS["skeptic"], font_name="Inter")
    add_text_box(slide, Inches(2.8), y_pos, Inches(4.0), Inches(1.0),
                 desc, font_size=11, color=BODY_TEXT)
    y_pos += Inches(1.2)

# Right column — System Failure Stats
stats_card = add_shape(slide, Inches(7.5), Inches(1.6), Inches(5.0), Inches(5.5), WHITE, corner_radius=0.02)
add_accent_bar(slide, Inches(7.5), Inches(1.6), Inches(5.0), Pt(4), PERSONA_COLORS["skeptic"])

add_text_box(slide, Inches(7.8), Inches(1.8), Inches(4.5), Inches(0.4),
             "Why They Don't Trust the System", font_size=18, bold=True, color=DARK_TEXT, font_name="Playfair Display")

system_stats = [
    ("9.8 yrs", "average delay to mental health diagnosis"),
    ("34.8%", "therapy dropout rate (Swift & Greenberg 2012)"),
    ("4.3%", "of primary care visits include mental health screening"),
    ("4-6%", "EAP utilization rate (despite being free to employees)"),
    ("58%", "cite cost as #1 barrier to therapy"),
    ("$250+", "per session without insurance coverage"),
    ("100%", "of Lucid data stays on YOUR machine, never leaves"),
]

y_pos = Inches(2.35)
for stat, desc in system_stats:
    add_text_box(slide, Inches(7.8), y_pos, Inches(1.2), Inches(0.35),
                 stat, font_size=18, bold=True, color=PERSONA_COLORS["skeptic"], font_name="Inter")
    add_text_box(slide, Inches(9.1), y_pos + Inches(0.03), Inches(3.2), Inches(0.35),
                 desc, font_size=11, color=BODY_TEXT)
    y_pos += Inches(0.55)


# ═══════════════════════════════════════════════════════════
# SLIDE 10: Persona 3 — Pain Points & Messaging
# ═══════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BG_COLOR)

add_accent_bar(slide, Inches(0), Inches(0), Pt(6), SLIDE_H, PERSONA_COLORS["skeptic"])

add_text_box(slide, Inches(0.8), Inches(0.4), Inches(4), Inches(0.3),
             "PERSONA 3  |  THE SYSTEM SKEPTIC", font_size=11, bold=True, color=PERSONA_COLORS["skeptic"], font_name="Inter")
add_text_box(slide, Inches(0.8), Inches(0.7), Inches(8), Inches(0.5),
             "Pain Points, Hooks & Feature Value", font_size=26, bold=True, color=DARK_TEXT, font_name="Playfair Display")

add_card(slide, Inches(0.8), Inches(1.5), Inches(5.8), Inches(2.8),
         "Core Pain Points", [
             "  'I don't need a therapist to tell me I'm stressed. I need data.'",
             "  Afraid mental health data will be used against them (insurance, employers).",
             "  The system screens them 1x/year with a 9-question self-report — absurd.",
             "  Cloud-based wellness apps feel like surveillance, not self-care.",
             "  They've tried Calm/Headspace (3.3% 30-day retention) and abandoned them.",
             "  Want to understand their patterns without involving another human.",
         ], PERSONA_COLORS["skeptic"])

add_card(slide, Inches(7.0), Inches(1.5), Inches(5.8), Inches(2.8),
         "Messaging That Resonates", [
             "  '100% on-device. Your data never leaves your machine. Ever.'",
             "  'No therapist required. No diagnosis. Just clarity.'",
             "  'Download and start — no referral, no insurance, no waiting room.'",
             "  'Your data belongs to you. Not your insurance company.'",
             "  'The system screens you 1x/year. Lucid screens every conversation.'",
             "  Privacy-first, sovereignty language — not medical/clinical.",
         ], PERSONA_COLORS["skeptic"])

# Feature Value Ranking
add_text_box(slide, Inches(0.8), Inches(4.5), Inches(6), Inches(0.4),
             "Feature Value Ranking (Highest to Lowest)", font_size=16, bold=True, color=DARK_TEXT, font_name="Playfair Display")

features = [
    ("1. On-Device Processing", "100% local — no cloud, no server, no data leaving their Mac", "10/10"),
    ("2. No Audio Storage", "Audio is analyzed and discarded — never recorded, never stored", "10/10"),
    ("3. Daily Wellness Score", "Self-knowledge without a gatekeeper — objective, private data", "9/10"),
    ("4. Trend Lines & Trajectories", "See their own patterns over time — self-directed insight", "8/10"),
    ("5. Data Export (CSV/JSON)", "Full data ownership — export, delete, take it anywhere", "8/10"),
    ("6. Grove (Streak System)", "Personal ritual that replaces the 'appointment' they refuse to make", "7/10"),
]

y_pos = Inches(5.0)
for feat, desc, score in features:
    add_text_box(slide, Inches(0.8), y_pos, Inches(3.5), Inches(0.3),
                 feat, font_size=12, bold=True, color=DARK_TEXT)
    add_text_box(slide, Inches(4.5), y_pos, Inches(6.5), Inches(0.3),
                 desc, font_size=11, color=BODY_TEXT)
    add_text_box(slide, Inches(11.8), y_pos, Inches(1.0), Inches(0.3),
                 score, font_size=12, bold=True, color=PERSONA_COLORS["skeptic"], alignment=PP_ALIGN.RIGHT)
    y_pos += Inches(0.38)


# ═══════════════════════════════════════════════════════════
# SLIDE 11: Persona 3 — Acquisition Strategy
# ═══════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BG_COLOR)

add_accent_bar(slide, Inches(0), Inches(0), Pt(6), SLIDE_H, PERSONA_COLORS["skeptic"])

add_text_box(slide, Inches(0.8), Inches(0.4), Inches(4), Inches(0.3),
             "PERSONA 3  |  THE SYSTEM SKEPTIC", font_size=11, bold=True, color=PERSONA_COLORS["skeptic"], font_name="Inter")
add_text_box(slide, Inches(0.8), Inches(0.7), Inches(8), Inches(0.5),
             "Acquisition & Conversion Strategy", font_size=26, bold=True, color=DARK_TEXT, font_name="Playfair Display")

add_card(slide, Inches(0.8), Inches(1.5), Inches(3.7), Inches(2.8),
         "Top Acquisition Channels", [
             "  1. TikTok (anti-system, 'therapy is broken' content)",
             "  2. Reddit (r/mentalhealth, r/privacy, r/selfimprovement)",
             "  3. Twitter/X (privacy advocacy, mental health discourse)",
             "  4. Podcasts (independent health, privacy-focused shows)",
             "  5. Word of mouth (strongest for trust-based product)",
         ], PERSONA_COLORS["skeptic"])

add_card(slide, Inches(4.8), Inches(1.5), Inches(3.7), Inches(2.8),
         "Hook Messages That Convert", [
             "  'The system screens you once a year. Your voice can screen you every day.'",
             "  'No therapist, no insurance company, no cloud server ever sees your data.'",
             "  'I stopped going to therapy. I didn't stop monitoring my mental health.'",
             "  'Your voice already knows. It's been telling you the truth your entire life.'",
             "  Content comparing the current system vs Lucid (table format).",
         ], PERSONA_COLORS["skeptic"])

add_card(slide, Inches(8.8), Inches(1.5), Inches(3.7), Inches(2.8),
         "Conversion Strategy", [
             "  Privacy as the LEAD value proposition, not a checkbox",
             "  Transparent technical explanation of on-device processing",
             "  Open about what data exists and how to delete it",
             "  No account creation required for core features",
             "  Trust-building through transparency, not marketing polish",
         ], PERSONA_COLORS["skeptic"])

# Insight box
add_shape(slide, Inches(0.8), Inches(4.6), Inches(11.7), Inches(2.5), WHITE, corner_radius=0.02)
add_accent_bar(slide, Inches(0.8), Inches(4.6), Pt(5), Inches(2.5), PERSONA_COLORS["skeptic"])

add_text_box(slide, Inches(1.3), Inches(4.8), Inches(10.7), Inches(0.8),
             '"I don\'t need someone to tell me how I feel. I need data to show me what I can\'t see."',
             font_size=22, bold=False, color=DARK_TEXT, font_name="Playfair Display", italic=True)

add_text_box(slide, Inches(1.3), Inches(5.5), Inches(10.7), Inches(1.2),
             "This persona is the hardest to acquire but the most loyal once converted. They will NEVER tolerate a privacy breach or a cloud upload. The '100% on-device' message isn't a feature for them — it's the entire reason they'll consider Lucid. Every other mental health app failed them because it required trust in a third party.\n\nKey insight: Don't market to them like a health app. Market like a privacy tool that happens to measure wellness. Lead with sovereignty, not clinical accuracy.",
             font_size=13, color=BODY_TEXT)


# ═══════════════════════════════════════════════════════════
# SLIDE 12: Persona 4 — The Performance Athlete (Profile)
# ═══════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BG_COLOR)

add_shape(slide, Inches(0), Inches(0), SLIDE_W, Inches(1.2), DARK_BG)
add_accent_bar(slide, Inches(0), Inches(1.2), SLIDE_W, Pt(4), PERSONA_COLORS["athlete"])

add_text_box(slide, Inches(0.8), Inches(0.15), Inches(4), Inches(0.4),
             "PERSONA 4", font_size=12, bold=True, color=PERSONA_COLORS["athlete"], font_name="Inter")
add_text_box(slide, Inches(0.8), Inches(0.45), Inches(8), Inches(0.6),
             "The Performance Athlete", font_size=30, bold=True, color=WHITE, font_name="Playfair Display")

# Left column — Profile
add_text_box(slide, Inches(0.8), Inches(1.6), Inches(5.5), Inches(0.4),
             "Who They Are", font_size=20, bold=True, color=DARK_TEXT, font_name="Playfair Display")

profile_items = [
    ("Demographics", "25-50 years old  |  $75K-$200K+ income  |  Urban/suburban\nWeekend warriors, CrossFit athletes, marathon runners, triathletes"),
    ("Psychographics", "Competitive, disciplined, data-obsessed about performance. They\ntrack every rep, every split, every recovery metric. Mental health\nis the #1 limiter they don't measure. Identity: 'I train harder.'"),
    ("Current Stack", "WHOOP ($199-$399/yr), Garmin/Strava, Oura Ring, TrainingPeaks\nNutrition tracking, sleep optimization, periodization plans\nZero tools for mental readiness or stress impact on performance"),
    ("The Gap", "35% of athletes suffer mental health crises. 46.4% meet criteria\nfor at least one mental health problem. They track VO2max, HRV,\nand lactate — but not the anxiety that costs them race day."),
]

y_pos = Inches(2.1)
for label, desc in profile_items:
    add_text_box(slide, Inches(0.8), y_pos, Inches(1.8), Inches(0.3),
                 label, font_size=11, bold=True, color=PERSONA_COLORS["athlete"], font_name="Inter")
    add_text_box(slide, Inches(2.8), y_pos, Inches(4.0), Inches(1.0),
                 desc, font_size=11, color=BODY_TEXT)
    y_pos += Inches(1.2)

# Right column — Stats
stats_card = add_shape(slide, Inches(7.5), Inches(1.6), Inches(5.0), Inches(5.5), WHITE, corner_radius=0.02)
add_accent_bar(slide, Inches(7.5), Inches(1.6), Inches(5.0), Pt(4), PERSONA_COLORS["athlete"])

add_text_box(slide, Inches(7.8), Inches(1.8), Inches(4.5), Inches(0.4),
             "The Mental Performance Gap", font_size=18, bold=True, color=DARK_TEXT, font_name="Playfair Display")

athlete_stats = [
    ("35%", "of athletes suffer mental health crises (burnout, depression)"),
    ("46.4%", "of athletes meet criteria for a mental health problem"),
    ("60%+", "of recreational athletes report moderate-to-severe stress"),
    ("84.7%", "voice biomarker accuracy for stress detection"),
    ("6.7 min", "AI detects stress before athletes self-report it"),
    ("$260M+", "WHOOP revenue — athletes pay for recovery data"),
    ("$0", "extra hardware — Lucid adds the mental layer"),
]

y_pos = Inches(2.35)
for stat, desc in athlete_stats:
    add_text_box(slide, Inches(7.8), y_pos, Inches(1.2), Inches(0.35),
                 stat, font_size=18, bold=True, color=PERSONA_COLORS["athlete"], font_name="Inter")
    add_text_box(slide, Inches(9.1), y_pos + Inches(0.03), Inches(3.2), Inches(0.35),
                 desc, font_size=11, color=BODY_TEXT)
    y_pos += Inches(0.55)


# ═══════════════════════════════════════════════════════════
# SLIDE 13: Persona 4 — Pain Points & Messaging
# ═══════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BG_COLOR)

add_accent_bar(slide, Inches(0), Inches(0), Pt(6), SLIDE_H, PERSONA_COLORS["athlete"])

add_text_box(slide, Inches(0.8), Inches(0.4), Inches(4), Inches(0.3),
             "PERSONA 4  |  THE PERFORMANCE ATHLETE", font_size=11, bold=True, color=PERSONA_COLORS["athlete"], font_name="Inter")
add_text_box(slide, Inches(0.8), Inches(0.7), Inches(8), Inches(0.5),
             "Pain Points, Hooks & Feature Value", font_size=26, bold=True, color=DARK_TEXT, font_name="Playfair Display")

add_card(slide, Inches(0.8), Inches(1.5), Inches(5.8), Inches(2.8),
         "Core Pain Points", [
             "  'I train 6 days a week. I have no idea what stress costs me on race day.'",
             "  They track every physical metric but ZERO mental readiness data.",
             "  Anxiety causes distraction, execution errors, and performance breakdowns.",
             "  Mental health is the #1 limiter — but they treat it as a character flaw.",
             "  Recovery is physical only: 'Am I mentally ready?' has no answer.",
             "  Stress directly predicts injury risk and training quality — unmeasured.",
         ], PERSONA_COLORS["athlete"])

add_card(slide, Inches(7.0), Inches(1.5), Inches(5.8), Inches(2.8),
         "Messaging That Resonates", [
             "  'Your WHOOP tracks recovery. Nothing tracks what's holding you back.'",
             "  '35% of athletes hit a wall that has nothing to do with fitness.'",
             "  'You train your body 6 days a week. When did you last check your mind?'",
             "  'The mental edge is the only edge left to gain.'",
             "  'AI detects stress 6.7 minutes before you feel it.'",
             "  Competitive, data-driven language. Think: Strava, WHOOP, TrainingPeaks.",
         ], PERSONA_COLORS["athlete"])

# Feature Value Ranking
add_text_box(slide, Inches(0.8), Inches(4.5), Inches(6), Inches(0.4),
             "Feature Value Ranking (Highest to Lowest)", font_size=16, bold=True, color=DARK_TEXT, font_name="Playfair Display")

features = [
    ("1. Stress & Calm Scores", "Pre-workout/race mental readiness — objective performance signal", "10/10"),
    ("2. Correlation Explorer", "Cross-ref with WHOOP, Oura, Garmin — find mental-physical links", "10/10"),
    ("3. Trend Lines & Trajectories", "Training block mental load tracking over weeks/months", "9/10"),
    ("4. Morning Briefing", "'Yesterday's long run spiked stress — consider active recovery today'", "9/10"),
    ("5. Daily Wellness Score", "Mental readiness score — like WHOOP Recovery but for your mind", "8/10"),
    ("6. Echoes (Pattern Discovery)", "'Stress 35% higher during race week' — training cycle insights", "8/10"),
]

y_pos = Inches(5.0)
for feat, desc, score in features:
    add_text_box(slide, Inches(0.8), y_pos, Inches(3.5), Inches(0.3),
                 feat, font_size=12, bold=True, color=DARK_TEXT)
    add_text_box(slide, Inches(4.5), y_pos, Inches(6.5), Inches(0.3),
                 desc, font_size=11, color=BODY_TEXT)
    add_text_box(slide, Inches(11.8), y_pos, Inches(1.0), Inches(0.3),
                 score, font_size=12, bold=True, color=PERSONA_COLORS["athlete"], alignment=PP_ALIGN.RIGHT)
    y_pos += Inches(0.38)


# ═══════════════════════════════════════════════════════════
# SLIDE 14: Persona 4 — Acquisition Strategy
# ═══════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BG_COLOR)

add_accent_bar(slide, Inches(0), Inches(0), Pt(6), SLIDE_H, PERSONA_COLORS["athlete"])

add_text_box(slide, Inches(0.8), Inches(0.4), Inches(4), Inches(0.3),
             "PERSONA 4  |  THE PERFORMANCE ATHLETE", font_size=11, bold=True, color=PERSONA_COLORS["athlete"], font_name="Inter")
add_text_box(slide, Inches(0.8), Inches(0.7), Inches(8), Inches(0.5),
             "Acquisition & Conversion Strategy", font_size=26, bold=True, color=DARK_TEXT, font_name="Playfair Display")

add_card(slide, Inches(0.8), Inches(1.5), Inches(3.7), Inches(2.8),
         "Top Acquisition Channels", [
             "  1. Reddit (r/crossfit 400K+, r/running 2M+, r/triathlon)",
             "  2. Strava community / running clubs / CrossFit boxes",
             "  3. YouTube (training vlogs, race recaps, recovery content)",
             "  4. Podcast sponsorships (endurance/CrossFit/fitness shows)",
             "  5. Micro-influencers (CrossFit coaches, running coaches)",
         ], PERSONA_COLORS["athlete"])

add_card(slide, Inches(4.8), Inches(1.5), Inches(3.7), Inches(2.8),
         "Hook Messages That Convert", [
             "  'I tracked my mental readiness for 30 days of training.'",
             "  'My WHOOP said I was recovered. Lucid said I wasn't ready.'",
             "  'The data that explained why I bonked at mile 20.'",
             "  'Your stress on race week is 35% higher. Here's the proof.'",
             "  Before/after stress data around race day / competition.",
         ], PERSONA_COLORS["athlete"])

add_card(slide, Inches(8.8), Inches(1.5), Inches(3.7), Inches(2.8),
         "Conversion Strategy", [
             "  Free tier: Basic stress/wellness score (show the gap)",
             "  Upgrade trigger: Correlation Explorer + trends (paid)",
             "  'Mental readiness' framing — not 'mental health'",
             "  WHOOP/Garmin/Strava integration as premium hook",
             "  Competitive edge positioning, not self-care",
         ], PERSONA_COLORS["athlete"])

# Insight box
add_shape(slide, Inches(0.8), Inches(4.6), Inches(11.7), Inches(2.5), WHITE, corner_radius=0.02)
add_accent_bar(slide, Inches(0.8), Inches(4.6), Pt(5), Inches(2.5), PERSONA_COLORS["athlete"])

add_text_box(slide, Inches(1.3), Inches(4.8), Inches(10.7), Inches(0.8),
             '"You train your body 6 days a week. When did you last check your mind?"',
             font_size=22, bold=False, color=DARK_TEXT, font_name="Playfair Display", italic=True)

add_text_box(slide, Inches(1.3), Inches(5.5), Inches(10.7), Inches(1.2),
             "This persona already believes in data-driven performance. They don't need to be sold on tracking — they need to see mental readiness as the missing metric. WHOOP proved athletes will pay $260M+/year for recovery data. Lucid adds the layer WHOOP can't: how stress, anxiety, and mental load impact their training and race-day performance.\n\nKey insight: Frame Lucid as a COMPETITIVE EDGE tool, not a wellness tool. 'The mental edge is the only edge left to gain.' This persona responds to performance language, leaderboard thinking, and the idea that their competitors don't have this data yet.",
             font_size=13, color=BODY_TEXT)


# ═══════════════════════════════════════════════════════════
# SLIDE 15: Cross-Persona Feature Value Matrix
# ═══════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BG_COLOR)

add_text_box(slide, Inches(0.8), Inches(0.4), Inches(8), Inches(0.5),
             "Feature Value Matrix: All Personas", font_size=28, bold=True, color=DARK_TEXT, font_name="Playfair Display")
add_accent_bar(slide, Inches(0.8), Inches(1.0), Inches(1.5), Pt(3), STEEL_BLUE)

add_text_box(slide, Inches(0.8), Inches(1.2), Inches(8), Inches(0.4),
             "How each feature ranks across all four personas. Higher number = higher value to that persona.",
             font_size=13, color=BODY_TEXT)

# Table header
header_y = Inches(1.8)
col_positions = [Inches(0.8), Inches(4.5), Inches(6.5), Inches(8.3), Inches(10.1)]
col_widths = [Inches(3.5), Inches(1.8), Inches(1.6), Inches(1.6), Inches(1.6)]

# Header row
add_shape(slide, Inches(0.6), header_y, Inches(12.0), Inches(0.5), DARK_BG)
headers = ["Feature", "Builder", "Architect", "Skeptic", "Athlete"]
colors = [WHITE, PERSONA_COLORS["tech"], PERSONA_COLORS["longevity"], PERSONA_COLORS["skeptic"], PERSONA_COLORS["athlete"]]
for i, (header, col_color) in enumerate(zip(headers, colors)):
    if i == 0:
        c = WHITE
    else:
        c = ACCENT_LIGHT
    add_text_box(slide, col_positions[i], header_y + Inches(0.05), col_widths[i], Inches(0.4),
                 header, font_size=12, bold=True, color=c, font_name="Inter", alignment=PP_ALIGN.CENTER if i > 0 else PP_ALIGN.LEFT)

# Data rows
features_matrix = [
    ("Wellness Score (0-100 daily)", "9", "9", "9", "8"),
    ("Meeting Stress Correlation", "10", "6", "5", "4"),
    ("Correlation Explorer (cross-device)", "7", "10", "5", "10"),
    ("Raw Acoustic Data Export", "4", "10", "8", "5"),
    ("Stress & Calm Scores (real-time)", "9", "8", "8", "10"),
    ("Morning Briefing", "9", "7", "6", "9"),
    ("Weekly Wrapped", "8", "6", "5", "6"),
    ("Grove (Streak System)", "6", "7", "7", "5"),
    ("Echoes (Pattern Discovery)", "7", "8", "6", "8"),
    ("On-Device Privacy", "6", "7", "10", "6"),
    ("No Audio Storage", "5", "5", "10", "5"),
    ("Trend Lines & Trajectories", "8", "9", "8", "9"),
    ("API Access", "5", "8", "8", "6"),
    ("Data Export (CSV/JSON)", "5", "8", "8", "6"),
]

row_y = header_y + Inches(0.55)
for j, (feature, *scores) in enumerate(features_matrix):
    bg_color = WHITE if j % 2 == 0 else MID_GRAY
    add_shape(slide, Inches(0.6), row_y, Inches(12.0), Inches(0.35), bg_color)

    add_text_box(slide, col_positions[0], row_y + Inches(0.02), col_widths[0], Inches(0.3),
                 feature, font_size=11, color=DARK_TEXT, font_name="Inter")

    for i, score in enumerate(scores):
        score_val = int(score)
        if score_val >= 9:
            score_color = RGBColor(0x2D, 0x7D, 0x46)  # Dark green
        elif score_val >= 7:
            score_color = PERSONA_COLORS[["tech", "longevity", "skeptic", "athlete"][i]]
        else:
            score_color = RGBColor(0x99, 0xA0, 0xA8)  # Gray

        add_text_box(slide, col_positions[i + 1], row_y + Inches(0.02), col_widths[i + 1], Inches(0.3),
                     score, font_size=12, bold=score_val >= 9, color=score_color,
                     font_name="Inter", alignment=PP_ALIGN.CENTER)

    row_y += Inches(0.35)


# ═══════════════════════════════════════════════════════════
# SLIDE 16: Messaging Framework Summary
# ═══════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BG_COLOR)

add_text_box(slide, Inches(0.8), Inches(0.4), Inches(8), Inches(0.5),
             "Messaging Framework by Persona", font_size=28, bold=True, color=DARK_TEXT, font_name="Playfair Display")
add_accent_bar(slide, Inches(0.8), Inches(1.0), Inches(1.5), Pt(3), STEEL_BLUE)

# Four messaging cards
messaging_data = [
    ("The Burned-Out Builder", PERSONA_COLORS["tech"],
     "Lead With:", "Productivity impact — stress data as a work optimization tool",
     "Tone:", "Direct, data-driven, no-nonsense. Think: Notion, Linear, Arc.",
     "Avoid:", "Therapy language, vulnerability messaging, self-care framing",
     "Hero Stat:", '"23 meetings. 4 stress spikes. 1 that matters."',
     "CTA:", '"See Your Stress Patterns"'),

    ("The Optimization Architect", PERSONA_COLORS["longevity"],
     "Lead With:", "Data gap — the missing layer in their health stack",
     "Tone:", "Scientific, quantified, precise. Think: Oura, WHOOP, Levels.",
     "Avoid:", "Emotional language, therapy comparisons, clinical framing",
     "Hero Stat:", '"9.8 years. That\'s how far your mind is behind your body."',
     "CTA:", '"Complete Your Health Stack"'),

    ("The System Skeptic", PERSONA_COLORS["skeptic"],
     "Lead With:", "Privacy & sovereignty — your data, your machine, your rules",
     "Tone:", "Transparent, anti-establishment, empowering. Think: Signal, Proton.",
     "Avoid:", "Clinical claims, institutional partnerships, corporate language",
     "Hero Stat:", '"100% on-device. Zero cloud. Zero compromise."',
     "CTA:", '"Download and Start — No Account Required"'),

    ("The Performance Athlete", PERSONA_COLORS["athlete"],
     "Lead With:", "Mental edge — stress data as a performance optimization tool",
     "Tone:", "Competitive, aspirational, data-driven. Think: WHOOP, Strava, TrainingPeaks.",
     "Avoid:", "Wellness/self-care framing, therapy language, passive monitoring",
     "Hero Stat:", '"35% of athletes hit a wall that has nothing to do with fitness."',
     "CTA:", '"Train Your Mental Edge"'),
]

for i, (name, color, *pairs) in enumerate(messaging_data):
    x = Inches(0.8) + Inches(i * 3.1)
    y = Inches(1.3)
    w = Inches(2.9)
    h = Inches(5.8)

    card = add_shape(slide, x, y, w, h, WHITE, corner_radius=0.02)
    add_accent_bar(slide, x, y, w, Pt(5), color)

    add_text_box(slide, x + Inches(0.15), y + Inches(0.2), w - Inches(0.3), Inches(0.4),
                 name, font_size=14, bold=True, color=DARK_TEXT, font_name="Playfair Display")

    curr_y = y + Inches(0.6)
    for j in range(0, len(pairs), 2):
        label = pairs[j]
        value = pairs[j + 1]

        add_text_box(slide, x + Inches(0.15), curr_y, w - Inches(0.3), Inches(0.25),
                     label, font_size=10, bold=True, color=color, font_name="Inter")
        add_text_box(slide, x + Inches(0.15), curr_y + Inches(0.22), w - Inches(0.3), Inches(0.6),
                     value, font_size=10, color=BODY_TEXT)
        curr_y += Inches(0.85)


# ═══════════════════════════════════════════════════════════
# SLIDE 17: Landing Page Strategy
# ═══════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BG_COLOR)

add_text_box(slide, Inches(0.8), Inches(0.4), Inches(8), Inches(0.5),
             "Landing Page Strategy by Persona", font_size=28, bold=True, color=DARK_TEXT, font_name="Playfair Display")
add_accent_bar(slide, Inches(0.8), Inches(1.0), Inches(1.5), Pt(3), STEEL_BLUE)

add_text_box(slide, Inches(0.8), Inches(1.2), Inches(10), Inches(0.5),
             "Each persona gets a dedicated landing page with tailored messaging, hero copy, and feature emphasis. Two already exist.",
             font_size=13, color=BODY_TEXT)

# Landing page cards
lp_data = [
    ("burnout.html", "LIVE", PERSONA_COLORS["tech"],
     "The Burned-Out Builder",
     '"You Had 23 Meetings Last Week.\nWhich Ones Are Burning You Out?"',
     "Meeting stress correlation, workday stress scores, burnout trajectory, calendar integration focus. Tabs: Stress & Calm, Mood & Energy, Correlations, Raw Metrics."),

    ("quantified-self.html", "LIVE", PERSONA_COLORS["longevity"],
     "The Optimization Architect",
     '"You Compare Sleep Scores.\nImagine Comparing Stress Resilience."',
     "Health stack completion, Oura/WHOOP cross-correlation, raw data export, API access. 'Your Body Has Data. Your Mind Deserves It Too.' Physical vs Mental stack comparison."),

    ("privacy.html", "PLANNED", PERSONA_COLORS["skeptic"],
     "The System Skeptic",
     '"Your Data. Your Machine.\nYour Rules. No Exceptions."',
     "Lead with on-device architecture, no cloud processing, no audio storage, full data ownership. Technical transparency about how processing works. Compare vs cloud-based alternatives."),

    ("athletes.html", "PLANNED", PERSONA_COLORS["athlete"],
     "The Performance Athlete",
     '"You Train Your Body.\nNow Train Your Edge."',
     "Mental readiness scores, WHOOP/Garmin/Strava cross-correlation, training block stress tracking, race-day mental prep. Pre-workout stress checks, recovery mental load data. Performance athletes who already track everything — except their mind."),
]

for i, (url, status, color, persona, hero, desc) in enumerate(lp_data):
    x = Inches(0.8) + Inches(i * 3.1)
    y = Inches(1.8)
    w = Inches(2.9)
    h = Inches(5.3)

    card = add_shape(slide, x, y, w, h, WHITE, corner_radius=0.02)
    add_accent_bar(slide, x, y, w, Pt(5), color)

    # Status badge
    badge_color = RGBColor(0x2D, 0x7D, 0x46) if status == "LIVE" else RGBColor(0xC4, 0x7A, 0x5A)
    badge = add_shape(slide, x + Inches(0.15), y + Inches(0.2), Inches(0.7), Inches(0.25), badge_color, corner_radius=0.1)
    add_text_box(slide, x + Inches(0.15), y + Inches(0.2), Inches(0.7), Inches(0.25),
                 status, font_size=9, bold=True, color=WHITE, font_name="Inter", alignment=PP_ALIGN.CENTER)

    add_text_box(slide, x + Inches(0.95), y + Inches(0.2), w - Inches(1.1), Inches(0.25),
                 url, font_size=10, color=BODY_TEXT, font_name="Inter")

    add_text_box(slide, x + Inches(0.15), y + Inches(0.6), w - Inches(0.3), Inches(0.3),
                 persona, font_size=12, bold=True, color=DARK_TEXT, font_name="Playfair Display")

    add_text_box(slide, x + Inches(0.15), y + Inches(1.0), w - Inches(0.3), Inches(1.2),
                 hero, font_size=11, bold=True, color=color, font_name="Playfair Display", italic=True)

    add_text_box(slide, x + Inches(0.15), y + Inches(2.4), w - Inches(0.3), Inches(2.6),
                 desc, font_size=10, color=BODY_TEXT)


# ═══════════════════════════════════════════════════════════
# SLIDE 18: Unified Value Proposition
# ═══════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, DARK_BG)

add_accent_bar(slide, Inches(0), Inches(0), Inches(0.15), SLIDE_H, STEEL_BLUE)

add_text_box(slide, Inches(1.2), Inches(0.8), Inches(10), Inches(0.5),
             "The Unified Thread", font_size=32, bold=True, color=WHITE, font_name="Playfair Display")

add_accent_bar(slide, Inches(1.2), Inches(1.4), Inches(1.5), Pt(3), STEEL_BLUE)

add_text_box(slide, Inches(1.2), Inches(1.8), Inches(10), Inches(1.0),
             "Five different personas. Five different pain points. One shared truth:",
             font_size=18, color=ACCENT_LIGHT, font_name="Inter")

add_text_box(slide, Inches(1.2), Inches(2.6), Inches(10), Inches(1.2),
             "Your voice has been leaking your mental state your entire life.\nYou just never had a way to listen.",
             font_size=28, bold=True, color=WHITE, font_name="Playfair Display", italic=True)

# Universal features
add_text_box(slide, Inches(1.2), Inches(4.2), Inches(5), Inches(0.4),
             "What Unites All Four Personas", font_size=18, bold=True, color=ACCENT_LIGHT, font_name="Playfair Display")

universal = [
    ("Passive", "No habit change required. Lucid works in the background."),
    ("Private", "100% on-device processing. No cloud, no server, no compromise."),
    ("Objective", "Involuntary biomarkers — can't be faked, can't be suppressed."),
    ("Continuous", "Every conversation is a data point. Not 1x/year screening."),
    ("Affordable", "$14.99/month — a fraction of one therapy session."),
]

y_pos = Inches(4.7)
for label, desc in universal:
    add_text_box(slide, Inches(1.2), y_pos, Inches(1.5), Inches(0.3),
                 label, font_size=14, bold=True, color=STEEL_BLUE, font_name="Inter")
    add_text_box(slide, Inches(2.8), y_pos, Inches(5), Inches(0.3),
                 desc, font_size=13, color=ACCENT_LIGHT)
    y_pos += Inches(0.42)

# Right side — the promise
add_shape(slide, Inches(7.5), Inches(4.0), Inches(5.0), Inches(3.0), RGBColor(0x22, 0x4A, 0x60), corner_radius=0.02)
add_accent_bar(slide, Inches(7.5), Inches(4.0), Inches(5.0), Pt(4), STEEL_BLUE)

add_text_box(slide, Inches(7.8), Inches(4.3), Inches(4.5), Inches(0.4),
             "The Promise", font_size=18, bold=True, color=WHITE, font_name="Playfair Display")

add_text_box(slide, Inches(7.8), Inches(4.8), Inches(4.5), Inches(1.8),
             "Lucid doesn't replace therapy, coaches, or wearables.\n\nIt fills the gap none of them can — continuous, passive, private insight into the one thing that connects every aspect of your performance and wellbeing:\n\nYour state of mind.",
             font_size=13, color=ACCENT_LIGHT)

# Tagline
add_text_box(slide, Inches(1.2), Inches(6.8), Inches(10), Inches(0.5),
             "Clarity through voice.", font_size=20, color=STEEL_BLUE, font_name="Playfair Display", italic=True)


# ═══════════════════════════════════════════════════════════
# SLIDE 19: Persona 5 — The Disconnected Dude (New)
# ═══════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, RGBColor(0x0E, 0x0E, 0x10))  # Dark bg matching website

# Header bar — dark with electric green accent
add_shape(slide, Inches(0), Inches(0), SLIDE_W, Inches(1.2), RGBColor(0x18, 0x18, 0x1B))
add_accent_bar(slide, Inches(0), Inches(1.2), SLIDE_W, Pt(4), PERSONA_COLORS["disconnected"])

add_text_box(slide, Inches(0.8), Inches(0.15), Inches(4), Inches(0.4),
             "PERSONA 5", font_size=12, bold=True, color=PERSONA_COLORS["disconnected"], font_name="Inter")
add_text_box(slide, Inches(0.8), Inches(0.45), Inches(8), Inches(0.6),
             "The Disconnected Dude", font_size=30, bold=True, color=WHITE, font_name="Inter")

# Left column — Profile
add_text_box(slide, Inches(0.8), Inches(1.6), Inches(5.5), Inches(0.4),
             "Who They Are", font_size=20, bold=True, color=WHITE, font_name="Inter")

profile_items_dd = [
    ("Demographics", "18-30 years old  |  $30K-$90K income  |  Urban/suburban\nGamers, Discord users, remote workers, trades, service industry, students"),
    ("Psychographics", "Deeply alienated from how mental health is discussed.\nWon't touch therapy apps. Distrusts institutions. Values privacy, autonomy.\nIdentity: 'I'm fine' (is not fine)."),
    ("Current Tools", "Discord, Steam, Reddit, YouTube, Twitch\nNo health tracking. Maybe a Fitbit they stopped wearing.\nTried Headspace once. Deleted after 2 sessions."),
    ("Spending", "$20-50/mo on game skins, subscriptions, energy drinks\nWilling to pay for tools, not 'wellness experiences.'"),
]

y_pos = Inches(2.1)
for label, desc in profile_items_dd:
    add_text_box(slide, Inches(0.8), y_pos, Inches(1.8), Inches(0.3),
                 label, font_size=11, bold=True, color=PERSONA_COLORS["disconnected"], font_name="Inter")
    add_text_box(slide, Inches(2.8), y_pos, Inches(4.0), Inches(1.0),
                 desc, font_size=11, color=RGBColor(0xB5, 0xBA, 0xC1))
    y_pos += Inches(1.2)

# Right column — Key Stats
stats_card = add_shape(slide, Inches(7.5), Inches(1.6), Inches(5.0), Inches(5.5),
                       RGBColor(0x26, 0x26, 0x2C), corner_radius=0.02)
add_accent_bar(slide, Inches(7.5), Inches(1.6), Inches(5.0), Pt(4), PERSONA_COLORS["disconnected"])

add_text_box(slide, Inches(7.8), Inches(1.8), Inches(4.5), Inches(0.4),
             "The Crisis Nobody's Solving", font_size=18, bold=True, color=WHITE, font_name="Inter")

pain_stats_dd = [
    ("4x", "Men die by suicide at nearly 4x the rate of women"),
    ("60%", "Of men with mental health symptoms never seek help"),
    ("48 days", "Average wait to see a new therapist in the US"),
    ("$250+", "Average therapy session cost without insurance"),
    ("36%", "Of therapy patients are male — despite equal prevalence"),
    ("50%", "Of US counties have zero practicing psychiatrists"),
    ("73%", "Of 18-29 year-olds report significant stress (APA)"),
]

y_pos = Inches(2.35)
for stat, desc in pain_stats_dd:
    add_text_box(slide, Inches(7.8), y_pos, Inches(1.2), Inches(0.35),
                 stat, font_size=18, bold=True, color=PERSONA_COLORS["disconnected"], font_name="Inter")
    add_text_box(slide, Inches(9.1), y_pos + Inches(0.03), Inches(3.2), Inches(0.35),
                 desc, font_size=11, color=RGBColor(0xB5, 0xBA, 0xC1))
    y_pos += Inches(0.55)


# ═══════════════════════════════════════════════════════════
# SLIDE 20: Persona 5 — Ranking vs Existing Personas
# ═══════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, RGBColor(0x0E, 0x0E, 0x10))

add_accent_bar(slide, Inches(0), Inches(0), Pt(6), SLIDE_H, PERSONA_COLORS["disconnected"])

add_text_box(slide, Inches(0.8), Inches(0.4), Inches(8), Inches(0.3),
             "PERSONA 5  |  THE DISCONNECTED DUDE", font_size=11, bold=True,
             color=PERSONA_COLORS["disconnected"], font_name="Inter")
add_text_box(slide, Inches(0.8), Inches(0.7), Inches(10), Inches(0.5),
             "Ranking Against Existing Personas", font_size=26, bold=True, color=WHITE, font_name="Inter")

# Ranking table
rank_y = Inches(1.5)
rank_cols = [Inches(0.8), Inches(3.8), Inches(5.6), Inches(7.4), Inches(9.2), Inches(11.0)]
rank_w = [Inches(2.8), Inches(1.7), Inches(1.7), Inches(1.7), Inches(1.7), Inches(1.7)]
rank_headers = ["Persona", "TAM\nPotential", "Acquisition\nDifficulty", "Retention\nPotential", "Messaging\nClarity", "Product-\nMarket Fit"]

# Table header row
add_shape(slide, Inches(0.8), rank_y, Inches(11.8), Inches(0.6), RGBColor(0x18, 0x18, 0x1B))
for i, h in enumerate(rank_headers):
    add_text_box(slide, rank_cols[i], rank_y + Inches(0.05), rank_w[i], Inches(0.5),
                 h, font_size=10, bold=True, color=PERSONA_COLORS["disconnected"], font_name="Inter",
                 alignment=PP_ALIGN.CENTER if i > 0 else PP_ALIGN.LEFT)

# Table data
rank_data = [
    ("1. The Burned-Out Builder", "9/10", "6/10", "9/10", "9/10", "10/10", PERSONA_COLORS["tech"]),
    ("2. The Disconnected Dude", "8/10", "8/10", "7/10", "9/10", "8/10", PERSONA_COLORS["disconnected"]),
    ("3. The Optimization Architect", "7/10", "5/10", "8/10", "7/10", "8/10", PERSONA_COLORS["longevity"]),
    ("4. The System Skeptic", "6/10", "7/10", "8/10", "8/10", "7/10", PERSONA_COLORS["skeptic"]),
    ("5. The Performance Athlete", "5/10", "5/10", "7/10", "6/10", "6/10", PERSONA_COLORS["athlete"]),
]

row_y = rank_y + Inches(0.65)
for j, (name, tam, acq, ret, msg, fit, color) in enumerate(rank_data):
    bg = RGBColor(0x1F, 0x1F, 0x23) if j % 2 == 0 else RGBColor(0x18, 0x18, 0x1B)
    row_h = Inches(0.5)
    add_shape(slide, Inches(0.8), row_y, Inches(11.8), row_h, bg)

    # Highlight the Disconnected Dude row
    is_dd = j == 1
    name_color = PERSONA_COLORS["disconnected"] if is_dd else WHITE
    add_text_box(slide, rank_cols[0], row_y + Inches(0.1), rank_w[0], Inches(0.3),
                 name, font_size=12, bold=is_dd, color=name_color, font_name="Inter")

    scores = [tam, acq, ret, msg, fit]
    for i, s in enumerate(scores):
        score_val = int(s.split("/")[0])
        if score_val >= 9:
            sc = PERSONA_COLORS["disconnected"]
        elif score_val >= 7:
            sc = RGBColor(0xFE, 0xE7, 0x5C)
        elif score_val >= 5:
            sc = RGBColor(0xB5, 0xBA, 0xC1)
        else:
            sc = RGBColor(0xED, 0x42, 0x45)
        add_text_box(slide, rank_cols[i + 1], row_y + Inches(0.1), rank_w[i + 1], Inches(0.3),
                     s, font_size=12, bold=True, color=sc, font_name="Inter",
                     alignment=PP_ALIGN.CENTER)

    row_y += Inches(0.5)

# Recommendation box
rec_y = row_y + Inches(0.5)
add_shape(slide, Inches(0.8), rec_y, Inches(11.8), Inches(2.8),
          RGBColor(0x1A, 0x2F, 0x1E), corner_radius=0.02)
add_accent_bar(slide, Inches(0.8), rec_y, Inches(11.8), Pt(4), PERSONA_COLORS["disconnected"])

add_text_box(slide, Inches(1.1), rec_y + Inches(0.25), Inches(5), Inches(0.4),
             "Recommendation: Priority #2 Persona", font_size=18, bold=True, color=WHITE, font_name="Inter")

rec_lines = [
    ("Why #2:", 12, True, PERSONA_COLORS["disconnected"], "Inter", PP_ALIGN.LEFT, False, Pt(4)),
    ("Massive underserved TAM — young men are the largest demographic actively avoiding mental health tools. "
     "They won't use Calm, won't journal, won't book a therapist. But they WILL use a tool that sits in "
     "their menubar, gives them numbers, respects their privacy, and doesn't lecture them.", 11, False,
     RGBColor(0xB5, 0xBA, 0xC1), "Inter", PP_ALIGN.LEFT, False, Pt(12)),
    ("Key risks:", 12, True, PERSONA_COLORS["disconnected"], "Inter", PP_ALIGN.LEFT, False, Pt(4)),
    ("Higher acquisition difficulty (8/10) — this audience is ad-resistant, skeptical, and spread across "
     "fragmented platforms (Discord, Reddit, Twitch). Requires authentic, irreverent messaging that's easy "
     "to get wrong. Retention is moderate (7/10) — they may lose interest if the tool doesn't surface "
     "genuinely surprising insights.", 11, False,
     RGBColor(0xB5, 0xBA, 0xC1), "Inter", PP_ALIGN.LEFT, False, Pt(12)),
    ("Action: Build a dedicated landing page with dark-mode, gaming-adjacent design and no-BS copy. "
     "Test acquisition via Discord communities and Reddit before scaling.", 11, False,
     WHITE, "Inter", PP_ALIGN.LEFT, False, Pt(6)),
]

add_multi_text(slide, Inches(1.1), rec_y + Inches(0.6), Inches(11.2), Inches(2.0), rec_lines)


# ═══════════════════════════════════════════════════════════
# SLIDE 21: Personas Considered & Evaluated (Appendix)
# ═══════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BG_COLOR)

add_text_box(slide, Inches(0.8), Inches(0.4), Inches(10), Inches(0.5),
             "Appendix: Personas Considered & Evaluated", font_size=28, bold=True, color=DARK_TEXT, font_name="Playfair Display")
add_accent_bar(slide, Inches(0.8), Inches(1.0), Inches(1.5), Pt(3), STEEL_BLUE)

add_text_box(slide, Inches(0.8), Inches(1.15), Inches(11), Inches(0.5),
             "We evaluated 7 additional candidate personas before finalizing the top five. Each was assessed for pain depth and Lucid product-market fit.",
             font_size=12, color=BODY_TEXT)

# Table header
table_y = Inches(1.7)
col_x = [Inches(0.6), Inches(3.0), Inches(4.6), Inches(6.4)]
col_w = [Inches(2.3), Inches(1.5), Inches(1.7), Inches(6.0)]

add_shape(slide, Inches(0.6), table_y, Inches(11.8), Inches(0.4), DARK_BG)
table_headers = ["Candidate Persona", "Pain Depth", "Lucid Fit", "Why Not a Top-4 Persona"]
for i, h in enumerate(table_headers):
    add_text_box(slide, col_x[i], table_y + Inches(0.05), col_w[i], Inches(0.3),
                 h, font_size=11, bold=True, color=WHITE, font_name="Inter")

# Table rows
candidates = [
    ("Caregivers", "Very High", "Low",
     "78% burnout, 87% stress — but time-poor, budget-constrained, unlikely Mac users. Need relief, not data."),
    ("CEOs / Founders", "High", "Medium",
     "71% burnout, 26% depression — already captured within Persona 1 (burnout) and Persona 2 (biohacker)."),
    ("New Parents / Postpartum", "High", "Low",
     "20% affected by perinatal conditions — but macOS desktop doesn't fit new parent lifestyle. Mobile phase."),
    ("Creative Freelancers", "Medium-High", "Medium",
     "70% burnout, 44% depression — significant overlap with Persona 1 (burnout) and Persona 3 (independent)."),
    ("Teachers / Professors", "Medium", "Low",
     "60% voice disorder risk, 3.7M US teachers — but Lucid measures mental state, not vocal cord health."),
    ("Vocal Professionals", "Medium", "Low",
     "Singers, coaches, podcasters — care about voice-as-instrument, not mental wellness via voice."),
    ("Remote Workers", "Medium", "Medium",
     "Isolation, always-on culture — subsumed by Persona 1 (Burned-Out Builder). Not distinct enough."),
]

row_y = table_y + Inches(0.45)
pain_colors = {
    "Very High": RGBColor(0xC0, 0x39, 0x2B),
    "High": RGBColor(0xD4, 0x6B, 0x08),
    "Medium-High": RGBColor(0xD4, 0x9B, 0x08),
    "Medium": RGBColor(0x99, 0xA0, 0xA8),
}
fit_colors = {
    "Low": RGBColor(0xC0, 0x39, 0x2B),
    "Medium": RGBColor(0xD4, 0x9B, 0x08),
}

for j, (name, pain, fit, reason) in enumerate(candidates):
    bg = WHITE if j % 2 == 0 else MID_GRAY
    add_shape(slide, Inches(0.6), row_y, Inches(11.8), Inches(0.55), bg)

    add_text_box(slide, col_x[0], row_y + Inches(0.05), col_w[0], Inches(0.45),
                 name, font_size=11, bold=True, color=DARK_TEXT, font_name="Inter")
    add_text_box(slide, col_x[1], row_y + Inches(0.05), col_w[1], Inches(0.45),
                 pain, font_size=11, bold=True, color=pain_colors.get(pain, BODY_TEXT), font_name="Inter")
    add_text_box(slide, col_x[2], row_y + Inches(0.05), col_w[2], Inches(0.45),
                 fit, font_size=11, bold=True, color=fit_colors.get(fit, BODY_TEXT), font_name="Inter")
    add_text_box(slide, col_x[3], row_y + Inches(0.05), col_w[3], Inches(0.45),
                 reason, font_size=10, color=BODY_TEXT)

    row_y += Inches(0.55)

# Bottom insight
insight_y = row_y + Inches(0.3)
add_shape(slide, Inches(0.6), insight_y, Inches(11.8), Inches(0.8), RGBColor(0xE8, 0xEE, 0xF4), corner_radius=0.02)
add_text_box(slide, Inches(0.8), insight_y + Inches(0.1), Inches(11.4), Inches(0.6),
             "Conclusion: No candidate persona has both a deeper pain point than the existing five AND strong product-market fit "
             "with Lucid's current form (macOS desktop, voice-based mental wellness, $14.99/month). The four selected personas represent "
             "the optimal balance of pain depth, channel accessibility, and willingness to pay. The Disconnected Dude (Persona 5) was added based on the massive underserved male mental health TAM.",
             font_size=11, color=DARK_TEXT, font_name="Inter")


# ═══════════════════════════════════════════════════════════
# SAVE
# ═══════════════════════════════════════════════════════════
output_path = "/Users/zacharypoll/Desktop/Documents/Claude Code/Lucid/Business Documents/Lucid Persona Analysis.pptx"
prs.save(output_path)
print(f"Saved to: {output_path}")
print(f"Total slides: {len(prs.slides)}")
