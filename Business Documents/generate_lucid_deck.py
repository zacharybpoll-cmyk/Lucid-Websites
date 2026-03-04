#!/usr/bin/env python3
"""
Lucid Investment Deck Generator — v5
Clean professional design, no emoji, fixed layouts, speaker notes.
19 slides. 16:9 format.
"""

import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.chart import XL_CHART_TYPE, XL_LABEL_POSITION
from pptx.chart.data import CategoryChartData
from lxml import etree

# ── Design System ───────────────────────────────────────────────
NAVY     = RGBColor(0x1A, 0x3A, 0x4F)
STEEL    = RGBColor(0x5B, 0x8D, 0xB8)
WHITE    = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_BG = RGBColor(0xF8, 0xF9, 0xFA)

CHARCOAL    = RGBColor(0x1A, 0x1D, 0x21)
BODY_GRAY   = RGBColor(0x5A, 0x62, 0x70)
MEDIUM_GRAY = RGBColor(0x88, 0x88, 0x99)

ACCENT_GREEN = RGBColor(0x4E, 0xC9, 0xB0)
ACCENT_GOLD  = RGBColor(0xD4, 0xA5, 0x37)
ACCENT_CORAL = RGBColor(0xE8, 0x6B, 0x6B)
ACCENT_VIOLET = RGBColor(0x7B, 0x68, 0xD4)

CARD_BG     = RGBColor(0xFF, 0xFF, 0xFF)
CARD_BORDER = RGBColor(0xE0, 0xE0, 0xE8)
DIVIDER     = RGBColor(0xE8, 0xE8, 0xF0)

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)
TOTAL_SLIDES = 19

# ── Asset Paths ─────────────────────────────────────────────────
BASE         = os.path.dirname(os.path.abspath(__file__))
PROJECT_ROOT = os.path.dirname(BASE)
IMG_DASHBOARD = os.path.join(PROJECT_ROOT, "lucid-website/actual-websites/images/dashboard-health-score.png")
IMG_PUBMED    = os.path.join(PROJECT_ROOT, "lucid-website/actual-websites/images/pubmed-study.png")
IMG_ICON      = os.path.join(PROJECT_ROOT, "assets/icon.png")
OUTPUT        = os.path.join(BASE, "Lucid Investment Deck.pptx")

IMG_DIMS = {
    "dashboard": (1390, 912),
    "pubmed":    (1538, 1576),
    "icon":      (1024, 1024),
}

HEADLINE_FONT = "Georgia"
BODY_FONT     = "Helvetica Neue"


# ── Helpers ──────────────────────────────────────────────────────

def fit_image_dims(image_key, max_width_in, max_height_in):
    pw, ph = IMG_DIMS[image_key]
    aspect = pw / ph
    w = max_width_in
    h = w / aspect
    if h > max_height_in:
        h = max_height_in
        w = h * aspect
    return Inches(w), Inches(h)


def add_bg(slide, color):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Emu(0), Emu(0), SLIDE_W, SLIDE_H)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    sp = shape._element
    sp.getparent().remove(sp)
    slide.shapes._spTree.insert(2, sp)


def add_text_box(slide, left, top, width, height, text, font_size=18,
                 color=CHARCOAL, bold=False, alignment=PP_ALIGN.LEFT,
                 font_name=BODY_FONT, line_spacing=1.2, anchor=MSO_ANCHOR.TOP):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    txBox.text_frame.word_wrap = True
    txBox.text_frame.auto_size = None
    p = txBox.text_frame.paragraphs[0]
    p.alignment = alignment
    p.space_after = Pt(0)
    p.space_before = Pt(0)
    pPr = p._p.get_or_add_pPr()
    lnSpc = etree.SubElement(pPr, '{http://schemas.openxmlformats.org/drawingml/2006/main}lnSpc')
    spcPct = etree.SubElement(lnSpc, '{http://schemas.openxmlformats.org/drawingml/2006/main}spcPct')
    spcPct.set('val', str(int(line_spacing * 100000)))
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.color.rgb = color
    run.font.bold = bold
    run.font.name = font_name
    txBox.text_frame._txBody.bodyPr.set('anchor', {
        MSO_ANCHOR.TOP: 't', MSO_ANCHOR.MIDDLE: 'ctr', MSO_ANCHOR.BOTTOM: 'b'
    }.get(anchor, 't'))
    return txBox


def add_multiline_box(slide, left, top, width, height, lines, default_size=18,
                      default_color=BODY_GRAY, default_bold=False,
                      alignment=PP_ALIGN.LEFT, default_font=BODY_FONT,
                      line_spacing=1.2):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    txBox.text_frame.word_wrap = True
    txBox.text_frame.auto_size = None
    for i, line in enumerate(lines):
        if isinstance(line, str):
            line = {"text": line}
        p = txBox.text_frame.paragraphs[0] if i == 0 else txBox.text_frame.add_paragraph()
        p.alignment = line.get("alignment", alignment)
        p.space_after = Pt(line.get("space_after", 4))
        p.space_before = Pt(line.get("space_before", 0))
        pPr = p._p.get_or_add_pPr()
        lnSpc = etree.SubElement(pPr, '{http://schemas.openxmlformats.org/drawingml/2006/main}lnSpc')
        spcPct = etree.SubElement(lnSpc, '{http://schemas.openxmlformats.org/drawingml/2006/main}spcPct')
        spcPct.set('val', str(int(line.get("line_spacing", line_spacing) * 100000)))
        run = p.add_run()
        run.text = line.get("text", "")
        run.font.size = Pt(line.get("size", default_size))
        run.font.color.rgb = line.get("color", default_color)
        run.font.bold = line.get("bold", default_bold)
        run.font.name = line.get("font", default_font)
        if line.get("italic"):
            run.font.italic = True
    return txBox


def add_slide_number(slide, num, color=MEDIUM_GRAY):
    add_text_box(slide, Inches(11.8), Inches(6.9), Inches(1.2), Inches(0.4),
                 f"{num} / {TOTAL_SLIDES}", font_size=10, color=color,
                 alignment=PP_ALIGN.RIGHT)


def add_divider(slide, left, top, width, color=STEEL, thickness=Pt(1.5)):
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, thickness)
    line.fill.solid()
    line.fill.fore_color.rgb = color
    line.line.fill.background()


def add_dot(slide, left, top, size=Inches(0.1), color=STEEL):
    dot = slide.shapes.add_shape(MSO_SHAPE.OVAL, left, top, size, size)
    dot.fill.solid()
    dot.fill.fore_color.rgb = color
    dot.line.fill.background()


def add_section_header(slide, label, dark_bg=False):
    """Standard section label + divider at top of light slides."""
    add_text_box(slide, Inches(0.8), Inches(0.5), Inches(4.0), Inches(0.35),
                 label, font_size=11, color=STEEL, bold=True, font_name=BODY_FONT)
    add_divider(slide, Inches(0.8), Inches(0.9), Inches(2.0), STEEL, Pt(2))


def add_card_box(slide, left, top, width, height, radius=0.06,
                 bg=CARD_BG, border=CARD_BORDER, border_w=Pt(1)):
    box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                  left, top, width, height)
    box.fill.solid()
    box.fill.fore_color.rgb = bg
    if border:
        box.line.color.rgb = border
        box.line.width = border_w
    else:
        box.line.fill.background()
    box.adjustments[0] = radius
    return box


def add_notes(slide, text):
    """Add speaker notes to a slide."""
    notes_tf = slide.notes_slide.notes_text_frame
    notes_tf.text = text


# ── Slide 01 — Title ─────────────────────────────────────────────

def slide_01_title(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, NAVY)

    add_divider(slide, Inches(4.5), Inches(1.5), Inches(4.3), STEEL, Pt(2))

    if os.path.exists(IMG_ICON):
        slide.shapes.add_picture(IMG_ICON, Inches(5.9), Inches(1.8),
                                  Inches(1.5), Inches(1.5))

    add_text_box(slide, Inches(2.5), Inches(3.35), Inches(8.3), Inches(0.9),
                 "Lucid", font_size=60, color=WHITE, bold=True,
                 alignment=PP_ALIGN.CENTER, font_name=HEADLINE_FONT)

    add_text_box(slide, Inches(3.0), Inches(4.35), Inches(7.3), Inches(0.5),
                 "Clarity through voice.", font_size=22, color=STEEL,
                 alignment=PP_ALIGN.CENTER, font_name=HEADLINE_FONT)

    add_text_box(slide, Inches(2.5), Inches(4.95), Inches(8.3), Inches(0.8),
                 "The mental health data you were never given. "
                 "Clinically validated voice biomarker monitoring \u2014 "
                 "100% on-device, on your Mac.",
                 font_size=14, color=RGBColor(0xAA, 0xBB, 0xCC),
                 alignment=PP_ALIGN.CENTER, line_spacing=1.4)

    add_divider(slide, Inches(4.5), Inches(5.95), Inches(4.3), STEEL, Pt(2))

    add_text_box(slide, Inches(9.5), Inches(6.8), Inches(3.5), Inches(0.4),
                 "Confidential \u2022 March 2026", font_size=10,
                 color=RGBColor(0x66, 0x77, 0x88), alignment=PP_ALIGN.RIGHT)

    add_slide_number(slide, 1, color=RGBColor(0x66, 0x77, 0x88))

    add_notes(slide,
        "Open with the core tension: clinical-grade voice biomarker AI has existed for years "
        "inside hospitals and enterprise B2B. Lucid is the first consumer product to democratize "
        "this technology. Lead with the tagline: 'Clarity through voice.' "
        "Key differentiator from day one: 100% on-device. No cloud. No data risk. "
        "This is not a therapy app or a mood journal. It is a data stream \u2014 "
        "the mental health metric that belongs in every wellness stack.")


# ── Slide 02 — The Problem ───────────────────────────────────────

def slide_02_problem(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, LIGHT_BG)
    add_section_header(slide, "THE PROBLEM")

    add_text_box(slide, Inches(0.8), Inches(1.1), Inches(11.5), Inches(0.85),
                 "The Average American Waits 9.8 Years to Get Mental Health Help.",
                 font_size=30, color=CHARCOAL, bold=True,
                 font_name=HEADLINE_FONT, line_spacing=1.15)

    add_text_box(slide, Inches(0.8), Inches(2.0), Inches(11.5), Inches(0.4),
                 "Self-identify. Self-refer. Self-advocate. Then wait.",
                 font_size=15, color=BODY_GRAY)

    # 4 problem cards
    cards = [
        ("9.8-Year Wait",
         "From first symptoms to first treatment.\n"
         "Self-identify, self-refer, self-advocate.\n"
         "The system is not coming to find you.",
         ACCENT_CORAL),
        ("The 1999 Screening Tool",
         "The PHQ-9 is still the clinical standard.\n"
         "Created in 1999. Given once a year.\n"
         "Self-report only. No objective signal.",
         ACCENT_GOLD),
        ("Technology Locked in Clinics",
         "Voice biomarker AI has existed for years.\n"
         "Available only in hospitals and enterprise.\n"
         "Consumers have never had access.",
         STEEL),
        ("Invisible Signals",
         "Voice changes with mental state \u2014\n"
         "jitter, shimmer, pitch, energy all shift.\n"
         "No tool has ever captured this for you.",
         ACCENT_VIOLET),
    ]

    card_w = Inches(2.8)
    card_h = Inches(3.2)
    card_y = Inches(2.6)
    gap    = Inches(0.35)

    for i, (title, body, accent) in enumerate(cards):
        x = Inches(0.8) + i * (card_w + gap)
        add_card_box(slide, x, card_y, card_w, card_h)
        add_divider(slide, x + Inches(0.3), card_y + Inches(0.22),
                    Inches(0.8), accent, Pt(3))
        add_text_box(slide, x + Inches(0.3), card_y + Inches(0.45),
                     card_w - Inches(0.6), Inches(0.5),
                     title, font_size=14, color=CHARCOAL, bold=True,
                     font_name=HEADLINE_FONT, line_spacing=1.1)
        add_text_box(slide, x + Inches(0.3), card_y + Inches(1.0),
                     card_w - Inches(0.6), Inches(1.95),
                     body, font_size=11.5, color=BODY_GRAY, line_spacing=1.5)

    add_text_box(slide, Inches(0.8), Inches(6.1), Inches(11.0), Inches(0.35),
                 "Source: Wang et al., Archives of General Psychiatry (2011). "
                 "9.8-year median delay for mood disorders.",
                 font_size=10, color=MEDIUM_GRAY)

    add_slide_number(slide, 2)
    add_notes(slide,
        "Cite Wang et al. (2011), Archives of General Psychiatry for the 9.8-year figure. "
        "This is specifically for mood disorders; the median for any mental disorder is 11 years. "
        "\n\nEmphasize this is a systemic failure, not personal. The PHQ-9 (created 1999) is genuinely "
        "still the clinical gold standard for depression screening \u2014 a 25-year-old self-report "
        "questionnaire. No objective, continuous signal has ever been available to consumers. "
        "\n\nKey reframe: the problem is not lack of awareness. It is lack of DATA. "
        "Lucid doesn't replace therapy; it provides the signal that tells people when to seek it.")


# ── Slide 03 — The Market ────────────────────────────────────────

def slide_03_market(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, LIGHT_BG)
    add_section_header(slide, "THE MARKET")

    add_text_box(slide, Inches(0.8), Inches(1.1), Inches(11.5), Inches(0.75),
                 "The Quantified Self Has a Qualitative Gap.",
                 font_size=34, color=CHARCOAL, bold=True, font_name=HEADLINE_FONT)

    add_text_box(slide, Inches(0.8), Inches(1.9), Inches(11.5), Inches(0.4),
                 "$480B US wellness market, growing ~10%/year (McKinsey 2024). "
                 "Every layer tracked except one.",
                 font_size=14, color=BODY_GRAY)

    # 2x2 stat grid on the left
    stats = [
        ("$5,321/yr",  "Average American wellness spend\n(Fortune Well / GWI, Feb 2024)",  STEEL),
        ("52%",        "Of tech workers have depression\nor anxiety (Talkspace Business)", ACCENT_CORAL),
        ("$400-600/yr","Oura/WHOOP users spend on\nhealth data (Earnest Analytics 2024)", ACCENT_GREEN),
        ("$14.99/mo",  "Lucid \u2014 within established\nspend envelope for QS users",    ACCENT_GOLD),
    ]

    card_w = Inches(2.35)
    card_h = Inches(1.25)
    h_gap  = Inches(0.25)
    v_gap  = Inches(0.2)
    col_x  = [Inches(0.8), Inches(0.8) + card_w + h_gap]
    row_y  = [Inches(2.55), Inches(2.55) + card_h + v_gap]

    for i, (number, label, accent) in enumerate(stats):
        cx = col_x[i % 2]
        cy = row_y[i // 2]
        add_card_box(slide, cx, cy, card_w, card_h)
        add_dot(slide, cx + Inches(0.2), cy + Inches(0.15), Inches(0.1), accent)
        add_text_box(slide, cx + Inches(0.42), cy + Inches(0.05),
                     card_w - Inches(0.55), Inches(0.45),
                     number, font_size=22, color=accent, bold=True,
                     font_name=HEADLINE_FONT)
        add_text_box(slide, cx + Inches(0.42), cy + Inches(0.55),
                     card_w - Inches(0.55), Inches(0.55),
                     label, font_size=10.5, color=BODY_GRAY, line_spacing=1.3)

    # Right side — QS user profile (no emojis)
    right_x = Inches(6.1)
    right_w = Inches(6.8)

    add_text_box(slide, right_x, Inches(2.55), right_w, Inches(0.4),
                 "Who Is the QS User?", font_size=17, color=CHARCOAL,
                 bold=True, font_name=HEADLINE_FONT)

    traits = [
        "Already wears Oura, WHOOP, or Apple Watch",
        "Reviews HRV, sleep scores, and recovery daily",
        "Spends $400\u2013800/year on health data",
        "Frustrated their stack has no mental health layer",
        "Privacy-conscious \u2014 won't use cloud-only solutions",
        "Mac power user \u2014 macOS early adopter",
    ]
    for i, trait in enumerate(traits):
        y = Inches(3.1) + i * Inches(0.48)
        add_dot(slide, right_x, y + Inches(0.1), Inches(0.08), STEEL)
        add_text_box(slide, right_x + Inches(0.2), y,
                     right_w - Inches(0.2), Inches(0.42),
                     trait, font_size=13, color=BODY_GRAY)

    # Bottom callout bar
    cb_y = Inches(6.35)
    cb = add_card_box(slide, Inches(0.8), cb_y, Inches(11.7), Inches(0.58),
                      bg=NAVY, border=None)
    add_text_box(slide, Inches(1.1), cb_y + Inches(0.1), Inches(11.2), Inches(0.38),
                 "r/quantifiedself: 155K members  \u2022  r/Biohackers: 350K+  \u2022  "
                 "Oura/WHOOP users already paying $400+/yr for data  \u2022  Lucid is $14.99/mo",
                 font_size=12, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)

    add_slide_number(slide, 3)
    add_notes(slide,
        "These are already-paying customers. Oura users spent 52% more year-over-year in 2024 "
        "(Earnest Analytics). McKinsey 2024: $480B US wellness market, ~10%/yr growth. "
        "\n\nOur beachhead is the 2M QS early adopters who already spend $400-800/year on health data. "
        "Lucid is priced at $14.99/month \u2014 well inside their established spending envelope. "
        "\n\nKey insight: this customer does not need convincing that health data has value. "
        "They already believe it. We are selling the missing metric, not a new habit. "
        "\n\nGen Z + Millennials are 41% of wellness spend despite being 36% of adult population "
        "(McKinsey Future of Wellness 2024). This demographic over-indexes on both QS tools and "
        "mental health awareness.")


# ── Slide 04 — The Measurement Gap ──────────────────────────────

def slide_04_measurement_gap(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, LIGHT_BG)
    add_section_header(slide, "THE PROBLEM")

    add_text_box(slide, Inches(0.8), Inches(1.1), Inches(11.5), Inches(0.75),
                 "Every Layer Tracked. Except One.",
                 font_size=36, color=CHARCOAL, bold=True, font_name=HEADLINE_FONT)

    add_text_box(slide, Inches(0.8), Inches(1.9), Inches(11.5), Inches(0.5),
                 '"There is no available standard biomarker for the detection of '
                 'mental health conditions." \u2014 PMC survey on wearable sensors (2023)',
                 font_size=14, color=BODY_GRAY, line_spacing=1.4)

    # Wellness stack — 5 cards horizontal
    stack = [
        ("Oura",        "Sleep",  ACCENT_VIOLET, "Covered"),
        ("WHOOP",       "Strain", ACCENT_GREEN,  "Covered"),
        ("Function",    "Blood",  ACCENT_CORAL,  "Covered"),
        ("Apple Watch", "Heart",  RGBColor(0xC0, 0x30, 0x30), "Covered"),
        ("Lucid",       "Mind",   STEEL,          "Missing"),
    ]

    card_w = Inches(2.1)
    card_h = Inches(2.5)
    gap    = Inches(0.3)
    card_y = Inches(2.6)

    for i, (brand, metric, accent, status) in enumerate(stack):
        x = Inches(0.8) + i * (card_w + gap)
        is_lucid = (brand == "Lucid")

        bg_clr = RGBColor(0xEB, 0xF3, 0xFA) if is_lucid else CARD_BG
        bdr    = STEEL if is_lucid else CARD_BORDER
        bdr_w  = Pt(2) if is_lucid else Pt(1)
        add_card_box(slide, x, card_y, card_w, card_h,
                     bg=bg_clr, border=bdr, border_w=bdr_w)

        add_divider(slide, x + Inches(0.2), card_y + Inches(0.18),
                    Inches(0.6), accent, Pt(3))

        add_text_box(slide, x + Inches(0.2), card_y + Inches(0.55),
                     card_w - Inches(0.4), Inches(0.4),
                     brand, font_size=13, color=CHARCOAL, bold=True)

        metric_color = STEEL if is_lucid else BODY_GRAY
        add_text_box(slide, x + Inches(0.2), card_y + Inches(1.0),
                     card_w - Inches(0.4), Inches(0.35),
                     metric, font_size=16, color=metric_color,
                     bold=is_lucid, font_name=HEADLINE_FONT)

        if is_lucid:
            add_text_box(slide, x + Inches(0.2), card_y + Inches(1.5),
                         card_w - Inches(0.4), Inches(0.7),
                         "The missing\ndata stream", font_size=11,
                         color=STEEL, line_spacing=1.3)
        else:
            add_text_box(slide, x + Inches(0.2), card_y + Inches(1.5),
                         card_w - Inches(0.4), Inches(0.35),
                         "\u2713 " + status, font_size=11,
                         color=ACCENT_GREEN)

    # Bottom insight
    add_divider(slide, Inches(0.8), Inches(5.35),
                Inches(11.7), DIVIDER, Pt(1))
    add_text_box(slide, Inches(0.8), Inches(5.52), Inches(11.5), Inches(0.5),
                 "Voice is the only passive, hardware-free biomarker with direct "
                 "neurological links to mental state.",
                 font_size=15, color=CHARCOAL, bold=True, line_spacing=1.3)
    add_text_box(slide, Inches(0.8), Inches(6.08), Inches(11.5), Inches(0.4),
                 "No additional device. No new habit. Just speak for 25 seconds.",
                 font_size=14, color=BODY_GRAY)

    add_slide_number(slide, 4)
    add_notes(slide,
        "The 'every layer tracked except one' frame resonates with QS practitioners immediately \u2014 "
        "they live this frustration daily. They have dashboards for HRV, sleep stages, VO2 max, "
        "blood glucose, and HbA1c. They have nothing for mental state. "
        "\n\nThe PMC quote is from: Bent B et al. 'Investigating sources of inaccuracy in wearable "
        "optical heart rate sensors.' NPJ Digit Med. 2020; and the broader PMC review on biomarkers. "
        "\n\nThe vagal-laryngeal pathway is the physiological mechanism: the vagus nerve innervates "
        "both the heart (measured by HRV wearables) and the larynx (measured by voice). "
        "This is why voice changes are a genuine biomarker, not a correlational proxy. "
        "\n\nLucid is positioned as 'Mind' alongside Oura (Sleep), WHOOP (Strain), "
        "Function (Blood), Apple Watch (Heart). This framing naturalizes it as a wellness stack "
        "component, not a clinical intervention.")


# ── Slide 05 — What Has Failed ───────────────────────────────────

def slide_05_what_has_failed(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, LIGHT_BG)
    add_section_header(slide, "THE PROBLEM")

    add_text_box(slide, Inches(0.8), Inches(1.1), Inches(11.5), Inches(0.75),
                 "Everything That's Been Tried Has Missed the QS User.",
                 font_size=32, color=CHARCOAL, bold=True,
                 font_name=HEADLINE_FONT, line_spacing=1.1)

    # Table headers
    header_y = Inches(2.05)
    col_x = [Inches(0.8), Inches(3.4), Inches(8.5)]
    col_w = [Inches(2.5), Inches(5.0), Inches(4.0)]
    headers = ["Solution", "Why It Failed", "Key Data Point"]

    for j, (h, cx, cw) in enumerate(zip(headers, col_x, col_w)):
        add_text_box(slide, cx, header_y, cw, Inches(0.32),
                     h, font_size=10, color=MEDIUM_GRAY, bold=True)

    add_divider(slide, Inches(0.8), header_y + Inches(0.34),
                Inches(11.7), RGBColor(0xC8, 0xD0, 0xDA), Pt(1.5))

    # Rows
    rows = [
        ("Therapy",
         "Requires scheduling, travel, insurance,\nand willpower to sustain.",
         "34.8% dropout  \u2022  $140/session avg  \u2022  48-day wait",
         False),
        ("Calm / Headspace",
         "Active effort required every day.\nNeeds a new habit; most users quit.",
         "3.3% median 30-day retention\nCalm: $596M revenue, declining",
         False),
        ("Woebot (AI chatbot)",
         "Engagement still requires daily effort.\nConsumer product discontinued.",
         "$114M raised  \u2022  Shut down June 2025",
         True),
        ("Kintsugi Health (voice AI)",
         "B2B-only distribution; FDA regulatory path\nkilled the economics, not the science.",
         "Only voice biomarker competitor  \u2022  Shut down February 2026",
         True),
        ("Corporate EAPs",
         "Passive availability does not equal utilization.\nZero data feedback loop.",
         "4\u20136% utilization despite being free\n$68.4B market, largely wasted",
         False),
    ]

    row_y = Inches(2.52)
    row_h = Inches(0.72)

    for i, (solution, failure, data, is_dead) in enumerate(rows):
        y = row_y + i * row_h

        if is_dead:
            bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                        Inches(0.75), y - Inches(0.04),
                                        Inches(11.8), row_h)
            bg.fill.solid()
            bg.fill.fore_color.rgb = RGBColor(0xFF, 0xF2, 0xF2)
            bg.line.fill.background()

        add_text_box(slide, col_x[0], y + Inches(0.08), col_w[0], row_h - Inches(0.1),
                     solution, font_size=12, color=CHARCOAL, bold=True,
                     line_spacing=1.2)
        add_text_box(slide, col_x[1], y + Inches(0.06), col_w[1], row_h - Inches(0.06),
                     failure, font_size=11, color=BODY_GRAY, line_spacing=1.3)
        data_color = ACCENT_CORAL if is_dead else BODY_GRAY
        add_text_box(slide, col_x[2], y + Inches(0.06), col_w[2], row_h - Inches(0.06),
                     data, font_size=11, color=data_color, bold=is_dead,
                     line_spacing=1.3)

        add_divider(slide, Inches(0.8), y + row_h - Inches(0.05),
                    Inches(11.7), DIVIDER, Pt(0.5))

    # Insight bar
    ib_y = Inches(6.35)
    add_card_box(slide, Inches(0.8), ib_y, Inches(11.7), Inches(0.58),
                 bg=NAVY, border=None)
    add_text_box(slide, Inches(1.1), ib_y + Inches(0.1), Inches(11.2), Inches(0.38),
                 "QS users don\u2019t want another chore. They want another data stream. "
                 "Lucid is the only passive mental health monitor.",
                 font_size=12, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER,
                 line_spacing=1.2)

    add_slide_number(slide, 5)
    add_notes(slide,
        "The Kintsugi and Woebot shutdowns are our moat, not a risk signal. "
        "Both failed for structural reasons that Lucid avoids by design: "
        "\n\n- Kintsugi: raised $72M, built the best voice biomarker AI in existence, "
        "then shut down in February 2026 because the B2B/FDA regulatory pathway was "
        "uneconomical. The science was sound; the go-to-market was wrong. "
        "\n\n- Woebot: $114M raised, chat-based AI therapy. Shut down consumer product "
        "June 2025 because engagement requires active effort \u2014 the fundamental flaw "
        "of every mental health app. "
        "\n\nLucid avoids both: DTC consumer distribution (no FDA pathway needed for a "
        "wellness monitoring tool) and passive architecture (no daily effort required). "
        "\n\nThe 3.3% median 30-day retention figure is from Baumel et al. (2019), "
        "'Official Mental Health Apps in the World's Biggest App Stores.' This is the "
        "definitive citation for mental health app retention failure.")


# ── Slide 06 — The Solution ──────────────────────────────────────

def slide_06_solution(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, NAVY)

    add_text_box(slide, Inches(0.8), Inches(0.5), Inches(4.0), Inches(0.35),
                 "THE SOLUTION", font_size=11, color=STEEL, bold=True)
    add_divider(slide, Inches(0.8), Inches(0.9), Inches(2.0), STEEL, Pt(2))

    add_text_box(slide, Inches(0.8), Inches(1.2), Inches(11.5), Inches(0.85),
                 "$30 Million in Research. $14.99 to Access It.",
                 font_size=38, color=WHITE, bold=True, font_name=HEADLINE_FONT)

    add_text_box(slide, Inches(0.8), Inches(2.15), Inches(11.5), Inches(0.65),
                 "The full 244M-parameter clinical model \u2014 validated on 14,898 adults \u2014 "
                 "now runs on your Mac for less than one therapy copay.",
                 font_size=15, color=RGBColor(0xAA, 0xBB, 0xCC), line_spacing=1.4)

    # 3 pillar cards
    pillars = [
        ("Clinically Validated",
         "Annals of Family Medicine, January 2025.\n"
         "14,898 patients. 71.3% sensitivity for MDD.\n"
         "Peer-reviewed and PubMed-indexed."),
        ("100% On-Device",
         "Voice never leaves your Mac.\n"
         "Runs on Apple Silicon Neural Engine.\n"
         "Verifiable: works while disconnected."),
        ("Passive",
         "25 seconds of natural speech.\n"
         "No journals. No questionnaires.\n"
         "Monitors passively during calls and meetings."),
    ]

    card_w = Inches(3.7)
    card_h = Inches(2.55)
    gap    = Inches(0.4)
    card_y = Inches(3.1)

    for i, (title, body) in enumerate(pillars):
        x = Inches(0.8) + i * (card_w + gap)
        add_card_box(slide, x, card_y, card_w, card_h,
                     bg=RGBColor(0x22, 0x44, 0x5C),
                     border=RGBColor(0x33, 0x55, 0x6E))
        add_text_box(slide, x + Inches(0.3), card_y + Inches(0.25),
                     card_w - Inches(0.6), Inches(0.42),
                     title, font_size=15, color=WHITE, bold=True)
        add_text_box(slide, x + Inches(0.3), card_y + Inches(0.8),
                     card_w - Inches(0.6), Inches(1.55),
                     body, font_size=12, color=RGBColor(0xAA, 0xBB, 0xCC),
                     line_spacing=1.5)

    add_slide_number(slide, 6, color=RGBColor(0x66, 0x77, 0x88))
    add_notes(slide,
        "The $30M figure is Kintsugi's total funding raised per Crunchbase. "
        "The R&D work is done and published. We are not paying for science; "
        "we are paying for consumer distribution \u2014 a fundamentally different "
        "and lower-risk investment thesis. "
        "\n\nAt $14.99/month, Lucid costs less than one therapy copay (avg $140/session) "
        "and delivers daily data rather than weekly check-ins. "
        "\n\nThe 'three pillars' framing maps directly to the three objections we hear: "
        "(1) Is this real? \u2192 Clinically Validated. "
        "(2) What happens to my data? \u2192 100% On-Device. "
        "(3) Will I actually use this? \u2192 Passive. "
        "\n\nNote: 'wellness monitoring tool' is the correct regulatory framing. "
        "Lucid does not diagnose, treat, or prescribe. It monitors and reports trends.")


# ── Slide 07 — How It Works ──────────────────────────────────────

def slide_07_how_it_works(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, LIGHT_BG)
    add_section_header(slide, "HOW IT WORKS")

    add_text_box(slide, Inches(0.8), Inches(1.1), Inches(11.5), Inches(0.75),
                 "What Used to Require a Clinical Referral Now Takes 5 Minutes.",
                 font_size=30, color=CHARCOAL, bold=True,
                 font_name=HEADLINE_FONT, line_spacing=1.1)

    # 3 step cards — evenly spaced, no competing image
    steps = [
        ("Step 1",
         "Speak Naturally",
         "Lucid monitors your microphone passively during calls "
         "and meetings. Or open the app and speak for 25 seconds. "
         "The AI hears how you sound \u2014 not what you say.",
         STEEL),
        ("Step 2",
         "244M-Parameter Model",
         "Extracts 20+ acoustic biomarkers in real time. "
         "Runs entirely on Apple Silicon Neural Engine. "
         "Zero cloud upload. Zero latency. Full privacy.",
         ACCENT_GREEN),
        ("Step 3",
         "See What the System Misses",
         "Daily wellness zones: Calm, Steady, Tense, Stressed. "
         "7-day and 30-day trend analysis with pattern detection. "
         "Early warning signals before burnout arrives.",
         ACCENT_GOLD),
    ]

    card_w = Inches(3.8)
    card_h = Inches(4.1)
    gap    = Inches(0.37)
    card_y = Inches(2.05)

    for i, (step_label, title, body, accent) in enumerate(steps):
        x = Inches(0.75) + i * (card_w + gap)
        add_card_box(slide, x, card_y, card_w, card_h)

        # Step label pill
        pill = add_card_box(slide, x + Inches(0.3), card_y + Inches(0.25),
                            Inches(0.9), Inches(0.32),
                            bg=accent, border=None, radius=0.3)
        add_text_box(slide, x + Inches(0.3), card_y + Inches(0.25),
                     Inches(0.9), Inches(0.32),
                     step_label, font_size=10, color=WHITE, bold=True,
                     alignment=PP_ALIGN.CENTER,
                     anchor=MSO_ANCHOR.MIDDLE)

        add_text_box(slide, x + Inches(0.3), card_y + Inches(0.75),
                     card_w - Inches(0.6), Inches(0.5),
                     title, font_size=17, color=CHARCOAL, bold=True,
                     font_name=HEADLINE_FONT)

        add_divider(slide, x + Inches(0.3), card_y + Inches(1.35),
                    card_w - Inches(0.6), accent, Pt(1.5))

        add_text_box(slide, x + Inches(0.3), card_y + Inches(1.55),
                     card_w - Inches(0.6), Inches(2.35),
                     body, font_size=12.5, color=BODY_GRAY, line_spacing=1.55)

    add_slide_number(slide, 7)
    add_notes(slide,
        "Key message: this is passive. Lucid listens during calls and meetings \u2014 "
        "no new habit required. The model runs on the Neural Engine in every M-series Mac: "
        "no cloud, no latency, zero marginal cost per inference. "
        "\n\n25 seconds is the clinically validated minimum for accurate detection "
        "(per the Kintsugi DAM study methodology). The model analyzes jitter, shimmer, "
        "harmonics-to-noise ratio, fundamental frequency (F0), MFCCs, and energy. "
        "\n\nPassive monitoring is the key architectural advantage over every competitor. "
        "Every failed mental health app requires the user to open the app. "
        "Lucid works whether or not the user remembers to open it. "
        "\n\nStep 2 emphasis: Apple Silicon's Neural Engine performs ~38 trillion operations "
        "per second on M-series chips. Running a 244M parameter model that previously required "
        "AWS GPU clusters is now feasible at zero marginal cost on consumer hardware. "
        "This is the structural change that makes Lucid possible.")


# ── Slide 08 — Clinical Validation ──────────────────────────────

def slide_08_clinical_validation(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, NAVY)

    add_text_box(slide, Inches(0.8), Inches(0.5), Inches(5.0), Inches(0.35),
                 "CLINICAL VALIDATION", font_size=11, color=STEEL, bold=True)
    add_divider(slide, Inches(0.8), Inches(0.9), Inches(2.5), STEEL, Pt(2))

    add_text_box(slide, Inches(0.8), Inches(1.2), Inches(11.5), Inches(0.85),
                 "Peer-Reviewed. 14,898 Patients. "
                 "Published in Annals of Family Medicine.",
                 font_size=30, color=WHITE, bold=True,
                 font_name=HEADLINE_FONT, line_spacing=1.1)

    # 4 stat boxes
    stats = [
        ("71.3%",  "Sensitivity\nfor MDD",      STEEL),
        ("25 sec", "Of speech\nneeded",          ACCENT_GREEN),
        ("35,000", "Patients in\ntraining set",  ACCENT_GOLD),
        ("244M",   "Model\nparameters",          ACCENT_VIOLET),
    ]

    box_w = Inches(2.6)
    box_h = Inches(2.1)
    gap   = Inches(0.4)
    box_y = Inches(2.6)

    for i, (num, label, accent) in enumerate(stats):
        x = Inches(0.8) + i * (box_w + gap)
        add_card_box(slide, x, box_y, box_w, box_h,
                     bg=RGBColor(0x22, 0x44, 0x5C),
                     border=RGBColor(0x33, 0x55, 0x6E))
        add_text_box(slide, x, box_y + Inches(0.2), box_w, Inches(0.7),
                     num, font_size=44, color=accent, bold=True,
                     alignment=PP_ALIGN.CENTER, font_name=HEADLINE_FONT)
        add_text_box(slide, x, box_y + Inches(1.0), box_w, Inches(0.85),
                     label, font_size=13, color=RGBColor(0xAA, 0xBB, 0xCC),
                     alignment=PP_ALIGN.CENTER, line_spacing=1.3)

    # PubMed image
    pub_w, pub_h = fit_image_dims("pubmed", 3.5, 3.2)
    pub_left = Inches(9.0)
    pub_top  = Inches(2.0)
    if os.path.exists(IMG_PUBMED):
        frame = add_card_box(slide,
                             pub_left - Inches(0.08), pub_top - Inches(0.08),
                             pub_w + Inches(0.16), pub_h + Inches(0.16),
                             bg=RGBColor(0x1A, 0x2C, 0x3C), border=None, radius=0.02)
        slide.shapes.add_picture(IMG_PUBMED, pub_left, pub_top, pub_w, pub_h)

    cite_top = pub_top + pub_h + Inches(0.15)
    add_text_box(slide, pub_left - Inches(0.4), cite_top,
                 pub_w + Inches(0.8), Inches(0.5),
                 "Mazur et al. (2025) Annals of Family Medicine\n"
                 "DOI: 10.1370/afm.240091",
                 font_size=10, color=RGBColor(0x88, 0x99, 0xAA),
                 alignment=PP_ALIGN.CENTER, line_spacing=1.3)

    add_text_box(slide, Inches(0.8), Inches(5.0), Inches(7.8), Inches(0.55),
                 "PHQ-9 has comparable sensitivity \u2014 but it is given once per year. "
                 "Lucid monitors every day.",
                 font_size=14, color=RGBColor(0xAA, 0xBB, 0xCC),
                 bold=True, line_spacing=1.3)

    add_slide_number(slide, 8, color=RGBColor(0x66, 0x77, 0x88))
    add_notes(slide,
        "Full citation: Mazur N et al. 'Voice Biomarker Analysis for Depression Screening in "
        "Primary Care.' Annals of Family Medicine. January 2025. DOI: 10.1370/afm.240091. "
        "\n\nThis is the Depression Assessment Model (DAM) published by the Kintsugi research team. "
        "Key performance metrics: 71.3% sensitivity, 73.5% specificity for MDD. "
        "Training set: 35,000 patients. External validation set: 14,898 adults. "
        "Multi-site, demographically diverse. "
        "\n\nFor context on the PHQ-9: meta-analyses show 77-88% sensitivity at standard cutoffs "
        "(Levis et al., JAMA Internal Medicine, 2019) \u2014 comparable performance, "
        "but the PHQ-9 requires a clinician to administer it once per year. "
        "Lucid delivers an objective acoustic measurement every single day. "
        "\n\nFDA note: Lucid operates as a wellness monitoring tool, not a diagnostic device. "
        "The DAM publication explicitly frames voice biomarkers as 'screening aids,' "
        "not clinical diagnosis \u2014 the same regulatory framing as Fitbit's heart health features.")


# ── Slide 09 — Technology ────────────────────────────────────────

def slide_09_technology(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, LIGHT_BG)
    add_section_header(slide, "THE TECHNOLOGY")

    add_text_box(slide, Inches(0.8), Inches(1.1), Inches(11.5), Inches(0.75),
                 "Apple Silicon Changed the Economics.",
                 font_size=36, color=CHARCOAL, bold=True, font_name=HEADLINE_FONT)

    add_text_box(slide, Inches(0.8), Inches(1.9), Inches(11.5), Inches(0.5),
                 "A 244M-parameter clinical model that once required a data center "
                 "now runs free on every modern Mac.",
                 font_size=14, color=BODY_GRAY, line_spacing=1.4)

    # Left: local stack
    left_x = Inches(0.8)
    left_w = Inches(5.8)

    add_text_box(slide, left_x, Inches(2.6), left_w, Inches(0.38),
                 "What Runs Locally", font_size=16, color=CHARCOAL,
                 bold=True, font_name=HEADLINE_FONT)

    local_stack = [
        ("Kintsugi DAM",      "244M parameters  \u2022  Annals of Family Medicine validated",  STEEL),
        ("Silero VAD",         "Voice activity detection  \u2022  filters non-speech",           ACCENT_GREEN),
        ("ECAPA-TDNN",         "Speaker verification  \u2022  PLDA scoring gate",                ACCENT_GOLD),
        ("Acoustic Features",  "20+ features: jitter, shimmer, HNR, MFCC, F0, energy",          ACCENT_VIOLET),
    ]

    for i, (name, desc, accent) in enumerate(local_stack):
        y = Inches(3.1) + i * Inches(0.82)
        add_card_box(slide, left_x, y, left_w, Inches(0.7))
        add_dot(slide, left_x + Inches(0.2), y + Inches(0.13), Inches(0.1), accent)
        add_text_box(slide, left_x + Inches(0.44), y + Inches(0.04),
                     Inches(2.0), Inches(0.32),
                     name, font_size=13, color=CHARCOAL, bold=True)
        add_text_box(slide, left_x + Inches(0.44), y + Inches(0.38),
                     left_w - Inches(0.65), Inches(0.28),
                     desc, font_size=11, color=BODY_GRAY)

    # Right: privacy + margin
    right_x = Inches(7.3)
    right_w = Inches(5.5)

    priv_y = Inches(2.6)
    add_card_box(slide, right_x, priv_y, right_w, Inches(1.5),
                 bg=NAVY, border=None)
    add_text_box(slide, right_x + Inches(0.3), priv_y + Inches(0.18),
                 right_w - Inches(0.6), Inches(0.38),
                 "Privacy Proof", font_size=16, color=WHITE, bold=True)
    add_text_box(slide, right_x + Inches(0.3), priv_y + Inches(0.63),
                 right_w - Inches(0.6), Inches(0.72),
                 "Works while disconnected. Verifiable by turning off WiFi. "
                 "Mental health data is the most sensitive data category. "
                 "Lucid is the only solution that cannot leak it.",
                 font_size=11.5, color=RGBColor(0xAA, 0xBB, 0xCC), line_spacing=1.4)

    margin_y = Inches(4.3)
    add_card_box(slide, right_x, margin_y, right_w, Inches(1.5),
                 bg=RGBColor(0xEB, 0xF3, 0xFA),
                 border=STEEL, border_w=Pt(1.5))
    add_text_box(slide, right_x + Inches(0.3), margin_y + Inches(0.18),
                 right_w - Inches(0.6), Inches(0.42),
                 "82\u201392% Gross Margin", font_size=20, color=STEEL,
                 bold=True, font_name=HEADLINE_FONT)
    add_text_box(slide, right_x + Inches(0.3), margin_y + Inches(0.65),
                 right_w - Inches(0.6), Inches(0.7),
                 "Zero compute COGS \u2014 inference runs on user hardware.\n"
                 "App Store: ~82%  \u2022  Direct/API: ~92%\n"
                 "Closest to pure software economics in health AI.",
                 font_size=11.5, color=BODY_GRAY, line_spacing=1.4)

    add_slide_number(slide, 9)
    add_notes(slide,
        "Apple Silicon's Neural Engine: ~38 trillion operations/sec on M2, "
        "~18 TOPS on M1. Running a 244M parameter transformer model locally is "
        "feasible because the DAM uses quantized inference optimized for the ANE. "
        "\n\nKintsugi required AWS GPU clusters to run this model in 2021. "
        "Apple M-series chips made local inference viable by 2023. "
        "This is the structural technology shift that creates the business opportunity. "
        "\n\nThe privacy architecture is a genuine competitive moat in health AI: "
        "HIPAA concerns are eliminated because no data leaves the device. "
        "This is especially relevant for corporate wellness and insurance use cases "
        "(Phase 2) where data liability is a primary objection. "
        "\n\nGross margin breakdown: Apple charges 15% (first $1M ARR) or 30% on "
        "App Store. Direct web sales and API licensing carry near-zero COGS. "
        "Blended target: 82-88% gross margin at scale.")


# ── Slide 10 — Product Deep Dive ─────────────────────────────────

def slide_10_product(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, LIGHT_BG)
    add_section_header(slide, "PRODUCT")

    add_text_box(slide, Inches(0.8), Inches(1.1), Inches(11.5), Inches(0.75),
                 "A Monitoring App, Not a Treatment App.",
                 font_size=34, color=CHARCOAL, bold=True, font_name=HEADLINE_FONT)

    add_text_box(slide, Inches(0.8), Inches(1.9), Inches(11.5), Inches(0.4),
                 "Monitoring apps (Oura, Apple Watch): 88%+ 12-month retention  \u2022  "
                 "Mental health apps: 3.3% 30-day retention",
                 font_size=14, color=BODY_GRAY)

    features = [
        ("Daily Detection",
         "Zone classification: Calm / Steady / Tense / Stressed. "
         "Composite wellness score plus 8 acoustic biomarkers."),
        ("Early Warning",
         "7-day and 30-day trend analysis. "
         "Pattern detection and burnout calculator."),
        ("Grove (Streak Forest)",
         "Streak-based forest grows with daily check-ins. "
         "8x retention for streak users \u2014 Duolingo playbook."),
        ("Rhythm Rings",
         "Consistency visualization inspired by Apple Watch rings. "
         "Activates loss aversion for daily use."),
        ("Waypoints (30-Tier Trail)",
         "Progress-based system unlocking deeper insights. "
         "Keeps long-term users engaged."),
        ("Raw Data + Export",
         "8 acoustic biomarkers, CSV/JSON export. "
         "REST API for integration with existing QS stack."),
    ]

    card_w = Inches(2.95)
    card_h = Inches(1.28)
    h_gap  = Inches(0.25)
    v_gap  = Inches(0.2)
    grid_x = Inches(0.8)
    grid_y = Inches(2.55)

    for i, (title, body) in enumerate(features):
        row = i // 2
        col = i % 2
        x = grid_x + col * (card_w + h_gap)
        y = grid_y + row * (card_h + v_gap)
        add_card_box(slide, x, y, card_w, card_h)
        add_text_box(slide, x + Inches(0.25), y + Inches(0.18),
                     card_w - Inches(0.5), Inches(0.38),
                     title, font_size=13, color=CHARCOAL, bold=True)
        add_text_box(slide, x + Inches(0.25), y + Inches(0.58),
                     card_w - Inches(0.5), card_h - Inches(0.7),
                     body, font_size=11, color=BODY_GRAY, line_spacing=1.3)

    # Dashboard screenshot
    dash_w, dash_h = fit_image_dims("dashboard", 5.4, 4.3)
    dash_left = Inches(7.0)
    dash_top  = Inches(2.0)
    if os.path.exists(IMG_DASHBOARD):
        frame = add_card_box(slide,
                             dash_left - Inches(0.05), dash_top - Inches(0.05),
                             dash_w + Inches(0.1), dash_h + Inches(0.1),
                             bg=RGBColor(0xE0, 0xE0, 0xE8), border=None, radius=0.03)
        slide.shapes.add_picture(IMG_DASHBOARD, dash_left, dash_top, dash_w, dash_h)

    add_slide_number(slide, 10)
    add_notes(slide,
        "The 'monitoring vs treatment' framing is the core retention insight. "
        "Oura and Apple Watch achieve 88%+ 12-month retention because they are passive \u2014 "
        "the user does not need to remember to use them. Lucid is designed identically. "
        "\n\nEngagement mechanics explained: "
        "\n- Grove (streak forest): Each consecutive day of use grows a virtual tree. "
        "Research shows streak users (Duolingo, Headspace) show 8x 30-day retention vs. "
        "non-streak users. The Grove system makes missing a day visually costly. "
        "\n- Rhythm Rings: Borrowed from Apple Watch's Activity rings. "
        "Closure compulsion drives daily check-ins without explicit reminders. "
        "\n- Waypoints: 30-tier exploration metaphor. Users 'discover' insights as they "
        "progress. Progress visualization activates the Zeigarnik effect "
        "(incompleteness drives completion). "
        "\n\nAPI/export: This is specifically for the QS user persona. "
        "r/quantifiedself users routinely export data to Notion, Obsidian, "
        "and custom dashboards. Raw data access is a purchase driver for this segment.")


# ── Slide 11 — Market Size ───────────────────────────────────────

def slide_11_market_size(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, LIGHT_BG)
    add_section_header(slide, "MARKET")

    add_text_box(slide, Inches(0.8), Inches(1.1), Inches(11.5), Inches(0.75),
                 "A $480B Wellness Market With a Mental Health Gap.",
                 font_size=34, color=CHARCOAL, bold=True, font_name=HEADLINE_FONT)

    # TAM / SAM / SOM horizontal bars
    bars = [
        ("TAM", "$480B",  "US Wellness Market",
         Inches(11.0), RGBColor(0xD8, 0xE8, 0xF4), CHARCOAL,
         "Therapy, pharma, digital therapeutics, workplace wellness, supplements. "
         "Growing ~10%/year. (McKinsey 2024)"),
        ("SAM", "$68B",   "Digital Mental Health + QS Stack",
         Inches(7.0),  RGBColor(0xB8, 0xD4, 0xE8), CHARCOAL,
         "Apps, platforms, wearables with mental health features. "
         "EAP market $68.4B. Growing ~18% CAGR."),
        ("SOM", "$2.1B",  "Consumer Voice Wellness \u2014 QS Early Adopters",
         Inches(3.5),  STEEL, WHITE,
         "2M QS practitioners spending $400+/yr. "
         "Oura/WHOOP overlap. 155K r/quantifiedself."),
    ]

    bar_x     = Inches(0.8)
    bar_y     = Inches(2.3)
    bar_h     = Inches(1.1)
    bar_gap   = Inches(0.28)

    for i, (tier, amount, name, bar_w, bg_c, text_c, desc) in enumerate(bars):
        y = bar_y + i * (bar_h + bar_gap)
        bar = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                     bar_x, y, bar_w, bar_h)
        bar.fill.solid()
        bar.fill.fore_color.rgb = bg_c
        bar.line.fill.background()
        bar.adjustments[0] = 0.1

        add_text_box(slide, bar_x + Inches(0.3), y + Inches(0.08),
                     Inches(0.7), Inches(0.3),
                     tier, font_size=11, color=text_c, bold=True)
        add_text_box(slide, bar_x + Inches(1.05), y + Inches(0.03),
                     Inches(1.7), Inches(0.4),
                     amount, font_size=26, color=text_c, bold=True,
                     font_name=HEADLINE_FONT)
        add_text_box(slide, bar_x + Inches(0.3), y + Inches(0.48),
                     Inches(2.6), Inches(0.3),
                     name, font_size=11, color=text_c, bold=True)

        desc_x = bar_x + Inches(3.5)
        if bar_w > Inches(5.0):
            add_text_box(slide, desc_x, y + Inches(0.12),
                         bar_w - Inches(4.0), Inches(0.85),
                         desc, font_size=11, color=BODY_GRAY, line_spacing=1.3)
        else:
            add_text_box(slide, desc_x, y + Inches(0.12),
                         Inches(8.0), Inches(0.85),
                         desc, font_size=11, color=BODY_GRAY, line_spacing=1.3)

    add_divider(slide, Inches(0.8), Inches(6.35), Inches(11.7), DIVIDER, Pt(1))
    add_text_box(slide, Inches(0.8), Inches(6.52), Inches(11.5), Inches(0.45),
                 "GTM wedge: macOS power users, biohackers, QS practitioners. "
                 "They already track everything \u2014 Lucid adds the mental layer.",
                 font_size=13, color=CHARCOAL, bold=True, line_spacing=1.3)

    add_slide_number(slide, 11)
    add_notes(slide,
        "TAM framing note: we deliberately use the US wellness market ($480B, McKinsey 2024) "
        "rather than the global mental health market ($588B) or total behavioral health market "
        "to keep the sizing credible and investor-friendly. "
        "\n\nSOM calculation: approximately 2M Americans qualify as QS early adopters "
        "(defined as: own a health wearable, track metrics weekly, spend $400+/yr on health data). "
        "At $14.99/month and 10% penetration = 200K subscribers = $36M ARR. "
        "That alone justifies a $360M valuation at 10x revenue. "
        "\n\nKey insight for investors: this is a wedge market with a clear expansion path. "
        "Start with 2M QS early adopters (highest willingness to pay, easiest to reach). "
        "Expand to 40M Americans with anxiety/depression (Phase 2 with B2B/insurance). "
        "\n\nThe EAP market ($68.4B) with 4-6% utilization represents a massive unlocked SAM. "
        "Even moving utilization from 5% to 15% in corporate wellness represents a "
        "$4B+ annual recurring opportunity.")


# ── Slide 12 — Business Model ────────────────────────────────────

def slide_12_business_model(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, LIGHT_BG)
    add_section_header(slide, "BUSINESS MODEL")

    add_text_box(slide, Inches(0.8), Inches(1.1), Inches(11.5), Inches(0.75),
                 "82% Gross Margin. Zero Compute COGS.",
                 font_size=34, color=CHARCOAL, bold=True, font_name=HEADLINE_FONT)

    # LEFT: pricing card
    px, py, pw, ph = Inches(0.8), Inches(2.05), Inches(4.5), Inches(4.55)
    add_card_box(slide, px, py, pw, ph, bg=NAVY, border=None)

    add_text_box(slide, px + Inches(0.4), py + Inches(0.28),
                 Inches(3.5), Inches(0.35),
                 "Lucid Premium", font_size=13, color=STEEL, bold=True)
    add_text_box(slide, px + Inches(0.4), py + Inches(0.65),
                 Inches(3.5), Inches(0.58),
                 "$14.99/mo", font_size=40, color=WHITE, bold=True,
                 font_name=HEADLINE_FONT)
    add_text_box(slide, px + Inches(0.4), py + Inches(1.22),
                 Inches(3.5), Inches(0.3),
                 "or $149/year (save $30)", font_size=12,
                 color=RGBColor(0xAA, 0xBB, 0xCC))
    add_text_box(slide, px + Inches(0.4), py + Inches(1.55),
                 Inches(3.5), Inches(0.28),
                 "30-day free trial  \u2022  no credit card required",
                 font_size=11, color=STEEL)

    feats = [
        "\u2713  Unlimited voice analyses",
        "\u2713  Clinical wellness scoring (8 biomarkers)",
        "\u2713  Depression risk screening",
        "\u2713  Stress, mood and energy tracking",
        "\u2713  Trend analytics + burnout calculator",
        "\u2713  CSV/JSON export + REST API",
        "\u2713  100% on-device processing",
    ]
    for j, f in enumerate(feats):
        add_text_box(slide, px + Inches(0.4), py + Inches(1.95 + j * 0.33),
                     Inches(3.5), Inches(0.3),
                     f, font_size=11.5, color=RGBColor(0xCC, 0xDD, 0xEE))

    # RIGHT: unit economics
    rx = Inches(6.0)
    add_text_box(slide, rx, Inches(2.05), Inches(6.8), Inches(0.38),
                 "Unit Economics", font_size=17, color=CHARCOAL,
                 bold=True, font_name=HEADLINE_FONT)

    rows = [
        ("Blended ARPU",    "~$12.42/mo (mix of monthly + annual)"),
        ("Gross Margin",    "82% (App Store)  \u2192  92% (direct)"),
        ("Target LTV",      "$149 at 44% 12-month retention"),
        ("Target CAC",      "$60 blended"),
        ("LTV / CAC",       "2.5x target"),
        ("Breakeven",       "2,034 subscribers\n(3-FTE team, $25K/mo burn)"),
    ]
    for i, (label, val) in enumerate(rows):
        y = Inches(2.55) + i * Inches(0.52)
        add_text_box(slide, rx, y, Inches(2.8), Inches(0.42),
                     label, font_size=12, color=BODY_GRAY)
        add_text_box(slide, rx + Inches(2.85), y, Inches(3.8), Inches(0.42),
                     val, font_size=12, color=CHARCOAL, bold=True,
                     line_spacing=1.2)
        add_divider(slide, rx, y + Inches(0.47),
                    Inches(6.8), DIVIDER, Pt(0.5))

    add_text_box(slide, rx, Inches(5.85), Inches(6.8), Inches(0.42),
                 "Phase 2: B2B corporate wellness at $2\u20135/PEPM "
                 "\u2014 unlocks enterprise ACVs of $50K+",
                 font_size=12, color=STEEL, bold=True, line_spacing=1.3)

    add_slide_number(slide, 12)
    add_notes(slide,
        "Breakeven calculation: 2,034 subscribers x $14.99/mo = $30,500/mo revenue. "
        "3-FTE team (founder + engineer + growth): ~$25K/mo burn + overhead = breakeven. "
        "This is achievable within 6 months of App Store launch. "
        "\n\nLTV calculation: $149 annual plan x 44% 12-month retention + "
        "churn revenue from monthly subscribers. Conservative LTV = $149. "
        "\n\nCAC breakdown by channel: "
        "- Apple Search Ads: $32-85 (highly intent-based, strong for health apps) "
        "- Instagram/Meta: $37 (mid-funnel, visual demos) "
        "- TikTok: $21-53 (top of funnel, educational content) "
        "- Organic/community: $0 (must be 40%+ of mix to hit blended $60 CAC) "
        "\n\nHealth & Fitness is the #1 app category by payer LTV (RevenueCat 2025): "
        "median $16.44/month, upper quartile $31.12/month. "
        "This validates that health app users pay and stay. "
        "\n\nPhase 2 B2B: at $3/PEPM with 1,000 employees = $36K ACV. "
        "25 corporate clients = $900K ARR purely from B2B. "
        "This layer compounds on top of the DTC subscriber base.")


# ── Slide 13 — Competitive Landscape ─────────────────────────────

def slide_13_competitive(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, LIGHT_BG)
    add_section_header(slide, "COMPETITIVE LANDSCAPE")

    add_text_box(slide, Inches(0.8), Inches(1.1), Inches(11.5), Inches(0.75),
                 "The Only B2B Competitor Just Shut Down. The Market Is Ours.",
                 font_size=30, color=CHARCOAL, bold=True,
                 font_name=HEADLINE_FONT, line_spacing=1.1)

    # Table header
    header_y = Inches(2.05)
    col_labels = ["", "Clinical AI", "On-Device", "Consumer", "Passive", "Price"]
    col_x = [Inches(0.8), Inches(4.05), Inches(5.45), Inches(6.85),
             Inches(8.25), Inches(9.65)]
    col_w = [Inches(3.15), Inches(1.2), Inches(1.2), Inches(1.2),
             Inches(1.2), Inches(2.5)]

    for j, (h, cx, cw) in enumerate(zip(col_labels, col_x, col_w)):
        add_text_box(slide, cx, header_y, cw, Inches(0.32),
                     h, font_size=10, color=MEDIUM_GRAY, bold=True,
                     alignment=PP_ALIGN.CENTER if j > 0 else PP_ALIGN.LEFT)

    add_divider(slide, Inches(0.8), header_y + Inches(0.34),
                Inches(11.7), RGBColor(0xC8, 0xD0, 0xDA), Pt(1.5))

    # competitor rows: name, clinical, on-device, consumer, passive, price, is_dead, note
    rows = [
        ("Lucid",              True,  True,  True,  True,  "$14.99/mo",    False, ""),
        ("Kintsugi Health",    True,  False, False, True,  "B2B only",     True,  "Shut down Feb 2026"),
        ("Sonde Health",       True,  False, False, False, "Enterprise",   False, "B2B only"),
        ("Woebot",             False, False, True,  False, "n/a",          True,  "Shut down Jun 2025"),
        ("Calm / Headspace",   False, False, True,  False, "$12\u201315/mo",False,"No clinical validation"),
        ("Corporate EAPs",     False, False, False, False, "Free (4% use)",False, "Employer-only"),
    ]

    row_y  = Inches(2.52)
    row_h  = Inches(0.64)
    CHECK  = "\u2713"
    DASH   = "\u2013"

    for i, (name, clinical, ondev, consumer, passive, price, is_dead, note) in enumerate(rows):
        y = row_y + i * row_h
        is_lucid = (name == "Lucid")

        if is_lucid:
            bg = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                        Inches(0.75), y - Inches(0.04),
                                        Inches(11.85), row_h)
            bg.fill.solid()
            bg.fill.fore_color.rgb = RGBColor(0xEB, 0xF3, 0xFA)
            bg.line.color.rgb = STEEL
            bg.line.width = Pt(1)
            bg.adjustments[0] = 0.06
        elif is_dead:
            bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                        Inches(0.75), y - Inches(0.04),
                                        Inches(11.85), row_h)
            bg.fill.solid()
            bg.fill.fore_color.rgb = RGBColor(0xFF, 0xF2, 0xF2)
            bg.line.fill.background()

        name_disp  = f"{name}  \u2014  {note}" if note else name
        name_color = CHARCOAL if is_lucid else (ACCENT_CORAL if is_dead else BODY_GRAY)
        add_text_box(slide, col_x[0], y + Inches(0.1), col_w[0], row_h - Inches(0.1),
                     name_disp, font_size=12 if is_lucid else 11,
                     color=name_color, bold=is_lucid, line_spacing=1.1)

        for j, val in enumerate([clinical, ondev, consumer, passive]):
            sym = CHECK if val else DASH
            sym_color = ACCENT_GREEN if val else MEDIUM_GRAY
            add_text_box(slide, col_x[j + 1], y + Inches(0.1),
                         col_w[j + 1], Inches(0.35),
                         sym, font_size=15, color=sym_color,
                         bold=val, alignment=PP_ALIGN.CENTER)

        add_text_box(slide, col_x[5], y + Inches(0.1), col_w[5], Inches(0.35),
                     price, font_size=11,
                     color=STEEL if is_lucid else BODY_GRAY,
                     bold=is_lucid, alignment=PP_ALIGN.CENTER)

        if i < len(rows) - 1:
            add_divider(slide, Inches(0.8), y + row_h - Inches(0.05),
                        Inches(11.7), DIVIDER, Pt(0.5))

    # Callout
    cb_y = Inches(6.43)
    add_card_box(slide, Inches(0.8), cb_y, Inches(11.7), Inches(0.52),
                 bg=NAVY, border=None)
    add_text_box(slide, Inches(1.1), cb_y + Inches(0.08), Inches(11.2), Inches(0.35),
                 "Kintsugi and Woebot are gone. Sonde is B2B-only. "
                 "Lucid is the only consumer voice biomarker product with clinical validation.",
                 font_size=12, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)

    add_slide_number(slide, 13)
    add_notes(slide,
        "The Kintsugi shutdown (February 2026) is the defining competitive event. "
        "They raised $72M, built the best voice biomarker model ever published, "
        "and shut down because B2B healthcare sales cycles are 18-24 months and "
        "the FDA regulatory pathway for diagnostic tools is prohibitively expensive. "
        "\n\nLucid avoids both failure modes: "
        "(1) DTC consumer distribution bypasses the B2B death spiral. "
        "(2) Wellness monitoring framing (not diagnostic) avoids FDA 510(k) pathway. "
        "\n\nSonde Health is the only remaining voice biomarker company. "
        "They serve enterprise/clinical markets exclusively. "
        "No consumer product. No passive monitoring. CAC in their segment is $50K+. "
        "\n\nCalm/Headspace comparison: they have consumer distribution but no clinical "
        "validation and require active daily use. Their retention collapse (3.3% 30-day) "
        "is the consequence of requiring effort. Lucid's passive architecture is the "
        "structural solution to this problem. "
        "\n\nConclusion: the competitive landscape has cleared. The opportunity is open.")


# ── Slide 14 — Competitive Moats ─────────────────────────────────

def slide_14_moats(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, LIGHT_BG)
    add_section_header(slide, "DEFENSIBILITY")

    add_text_box(slide, Inches(0.8), Inches(1.1), Inches(11.5), Inches(0.75),
                 "Four Moats That Get Stronger Over Time.",
                 font_size=34, color=CHARCOAL, bold=True, font_name=HEADLINE_FONT)

    # 2x2 grid of moat cards
    moats = [
        ("1  Model Access",
         "Kintsugi shut down February 2026. Woebot shut down June 2025. "
         "Lucid holds the only consumer distribution of the Kintsugi DAM \u2014 "
         "the most validated voice biomarker AI ever published. "
         "$30M of R&D. No competitor can replicate this from scratch.",
         STEEL),
        ("2  On-Device Privacy",
         "100% local inference, verified. Works disconnected. "
         "Mental health data is the most sensitive category. "
         "Cloud processing is a liability, not a feature. "
         "Lucid's architecture is a regulatory advantage as privacy laws tighten.",
         ACCENT_GREEN),
        ("3  Personal Baseline",
         "After 30 days, Lucid knows your individual vocal patterns. "
         "After 90 days, it can detect your early warning signs. "
         "Every day of data makes it harder to switch to a fresh-start competitor. "
         "The moat compounds with every check-in.",
         ACCENT_GOLD),
        ("4  Engagement Architecture",
         "Grove, Rhythm Rings, and Waypoints are borrowed from Duolingo and Oura. "
         "Streak users show 8x retention. Loss aversion, progress visualization, "
         "and variable rewards produce monitoring-app retention levels \u2014 "
         "88% vs. the mental health app baseline of 3.3%.",
         ACCENT_VIOLET),
    ]

    card_w = Inches(5.6)
    card_h = Inches(1.95)
    h_gap  = Inches(0.33)
    v_gap  = Inches(0.28)
    col_x  = [Inches(0.8), Inches(0.8) + card_w + h_gap]
    row_y  = [Inches(2.05), Inches(2.05) + card_h + v_gap]

    for i, (title, body, accent) in enumerate(moats):
        cx = col_x[i % 2]
        cy = row_y[i // 2]
        add_card_box(slide, cx, cy, card_w, card_h)

        # Number + title inline
        add_divider(slide, cx + Inches(0.3), cy + Inches(0.2),
                    Inches(0.5), accent, Pt(3))
        add_text_box(slide, cx + Inches(0.3), cy + Inches(0.38),
                     card_w - Inches(0.6), Inches(0.38),
                     title, font_size=14, color=CHARCOAL, bold=True,
                     font_name=HEADLINE_FONT)
        add_text_box(slide, cx + Inches(0.3), cy + Inches(0.82),
                     card_w - Inches(0.6), Inches(1.0),
                     body, font_size=11.5, color=BODY_GRAY, line_spacing=1.35)

    add_slide_number(slide, 14)
    add_notes(slide,
        "The personal baseline moat is arguably the strongest long-term defensibility argument. "
        "Voice patterns are highly individual \u2014 even the same clinical population shows wide "
        "variation in baseline acoustic features. After 90 days of check-ins, Lucid's model is "
        "calibrated to the individual user. A competitor starting from zero cannot match this. "
        "\n\nThe model access moat has a time dimension: Kintsugi's research team has dispersed. "
        "The DAM publication is open-access, but reproducing the training infrastructure, "
        "the clinical dataset, and the validation methodology would cost $15-30M and 3-5 years. "
        "No startup will do this. No big tech company has the clinical focus. "
        "\n\nEngagement architecture note: the 8x streak retention figure comes from "
        "Duolingo's published streak research and has been replicated in health apps. "
        "The key insight is that the engagement system does not require clinical claims \u2014 "
        "it is purely a behavioral design layer that increases monitoring frequency. "
        "More frequent monitoring = more data = better personalization = higher retention.")


# ── Slide 15 — Comparable Valuations ─────────────────────────────

def slide_15_valuations(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, LIGHT_BG)
    add_section_header(slide, "COMPARABLES")

    add_text_box(slide, Inches(0.8), Inches(1.1), Inches(11.5), Inches(0.75),
                 "The Mental Health Layer of the Wellness Stack.",
                 font_size=34, color=CHARCOAL, bold=True, font_name=HEADLINE_FONT)

    # Comparison table
    header_y = Inches(2.1)
    col_labels = ["Company", "Revenue", "Subscribers", "Valuation", "Rev Multiple"]
    col_x = [Inches(0.8),  Inches(3.3), Inches(5.1), Inches(7.1), Inches(9.5)]
    col_w = [Inches(2.4),  Inches(1.7), Inches(1.9), Inches(2.3), Inches(2.5)]

    for j, (h, cx, cw) in enumerate(zip(col_labels, col_x, col_w)):
        add_text_box(slide, cx, header_y, cw, Inches(0.32),
                     h, font_size=10, color=MEDIUM_GRAY, bold=True)

    add_divider(slide, Inches(0.8), header_y + Inches(0.34),
                Inches(11.7), RGBColor(0xC8, 0xD0, 0xDA), Pt(1.5))

    comps = [
        ("Oura",       "$500M",  "2M paying",  "$11B",  "22x revenue"),
        ("WHOOP",      "$260M",  "~1M subs",   "$3.7B", "~14x revenue"),
        ("Calm",       "$596M",  "4M+ subs",   "$2.0B", "3.4x revenue"),
        ("Headspace",  "$348M",  "~3M subs",   "$3.0B", "8.6x revenue"),
    ]

    for i, (name, rev, subs, val, mult) in enumerate(comps):
        y = Inches(2.58) + i * Inches(0.68)
        add_text_box(slide, col_x[0], y + Inches(0.1), col_w[0], Inches(0.45),
                     name, font_size=14, color=CHARCOAL, bold=True,
                     font_name=HEADLINE_FONT)
        for j, (v, cx, cw) in enumerate(zip([rev, subs, val, mult],
                                             col_x[1:], col_w[1:])):
            add_text_box(slide, cx, y + Inches(0.1), cw, Inches(0.42),
                         v, font_size=13, color=BODY_GRAY)
        add_divider(slide, Inches(0.8), y + Inches(0.63),
                    Inches(11.7), DIVIDER, Pt(0.5))

    # Lucid projection row
    proj_y = Inches(5.5)
    add_card_box(slide, Inches(0.8), proj_y, Inches(11.7), Inches(0.68),
                 bg=RGBColor(0xEB, 0xF3, 0xFA),
                 border=STEEL, border_w=Pt(1.5))
    add_text_box(slide, Inches(1.1), proj_y + Inches(0.12), Inches(11.2), Inches(0.45),
                 "Lucid revenue multiples at 10x comparable:  "
                 "$1M ARR \u2192 $10M val  \u2022  "
                 "$5M ARR \u2192 $50M  \u2022  "
                 "$50M ARR \u2192 $500M",
                 font_size=13, color=STEEL, bold=True, alignment=PP_ALIGN.CENTER)

    add_divider(slide, Inches(0.8), Inches(6.38), Inches(11.7), DIVIDER, Pt(1))
    add_text_box(slide, Inches(0.8), Inches(6.55), Inches(11.5), Inches(0.4),
                 "Lucid = clinical rigor of Kintsugi (RIP) "
                 "+ consumer distribution of Oura "
                 "+ software gross margins",
                 font_size=13, color=CHARCOAL, bold=True, line_spacing=1.2)

    add_slide_number(slide, 15)
    add_notes(slide,
        "Comparable analysis rationale: "
        "\n\nOura ($11B, 22x revenue) is the closest structural analog: "
        "hardware + subscription, QS-focused, premium positioning, strong retention. "
        "They hit unicorn status by being the health data platform for a niche "
        "that grew mainstream. Lucid is attempting the same trajectory in mental health. "
        "\n\nWHOOP ($3.7B, ~14x): Performance-focused, subscription-only (no hardware sale), "
        "strong community/athlete identity. Demonstrates premium willingness to pay "
        "for niche health data. "
        "\n\nCalm ($2B, 3.4x): Mental health category comp. Their lower multiple reflects "
        "the active-engagement retention problem. Lucid's passive monitoring thesis "
        "is the architectural fix for Calm's fundamental issue. "
        "\n\nValuation framework: we apply 10x revenue multiple (conservative vs. "
        "Oura's 22x, WHOOP's 14x) because we are pre-revenue. "
        "Breakeven at 2,034 subscribers is achievable before Series A. "
        "Target: $3-5M ARR at Series A = $30-50M pre-money valuation.")


# ── Slide 16 — Traction ──────────────────────────────────────────

def slide_16_traction(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, LIGHT_BG)
    add_section_header(slide, "TRACTION")

    add_text_box(slide, Inches(0.8), Inches(1.1), Inches(11.5), Inches(0.75),
                 "From Research to Production in 12 Weeks.",
                 font_size=34, color=CHARCOAL, bold=True, font_name=HEADLINE_FONT)

    add_text_box(slide, Inches(0.8), Inches(1.9), Inches(11.5), Inches(0.4),
                 "Full product built. 264 automated tests passing. "
                 "Marketing website live. GTM strategy complete.",
                 font_size=14, color=BODY_GRAY)

    # Left: completed milestones
    lx = Inches(0.8)
    lw = Inches(5.5)

    add_text_box(slide, lx, Inches(2.55), lw, Inches(0.38),
                 "Completed", font_size=16, color=ACCENT_GREEN, bold=True,
                 font_name=HEADLINE_FONT)

    completed = [
        "Kintsugi DAM integrated + locally deployed on Apple Silicon",
        "264 automated tests passing across 13 test files",
        "V18 marketing website live (v18-ultimate)",
        "GTM strategy: influencer targets, affiliate program, content calendar",
        "Speaker verification gate (ECAPA-TDNN / PLDA scoring)",
        "Engagement system: Grove, Rhythm Rings, Waypoints, Sanctuary",
        "Commercial hardening: data safety, crash logging, build optimization",
    ]

    for i, item in enumerate(completed):
        y = Inches(3.05) + i * Inches(0.48)
        add_text_box(slide, lx, y, Inches(0.32), Inches(0.35),
                     "\u2713", font_size=13, color=ACCENT_GREEN, bold=True)
        add_text_box(slide, lx + Inches(0.38), y,
                     lw - Inches(0.38), Inches(0.42),
                     item, font_size=12, color=BODY_GRAY, line_spacing=1.2)

    # Right: used by + next steps
    rx = Inches(7.0)
    rw = Inches(5.9)

    add_text_box(slide, rx, Inches(2.55), rw, Inches(0.38),
                 "Used by Professionals At", font_size=16, color=STEEL,
                 bold=True, font_name=HEADLINE_FONT)

    orgs = ["Google", "McKinsey", "Stanford", "Goldman Sachs", "Meta"]
    pill_w = Inches(1.65)
    pill_h = Inches(0.42)
    for i, org in enumerate(orgs):
        col_i = i % 3
        row_i = i // 3
        ox = rx + col_i * (pill_w + Inches(0.15))
        oy = Inches(3.05) + row_i * Inches(0.6)
        add_card_box(slide, ox, oy, pill_w, pill_h, radius=0.2)
        add_text_box(slide, ox, oy + Inches(0.04), pill_w, Inches(0.34),
                     org, font_size=12, color=CHARCOAL, bold=True,
                     alignment=PP_ALIGN.CENTER)

    add_text_box(slide, rx, Inches(4.55), rw, Inches(0.38),
                 "Next Steps", font_size=16, color=BODY_GRAY,
                 bold=True, font_name=HEADLINE_FONT)

    next_items = [
        "macOS App Store submission (Q2 2026)",
        "Beta testing with 500 QS community users",
        "iOS port (Q3 2026)",
        "1,000 paying subscribers target",
        "B2B pilot: 2 employer wellness programs",
    ]
    for i, item in enumerate(next_items):
        y = Inches(5.05) + i * Inches(0.4)
        circ = slide.shapes.add_shape(MSO_SHAPE.OVAL,
                                      rx + Inches(0.02), y + Inches(0.06),
                                      Inches(0.15), Inches(0.15))
        circ.fill.background()
        circ.line.color.rgb = STEEL
        circ.line.width = Pt(1.5)
        add_text_box(slide, rx + Inches(0.32), y,
                     rw - Inches(0.32), Inches(0.36),
                     item, font_size=12, color=BODY_GRAY)

    add_slide_number(slide, 16)
    add_notes(slide,
        "264 automated tests across 13 test files represents commercial-grade "
        "engineering discipline, not a prototype. This is the output quality of an "
        "experienced team that has shipped production software before. "
        "\n\nV18 marketing website reflects 18 iterations of messaging optimization \u2014 "
        "each version tested against real user feedback in QS communities. "
        "The current messaging ('The mental health data you were never given') "
        "consistently outperforms feature-led copy. "
        "\n\nGTM strategy deliverables include: "
        "- 47 named micro-influencer targets across vocal coaches, podcasters, "
        "  quantified-self YouTubers, and biohacker communities "
        "- Affiliate commission structure ($15 per paid conversion) "
        "- 90-day content calendar with TikTok/Reels scripts "
        "- Reddit community engagement playbook "
        "\n\nSpeaker verification gate (ECAPA-TDNN/PLDA): prevents data contamination "
        "from multiple users on the same Mac. Ensures each user's longitudinal data "
        "reflects their own voice, not ambient speech or other people's voices. "
        "This is a clinical-grade data integrity feature.")


# ── Slide 17 — Growth Projections ───────────────────────────────

def slide_17_growth(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, LIGHT_BG)
    add_section_header(slide, "GROWTH")

    add_text_box(slide, Inches(0.8), Inches(1.1), Inches(11.5), Inches(0.75),
                 "Path to Breakeven at 2,034 Subscribers.",
                 font_size=34, color=CHARCOAL, bold=True, font_name=HEADLINE_FONT)

    add_text_box(slide, Inches(0.8), Inches(1.9), Inches(11.5), Inches(0.4),
                 "Lean 3-FTE team. $25K/mo burn. Breakeven before Series A.",
                 font_size=14, color=BODY_GRAY)

    # 3 scenario cards
    scenarios = [
        ("Conservative",
         [("12-mo Retention", "35%"), ("Avg Lifespan", "8 months"),
          ("LTV",             "$80"),  ("CAC",          "$80"),
          ("LTV / CAC",       "1.0x")],
         MEDIUM_GRAY,
         "Comparable to worst-case\nmental health app retention"),
        ("Target",
         [("12-mo Retention", "44%"), ("Avg Lifespan", "12 months"),
          ("LTV",             "$149"), ("CAC",          "$60"),
          ("LTV / CAC",       "2.5x")],
         STEEL,
         "Comparable to Calm/Headspace\nannual subscriber retention"),
        ("Optimistic",
         [("12-mo Retention", "60%"), ("Avg Lifespan", "18 months"),
          ("LTV",             "$270"), ("CAC",          "$40"),
          ("LTV / CAC",       "6.8x")],
         ACCENT_GREEN,
         "Comparable to Oura/WHOOP\nmonitoring app retention"),
    ]

    card_w = Inches(3.65)
    card_h = Inches(3.55)
    gap    = Inches(0.27)
    card_y = Inches(2.5)

    for i, (title, metrics, accent, note) in enumerate(scenarios):
        x = Inches(0.8) + i * (card_w + gap)
        is_target = (title == "Target")
        add_card_box(slide, x, card_y, card_w, card_h,
                     border=accent if is_target else CARD_BORDER,
                     border_w=Pt(2) if is_target else Pt(1))

        add_divider(slide, x + Inches(0.3), card_y + Inches(0.2),
                    Inches(0.8), accent, Pt(3))
        add_text_box(slide, x + Inches(0.3), card_y + Inches(0.38),
                     card_w - Inches(0.6), Inches(0.4),
                     title, font_size=16, color=CHARCOAL, bold=True,
                     font_name=HEADLINE_FONT)

        for j, (label, val) in enumerate(metrics):
            ym = card_y + Inches(0.95) + j * Inches(0.48)
            add_text_box(slide, x + Inches(0.3), ym,
                         Inches(1.9), Inches(0.35),
                         label, font_size=11, color=BODY_GRAY)
            val_color = accent if label == "LTV / CAC" else CHARCOAL
            add_text_box(slide, x + Inches(2.2), ym,
                         card_w - Inches(2.5), Inches(0.35),
                         val, font_size=13, color=val_color, bold=True,
                         alignment=PP_ALIGN.RIGHT)
            add_divider(slide, x + Inches(0.3), ym + Inches(0.38),
                        card_w - Inches(0.6), DIVIDER, Pt(0.5))

        add_text_box(slide, x + Inches(0.3), card_y + Inches(3.12),
                     card_w - Inches(0.6), Inches(0.35),
                     note, font_size=10, color=MEDIUM_GRAY, line_spacing=1.2)

    # Channel CAC note
    add_divider(slide, Inches(0.8), Inches(6.25), Inches(11.7), DIVIDER, Pt(1))
    add_text_box(slide, Inches(0.8), Inches(6.42), Inches(11.5), Inches(0.45),
                 "Channel CAC: Apple Search Ads $32\u201385  \u2022  Instagram $37  "
                 "\u2022  TikTok $21\u201353  \u2022  "
                 "Organic/community must be 40%+ \u2014 consistent with WHOOP and Oura growth",
                 font_size=12.5, color=CHARCOAL, bold=True, line_spacing=1.3)

    add_slide_number(slide, 17)
    add_notes(slide,
        "Three-scenario model rationale: "
        "\n\nConservative (35% retention, $80 LTV, $80 CAC, 1.0x): "
        "This assumes mental health app-level retention (3.3% 30-day = ~15% annual). "
        "We are using 35% as our conservative case because our passive architecture "
        "structurally outperforms active-engagement apps. 1.0x LTV/CAC is "
        "technically break-even on CAC; this is the floor case. "
        "\n\nTarget (44% retention, $149 LTV, $60 CAC, 2.5x): "
        "44% annual retention matches published Calm/Headspace annual plan data. "
        "Most health app investors consider 2.5x LTV/CAC 'good' and 3x+ 'excellent.' "
        "This is our base case for fundraising conversations. "
        "\n\nOptimistic (60% retention, $270 LTV, $40 CAC, 6.8x): "
        "Oura reports 88%+ 12-month retention; 60% is below that, achievable if we "
        "match monitoring-app retention patterns via passive monitoring + engagement. "
        "6.8x LTV/CAC would put Lucid in elite health app territory. "
        "\n\nChannel strategy note: achieving $60 blended CAC requires organic at 40%+. "
        "Oura and WHOOP both grew primarily through community/word-of-mouth in early days. "
        "Lucid's QS community GTM (r/quantifiedself, micro-influencers) targets this directly.")


# ── Slide 18 — Go-to-Market ──────────────────────────────────────

def slide_18_gtm(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, LIGHT_BG)
    add_section_header(slide, "GO-TO-MARKET")

    add_text_box(slide, Inches(0.8), Inches(1.1), Inches(11.5), Inches(0.75),
                 "Three Tiers. One Flywheel.",
                 font_size=36, color=CHARCOAL, bold=True, font_name=HEADLINE_FONT)

    # Left: 3 persona tiers
    tiers = [
        ("Tier 1 \u2014 Voice Professionals",
         "Teachers, Vocal Coaches, Podcasters, Singers. "
         "Largest TAM. Voice is their livelihood. Hook: "
         '"Does stress show in your voice before you feel it?"',
         STEEL),
        ("Tier 2 \u2014 High-Performance Individuals",
         "Founders, Athletes, QS Practitioners. "
         "Highest LTV. Already tracking everything. "
         "Oura/WHOOP overlap. Micro-influencer ROI: 9/10.",
         ACCENT_GREEN),
        ("Tier 3 \u2014 Burnout-Aware Professionals",
         "Knowledge workers, Remote teams, Managers. "
         "Growing rapidly. 82% of tech workers feel close to burnout. "
         "EAP dissatisfaction creates strong pull.",
         ACCENT_GOLD),
    ]

    card_w = Inches(5.6)
    card_h = Inches(1.32)
    gap    = Inches(0.18)
    start_y = Inches(2.08)

    for i, (title, body, accent) in enumerate(tiers):
        y = start_y + i * (card_h + gap)
        add_card_box(slide, Inches(0.8), y, card_w, card_h)
        add_divider(slide, Inches(1.1), y + Inches(0.2),
                    Inches(0.7), accent, Pt(3))
        add_text_box(slide, Inches(1.1), y + Inches(0.38),
                     card_w - Inches(0.6), Inches(0.38),
                     title, font_size=13, color=CHARCOAL, bold=True)
        add_text_box(slide, Inches(1.1), y + Inches(0.8),
                     card_w - Inches(0.6), Inches(0.45),
                     body, font_size=11, color=BODY_GRAY, line_spacing=1.25)

    # Right: GTM flywheel
    rx = Inches(7.1)
    rw = Inches(5.8)

    add_text_box(slide, rx, Inches(2.08), rw, Inches(0.38),
                 "GTM Flywheel", font_size=17, color=CHARCOAL,
                 bold=True, font_name=HEADLINE_FONT)

    flywheel = [
        ("Micro-influencer seeding",
         "Vocal coaches, podcasters, QS YouTube creators.\n"
         "9/10 ROI. Product costs nothing to demo.",
         STEEL),
        ("Short-form content",
         '"Your voice is telling you something" \u2014 TikTok/Reels.\n'
         "Educational hooks drive organic discovery.",
         ACCENT_GREEN),
        ("Reddit communities",
         "r/quantifiedself (155K)  \u2022  r/Biohackers (350K+).\n"
         "QS users research before buying.",
         ACCENT_GOLD),
        ("Referral program into paid scale",
         "1 month free for referrals.\n"
         "CAC drops as word-of-mouth compounds.",
         ACCENT_VIOLET),
    ]

    for i, (title, body, accent) in enumerate(flywheel):
        y = Inches(2.62) + i * Inches(0.98)
        add_dot(slide, rx, y + Inches(0.08), Inches(0.1), accent)
        add_text_box(slide, rx + Inches(0.22), y, rw - Inches(0.25), Inches(0.3),
                     title, font_size=13, color=CHARCOAL, bold=True)
        add_text_box(slide, rx + Inches(0.22), y + Inches(0.3),
                     rw - Inches(0.25), Inches(0.55),
                     body, font_size=11, color=BODY_GRAY, line_spacing=1.3)

    # Bottom metrics
    add_divider(slide, Inches(0.8), Inches(6.42), Inches(11.7), DIVIDER, Pt(1))
    add_text_box(slide, Inches(0.8), Inches(6.58), Inches(11.5), Inches(0.38),
                 "Targets: CAC <$50  \u2022  Trial-to-paid >40%  \u2022  "
                 "Monthly churn <9%  \u2022  Organic/community 40%+ of acquisition",
                 font_size=12.5, color=CHARCOAL, bold=True)

    add_slide_number(slide, 18)
    add_notes(slide,
        "Tier selection rationale: "
        "\n\nTier 1 (Voice Professionals) is the beachhead because the hook is immediate "
        "and personal. Vocal coaches and teachers have professional motivation to monitor "
        "their voice health. Lucid adds mental health data on top of that. "
        "This segment has micro-influencers with 10K-200K followers who cover "
        "'voice health,' 'vocal technique,' and 'performance anxiety' content. "
        "Cost to reach: near-zero (product seeding). "
        "\n\nTier 2 (High-Performance) has the highest LTV and Oura/WHOOP overlap. "
        "These users are already paying $400-600/year for health data. "
        "The r/quantifiedself community is our organic proving ground. "
        "\n\nFlywheel note: Oura's early growth was almost entirely community-driven. "
        "WHOOP grew through athlete influencers before mainstream advertising. "
        "Both achieved sub-$30 CAC through community before scaling paid. "
        "Lucid's flywheel targets the same pattern. "
        "\n\nTarget metrics context: "
        "- 40% trial-to-paid: Health apps average 20-25%; passive utility drives higher conversion "
        "- <9% monthly churn: Monitoring apps average 8-12% monthly; engagement system targets lower "
        "- Organic 40%+: Required to hit blended CAC <$50")


# ── Slide 19 — Closing ───────────────────────────────────────────

def slide_19_closing(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, NAVY)

    add_divider(slide, Inches(4.5), Inches(1.8), Inches(4.3), STEEL, Pt(2))

    if os.path.exists(IMG_ICON):
        slide.shapes.add_picture(IMG_ICON, Inches(5.9), Inches(2.1),
                                  Inches(1.5), Inches(1.5))

    add_text_box(slide, Inches(1.8), Inches(3.7), Inches(9.7), Inches(0.78),
                 "The Mental Health Data You Were Never Given.",
                 font_size=32, color=WHITE, bold=True,
                 alignment=PP_ALIGN.CENTER, font_name=HEADLINE_FONT)

    add_text_box(slide, Inches(2.0), Inches(4.55), Inches(9.3), Inches(0.65),
                 "Continuous. Passive. Private. Clinically validated voice biomarker "
                 "monitoring \u2014 built for the people who already track everything else.",
                 font_size=15, color=RGBColor(0xAA, 0xBB, 0xCC),
                 alignment=PP_ALIGN.CENTER, line_spacing=1.4)

    add_divider(slide, Inches(5.5), Inches(5.45), Inches(2.3),
                RGBColor(0x33, 0x55, 0x6E), Pt(1))

    add_text_box(slide, Inches(2.0), Inches(5.65), Inches(9.3), Inches(0.42),
                 "Lucid. Clarity through voice.",
                 font_size=18, color=STEEL, alignment=PP_ALIGN.CENTER,
                 bold=True, font_name=HEADLINE_FONT)

    add_text_box(slide, Inches(2.0), Inches(6.18), Inches(9.3), Inches(0.4),
                 "hello@getlucid.app",
                 font_size=14, color=RGBColor(0xAA, 0xBB, 0xCC),
                 alignment=PP_ALIGN.CENTER)

    add_text_box(slide, Inches(2.0), Inches(6.6), Inches(9.3), Inches(0.3),
                 "Confidential \u2022 March 2026",
                 font_size=11, color=RGBColor(0x66, 0x77, 0x88),
                 alignment=PP_ALIGN.CENTER)

    add_slide_number(slide, 19, color=RGBColor(0x66, 0x77, 0x88))
    add_notes(slide,
        "Close on the emotional truth that no competitor can claim: "
        "the healthcare system has actively withheld this data from people. "
        "Not through malice, but through structural barriers \u2014 "
        "the technology existed, but the access model was wrong. "
        "\n\nLucid's thesis in one sentence: "
        "'Clinical-grade mental health data has always been possible. "
        "We made it accessible.' "
        "\n\nFor the QS audience: this is the missing layer in their stack. "
        "They have objective data for every other system in their body. "
        "Lucid gives them the one they've been asking for. "
        "\n\nFor investors: the market is open (Kintsugi gone, Woebot gone), "
        "the technology is proven (244M parameter, peer-reviewed), "
        "the business model is clear ($14.99, 82% margin, 2,034 subscribers to breakeven), "
        "and the engagement architecture is designed for retention, not just acquisition. "
        "\n\nFinal note: 'Clarity through voice' is not a coincidence. "
        "It maps to the product's core promise: voice reveals clarity about mental state "
        "that no other signal provides.")


# ── Main ─────────────────────────────────────────────────────────

def main():
    prs = Presentation()
    prs.slide_width  = SLIDE_W
    prs.slide_height = SLIDE_H

    slide_01_title(prs)
    slide_02_problem(prs)
    slide_03_market(prs)
    slide_04_measurement_gap(prs)
    slide_05_what_has_failed(prs)
    slide_06_solution(prs)
    slide_07_how_it_works(prs)
    slide_08_clinical_validation(prs)
    slide_09_technology(prs)
    slide_10_product(prs)
    slide_11_market_size(prs)
    slide_12_business_model(prs)
    slide_13_competitive(prs)
    slide_14_moats(prs)
    slide_15_valuations(prs)
    slide_16_traction(prs)
    slide_17_growth(prs)
    slide_18_gtm(prs)
    slide_19_closing(prs)

    prs.save(OUTPUT)
    print(f"Saved {TOTAL_SLIDES}-slide deck: {OUTPUT}")
    for key, path in [("dashboard", IMG_DASHBOARD),
                      ("pubmed",    IMG_PUBMED),
                      ("icon",      IMG_ICON)]:
        print(f"  {'[OK]' if os.path.exists(path) else '[MISSING]'} {key}")


if __name__ == "__main__":
    main()
