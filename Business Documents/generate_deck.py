#!/usr/bin/env python3
"""
Lucid Investment Deck Generator — v3
Light-gray-dominant design, proper image aspect ratios, QS framing.
19 slides (no team/ask). 16:9 format.
"""

import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION, XL_LABEL_POSITION
from pptx.chart.data import CategoryChartData
from lxml import etree

# ── Design System v3 ────────────────────────────────────────────
# Light gray dominant (~80% of slides), navy for 2-3 impact slides only

NAVY = RGBColor(0x1A, 0x3A, 0x4F)
STEEL = RGBColor(0x5B, 0x8D, 0xB8)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_BG = RGBColor(0xF8, 0xF9, 0xFA)

# Text hierarchy on light backgrounds
CHARCOAL = RGBColor(0x2D, 0x2D, 0x3A)       # Headlines
BODY_GRAY = RGBColor(0x55, 0x55, 0x66)       # Body text
MEDIUM_GRAY = RGBColor(0x88, 0x88, 0x99)     # Captions / slide numbers
LIGHT_GRAY = RGBColor(0xCC, 0xCC, 0xCC)      # Borders

# Text hierarchy on dark backgrounds
DARK_TEXT = RGBColor(0x1A, 0x1A, 0x2E)

# Accents — used sparingly
ACCENT_GREEN = RGBColor(0x4E, 0xC9, 0xB0)
ACCENT_GOLD = RGBColor(0xD4, 0xA5, 0x37)
ACCENT_CORAL = RGBColor(0xE8, 0x6B, 0x6B)

# Card colors for light backgrounds
CARD_BG = RGBColor(0xFF, 0xFF, 0xFF)          # White cards on light gray
CARD_BORDER = RGBColor(0xE0, 0xE0, 0xE8)     # Subtle border

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)
TOTAL_SLIDES = 19

# ── Asset Paths ─────────────────────────────────────────────────
BASE = os.path.dirname(os.path.abspath(__file__))
PROJECT_ROOT = os.path.dirname(BASE)
IMG_DASHBOARD = os.path.join(PROJECT_ROOT, "lucid-website/actual-websites/images/dashboard-health-score.png")
IMG_PUBMED = os.path.join(PROJECT_ROOT, "lucid-website/actual-websites/images/pubmed-study.png")
IMG_LIFESTYLE = os.path.join(PROJECT_ROOT, "lucid-website/actual-websites/images/lifestyle-attune-man.png")
IMG_ICON = os.path.join(BASE, "assets/icon.png")
OUTPUT = os.path.join(BASE, "Lucid Investment Deck.pptx")

# Native image dimensions (pixels) — for aspect ratio preservation
IMG_DIMS = {
    "dashboard": (1390, 912),   # 1.52:1 landscape
    "pubmed":    (1538, 1576),  # 0.98:1 portrait
    "lifestyle": (2816, 1536),  # 1.83:1 landscape
    "icon":      (1024, 1024),  # 1:1 square
}

HEADLINE_FONT = "Georgia"
BODY_FONT = "Helvetica Neue"


# ── Image Aspect Ratio Helper ───────────────────────────────────

def fit_image_dims(image_key, max_width_in, max_height_in):
    """Return (width, height) in Inches that fit within bounding box
    while preserving the native aspect ratio of the image."""
    pw, ph = IMG_DIMS[image_key]
    aspect = pw / ph  # width / height

    # Try fitting to max width first
    w = max_width_in
    h = w / aspect
    if h > max_height_in:
        # Constrained by height instead
        h = max_height_in
        w = h * aspect

    return Inches(w), Inches(h)


# ── Utility Functions ───────────────────────────────────────────

def add_bg(slide, color):
    """Add a full-slide background rectangle."""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Emu(0), Emu(0), SLIDE_W, SLIDE_H
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    sp = shape._element
    sp.getparent().remove(sp)
    slide.shapes._spTree.insert(2, sp)


def add_text_box(slide, left, top, width, height, text, font_size=18,
                 color=CHARCOAL, bold=False, alignment=PP_ALIGN.LEFT,
                 font_name=BODY_FONT, line_spacing=1.2, anchor=MSO_ANCHOR.TOP):
    """Add a text box with a single run of text."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    txBox.text_frame.word_wrap = True
    txBox.text_frame.auto_size = None
    p = txBox.text_frame.paragraphs[0]
    p.alignment = alignment
    p.space_after = Pt(0)
    p.space_before = Pt(0)
    pPr = p._p.get_or_add_pPr()
    lnSpc = etree.SubElement(pPr, '{http://schemas.openxmlformats.org/drawingml/2006/main}lnSpc')
    spcPct = etree.SubElement(lnSpc, '{http://schemas.openxmlformats.org/drawingml/2006/main}spcPct')
    spcPct.set('val', str(int(line_spacing * 100000)))

    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.color.rgb = color
    run.font.bold = bold
    run.font.name = font_name
    txBox.text_frame._txBody.bodyPr.set('anchor', {
        MSO_ANCHOR.TOP: 't', MSO_ANCHOR.MIDDLE: 'ctr', MSO_ANCHOR.BOTTOM: 'b'
    }.get(anchor, 't'))
    return txBox


def add_multiline_box(slide, left, top, width, height, lines, default_size=18,
                      default_color=BODY_GRAY, default_bold=False, alignment=PP_ALIGN.LEFT,
                      default_font=BODY_FONT, line_spacing=1.2):
    """Add a text box with multiple paragraphs."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    txBox.text_frame.word_wrap = True
    txBox.text_frame.auto_size = None

    for i, line in enumerate(lines):
        if isinstance(line, str):
            line = {"text": line}

        if i == 0:
            p = txBox.text_frame.paragraphs[0]
        else:
            p = txBox.text_frame.add_paragraph()

        p.alignment = line.get("alignment", alignment)
        p.space_after = Pt(line.get("space_after", 4))
        p.space_before = Pt(line.get("space_before", 0))

        pPr = p._p.get_or_add_pPr()
        lnSpc = etree.SubElement(pPr, '{http://schemas.openxmlformats.org/drawingml/2006/main}lnSpc')
        spcPct = etree.SubElement(lnSpc, '{http://schemas.openxmlformats.org/drawingml/2006/main}spcPct')
        spcPct.set('val', str(int(line.get("line_spacing", line_spacing) * 100000)))

        run = p.add_run()
        run.text = line.get("text", "")
        run.font.size = Pt(line.get("size", default_size))
        run.font.color.rgb = line.get("color", default_color)
        run.font.bold = line.get("bold", default_bold)
        run.font.name = line.get("font", default_font)
        if line.get("italic"):
            run.font.italic = True

    return txBox


def add_slide_number(slide, num, color=MEDIUM_GRAY):
    """Add slide number in bottom right."""
    add_text_box(
        slide, Inches(11.8), Inches(6.9), Inches(1.2), Inches(0.4),
        f"{num} / {TOTAL_SLIDES}", font_size=10, color=color,
        alignment=PP_ALIGN.RIGHT
    )


def add_stat_box(slide, left, top, width, height, number, label,
                 num_color=STEEL, label_color=BODY_GRAY, bg_color=None,
                 num_size=44, label_size=13, border_color=None):
    """Add a stat callout box with large number and label below."""
    if bg_color:
        box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
        box.fill.solid()
        box.fill.fore_color.rgb = bg_color
        if border_color:
            box.line.color.rgb = border_color
            box.line.width = Pt(1)
        else:
            box.line.fill.background()
        box.adjustments[0] = 0.05

    add_text_box(slide, left, top + Inches(0.2), width, Inches(0.8),
                 number, font_size=num_size, color=num_color,
                 bold=True, alignment=PP_ALIGN.CENTER, font_name=HEADLINE_FONT)
    add_text_box(slide, left, top + Inches(0.9), width, Inches(0.6),
                 label, font_size=label_size, color=label_color,
                 alignment=PP_ALIGN.CENTER, line_spacing=1.1)


def add_card(slide, left, top, width, height, title, body,
             bg_color=CARD_BG, title_color=CHARCOAL,
             body_color=BODY_GRAY, corner_radius=0.06, border_color=CARD_BORDER,
             body_size=12):
    """Add a rounded card with title and body text — light theme."""
    box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    box.fill.solid()
    box.fill.fore_color.rgb = bg_color
    if border_color:
        box.line.color.rgb = border_color
        box.line.width = Pt(1)
    else:
        box.line.fill.background()
    box.adjustments[0] = corner_radius

    add_text_box(slide, left + Inches(0.25), top + Inches(0.2),
                 width - Inches(0.5), Inches(0.4),
                 title, font_size=15, color=title_color, bold=True)
    add_text_box(slide, left + Inches(0.25), top + Inches(0.55),
                 width - Inches(0.5), height - Inches(0.7),
                 body, font_size=body_size, color=body_color, line_spacing=1.3)


def add_divider_line(slide, left, top, width, color=STEEL, thickness=Pt(1.5)):
    """Add a horizontal divider line."""
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, thickness)
    line.fill.solid()
    line.fill.fore_color.rgb = color
    line.line.fill.background()


def add_accent_dot(slide, left, top, size=Inches(0.12), color=STEEL):
    """Add a small colored dot as a bullet/accent."""
    dot = slide.shapes.add_shape(MSO_SHAPE.OVAL, left, top, size, size)
    dot.fill.solid()
    dot.fill.fore_color.rgb = color
    dot.line.fill.background()


def add_picture_safe(slide, img_path, left, top, image_key, max_w, max_h):
    """Add a picture preserving its native aspect ratio within a bounding box."""
    if not os.path.exists(img_path):
        return None
    w, h = fit_image_dims(image_key, max_w, max_h)
    return slide.shapes.add_picture(img_path, left, top, w, h)


# ── Slide Builders ──────────────────────────────────────────────

def slide_01_title(prs):
    """Slide 1 — Title (NAVY background)."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, NAVY)

    add_divider_line(slide, Inches(4.5), Inches(1.5), Inches(4.3), STEEL, Pt(2))

    if os.path.exists(IMG_ICON):
        slide.shapes.add_picture(IMG_ICON, Inches(5.9), Inches(1.8), Inches(1.5), Inches(1.5))

    add_text_box(slide, Inches(2.5), Inches(3.5), Inches(8.3), Inches(0.9),
                 "Lucid.", font_size=54, color=WHITE, bold=True,
                 alignment=PP_ALIGN.CENTER, font_name=HEADLINE_FONT)

    add_text_box(slide, Inches(3.0), Inches(4.5), Inches(7.3), Inches(0.6),
                 "Clinical-Grade Mental Health Monitoring", font_size=22,
                 color=STEEL, alignment=PP_ALIGN.CENTER, line_spacing=1.3)

    add_text_box(slide, Inches(3.0), Inches(5.1), Inches(7.3), Inches(0.5),
                 "From Your Voice. On Your Device. In 25 Seconds.",
                 font_size=16, color=RGBColor(0xAA, 0xBB, 0xCC),
                 alignment=PP_ALIGN.CENTER)

    add_divider_line(slide, Inches(4.5), Inches(5.9), Inches(4.3), STEEL, Pt(2))

    add_slide_number(slide, 1, color=RGBColor(0x66, 0x77, 0x88))


def slide_02_problem_blind_spot(prs):
    """Slide 2 — The Quantified-Self Blind Spot (LIGHT background)."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, LIGHT_BG)

    # Section label
    add_text_box(slide, Inches(0.8), Inches(0.5), Inches(3.0), Inches(0.35),
                 "THE PROBLEM", font_size=11, color=STEEL, bold=True,
                 font_name=BODY_FONT)

    add_divider_line(slide, Inches(0.8), Inches(0.9), Inches(2.0), STEEL, Pt(2))

    # Headline
    add_text_box(slide, Inches(0.8), Inches(1.1), Inches(11.5), Inches(0.8),
                 "The Quantified-Self Blind Spot", font_size=36, color=CHARCOAL,
                 bold=True, font_name=HEADLINE_FONT)

    # Subhead
    add_text_box(slide, Inches(0.8), Inches(1.85), Inches(11.5), Inches(0.6),
                 "You track your heart, your sleep, your steps, your recovery. What about your mind?",
                 font_size=18, color=BODY_GRAY, font_name=BODY_FONT, line_spacing=1.4)

    # Wearable cards row — what they track
    wearables = [
        ("Whoop", "Recovery Score", ACCENT_GREEN),
        ("Oura", "Sleep Quality", RGBColor(0x7B, 0x68, 0xD4)),
        ("Apple Watch", "HRV / Heart Rate", ACCENT_CORAL),
        ("Garmin", "VO2 Max / Training", ACCENT_GOLD),
    ]

    card_w = Inches(2.2)
    card_h = Inches(1.6)
    start_x = Inches(0.8)
    card_y = Inches(2.8)
    gap = Inches(0.25)

    for i, (brand, metric, accent) in enumerate(wearables):
        x = start_x + i * (card_w + gap)
        # Card background
        box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, card_y, card_w, card_h)
        box.fill.solid()
        box.fill.fore_color.rgb = CARD_BG
        box.line.color.rgb = CARD_BORDER
        box.line.width = Pt(1)
        box.adjustments[0] = 0.06

        # Accent line at top of card
        add_divider_line(slide, x + Inches(0.25), card_y + Inches(0.2), Inches(0.5), accent, Pt(3))

        # Brand name
        add_text_box(slide, x + Inches(0.25), card_y + Inches(0.4), card_w - Inches(0.5), Inches(0.4),
                     brand, font_size=15, color=CHARCOAL, bold=True)
        # Metric
        add_text_box(slide, x + Inches(0.25), card_y + Inches(0.85), card_w - Inches(0.5), Inches(0.5),
                     metric, font_size=12, color=BODY_GRAY)

        # Checkmark
        add_text_box(slide, x + card_w - Inches(0.5), card_y + Inches(0.35), Inches(0.35), Inches(0.4),
                     "\u2713", font_size=18, color=ACCENT_GREEN, bold=True, alignment=PP_ALIGN.CENTER)

    # The gap — question mark card
    gap_x = start_x + 4 * (card_w + gap)
    gap_card_w = Inches(1.8)
    box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, gap_x, card_y, gap_card_w, card_h)
    box.fill.solid()
    box.fill.fore_color.rgb = RGBColor(0xFF, 0xF3, 0xF3)
    box.line.color.rgb = ACCENT_CORAL
    box.line.width = Pt(1.5)
    box.adjustments[0] = 0.06

    add_text_box(slide, gap_x, card_y + Inches(0.15), gap_card_w, Inches(0.6),
                 "?", font_size=36, color=ACCENT_CORAL, bold=True,
                 alignment=PP_ALIGN.CENTER, font_name=HEADLINE_FONT)
    add_text_box(slide, gap_x + Inches(0.15), card_y + Inches(0.75), gap_card_w - Inches(0.3), Inches(0.6),
                 "Mental Health", font_size=14, color=ACCENT_CORAL, bold=True,
                 alignment=PP_ALIGN.CENTER)

    # Bottom insight
    add_text_box(slide, Inches(0.8), Inches(4.8), Inches(11.5), Inches(0.6),
                 "These are high-performers with disposable income who obsessively quantify every body metric.",
                 font_size=15, color=BODY_GRAY, line_spacing=1.4)

    add_text_box(slide, Inches(0.8), Inches(5.35), Inches(11.5), Inches(0.6),
                 "Their only mental health tools? Subjective self-report questionnaires, mood journals, and therapy once a week.",
                 font_size=15, color=BODY_GRAY, line_spacing=1.4)

    # Bottom accent bar
    add_divider_line(slide, Inches(0.8), Inches(6.2), Inches(11.7), RGBColor(0xE8, 0xE8, 0xF0), Pt(1))

    add_text_box(slide, Inches(0.8), Inches(6.35), Inches(11.5), Inches(0.4),
                 "They have zero way to objectively, continuously track their mental health.",
                 font_size=16, color=CHARCOAL, bold=True, line_spacing=1.3)

    add_slide_number(slide, 2)


def slide_03_qs_market(prs):
    """Slide 3 — The Quantified-Self Market (LIGHT background). NEW in v3."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, LIGHT_BG)

    add_text_box(slide, Inches(0.8), Inches(0.5), Inches(3.0), Inches(0.35),
                 "THE MARKET", font_size=11, color=STEEL, bold=True)
    add_divider_line(slide, Inches(0.8), Inches(0.9), Inches(2.0), STEEL, Pt(2))

    add_text_box(slide, Inches(0.8), Inches(1.1), Inches(11.5), Inches(0.8),
                 "The Quantified-Self Market", font_size=36, color=CHARCOAL,
                 bold=True, font_name=HEADLINE_FONT)

    add_text_box(slide, Inches(0.8), Inches(1.85), Inches(11.5), Inches(0.5),
                 "150-200M people are already paying for health data. Lucid adds the missing metric.",
                 font_size=17, color=BODY_GRAY, line_spacing=1.4)

    # Wearable brand grid — table-like layout
    brands = [
        ("Apple Watch", "200M+ global users (~32M US)", "$249 \u2013 $799", ACCENT_CORAL),
        ("Oura Ring", "5.5M+ sold \u2022 $11B valuation", "$349 \u2013 $499 + $5.99/mo", RGBColor(0x7B, 0x68, 0xD4)),
        ("WHOOP", "~1\u20132M subscribers", "$149 \u2013 $359/yr", ACCENT_GREEN),
        ("Garmin", "~20\u201330M fitness users", "$549 \u2013 $750", ACCENT_GOLD),
    ]

    card_w = Inches(5.5)
    card_h = Inches(0.75)
    grid_x = Inches(0.8)
    grid_y = Inches(2.6)
    v_gap = Inches(0.15)

    for i, (brand, users, price, accent) in enumerate(brands):
        y = grid_y + i * (card_h + v_gap)

        # Card background
        box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, grid_x, y, card_w, card_h)
        box.fill.solid()
        box.fill.fore_color.rgb = CARD_BG
        box.line.color.rgb = CARD_BORDER
        box.line.width = Pt(1)
        box.adjustments[0] = 0.08

        # Accent dot
        add_accent_dot(slide, grid_x + Inches(0.2), y + Inches(0.12), Inches(0.12), accent)

        # Brand name
        add_text_box(slide, grid_x + Inches(0.45), y + Inches(0.05), Inches(1.6), Inches(0.3),
                     brand, font_size=14, color=CHARCOAL, bold=True)

        # Users
        add_text_box(slide, grid_x + Inches(0.45), y + Inches(0.38), Inches(3.0), Inches(0.3),
                     users, font_size=11, color=BODY_GRAY)

        # Price — right side of card
        add_text_box(slide, grid_x + Inches(3.2), y + Inches(0.05), Inches(2.1), Inches(0.3),
                     price, font_size=12, color=STEEL, bold=True, alignment=PP_ALIGN.RIGHT)

    # RIGHT side — summary stats
    right_x = Inches(7.0)

    # Big number callout
    add_text_box(slide, right_x, Inches(2.7), Inches(5.5), Inches(0.7),
                 "150\u2013200M", font_size=48, color=STEEL, bold=True,
                 font_name=HEADLINE_FONT)

    add_text_box(slide, right_x, Inches(3.45), Inches(5.5), Inches(0.5),
                 "people actively tracking HRV, sleep, and recovery",
                 font_size=15, color=CHARCOAL, bold=True, line_spacing=1.3)

    add_text_box(slide, right_x, Inches(4.1), Inches(5.5), Inches(0.8),
                 "These users spend $249\u2013$799 on hardware + $6\u201330/mo on subscriptions. They are already paying for health data.",
                 font_size=13, color=BODY_GRAY, line_spacing=1.5)

    # Lucid positioning card
    lucid_y = Inches(5.2)
    lucid_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                         right_x, lucid_y, Inches(5.5), Inches(0.8))
    lucid_box.fill.solid()
    lucid_box.fill.fore_color.rgb = NAVY
    lucid_box.line.fill.background()
    lucid_box.adjustments[0] = 0.08

    add_text_box(slide, right_x + Inches(0.3), lucid_y + Inches(0.1), Inches(2.5), Inches(0.35),
                 "Lucid", font_size=18, color=WHITE, bold=True, font_name=HEADLINE_FONT)
    add_text_box(slide, right_x + Inches(3.0), lucid_y + Inches(0.1), Inches(2.2), Inches(0.35),
                 "$14.99/mo", font_size=18, color=STEEL, bold=True,
                 alignment=PP_ALIGN.RIGHT, font_name=HEADLINE_FONT)
    add_text_box(slide, right_x + Inches(0.3), lucid_y + Inches(0.45), Inches(4.8), Inches(0.3),
                 "The mental health metric their stack is missing", font_size=11,
                 color=RGBColor(0xAA, 0xBB, 0xCC))

    # Bottom callout
    add_divider_line(slide, Inches(0.8), Inches(6.3), Inches(11.7), RGBColor(0xE8, 0xE8, 0xF0), Pt(1))
    add_text_box(slide, Inches(0.8), Inches(6.45), Inches(11.5), Inches(0.4),
                 "Zero voice-first mental health players in this ecosystem. Lucid is the first.",
                 font_size=16, color=CHARCOAL, bold=True)

    add_slide_number(slide, 3)


def slide_04_measurement_gap(prs):
    """Slide 4 — The Measurement Gap (LIGHT background)."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, LIGHT_BG)

    add_text_box(slide, Inches(0.8), Inches(0.5), Inches(3.0), Inches(0.35),
                 "THE PROBLEM", font_size=11, color=STEEL, bold=True)
    add_divider_line(slide, Inches(0.8), Inches(0.9), Inches(2.0), STEEL, Pt(2))

    add_text_box(slide, Inches(0.8), Inches(1.1), Inches(11.5), Inches(0.8),
                 "You Can't Improve What You Can't Measure", font_size=36,
                 color=CHARCOAL, bold=True, font_name=HEADLINE_FONT)

    # Two columns: Body vs Mind
    col_w = Inches(5.2)
    col_h = Inches(3.6)
    left_x = Inches(0.8)
    right_x = Inches(6.8)
    col_y = Inches(2.2)

    # LEFT: Body — Quantified
    box_l = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left_x, col_y, col_w, col_h)
    box_l.fill.solid()
    box_l.fill.fore_color.rgb = CARD_BG
    box_l.line.color.rgb = ACCENT_GREEN
    box_l.line.width = Pt(1.5)
    box_l.adjustments[0] = 0.04

    add_text_box(slide, left_x + Inches(0.4), col_y + Inches(0.25), Inches(4.0), Inches(0.4),
                 "Body: Fully Quantified", font_size=20, color=ACCENT_GREEN, bold=True,
                 font_name=HEADLINE_FONT)

    body_metrics = [
        ("\u2713  Heart rate \u2014 continuous, real-time", ACCENT_GREEN),
        ("\u2713  Sleep stages \u2014 automatic, nightly", ACCENT_GREEN),
        ("\u2713  HRV \u2014 objective stress proxy", ACCENT_GREEN),
        ("\u2713  VO2 max \u2014 cardio fitness trend", ACCENT_GREEN),
        ("\u2713  Steps / activity \u2014 passive, 24/7", ACCENT_GREEN),
        ("\u2713  Blood oxygen \u2014 overnight SpO2", ACCENT_GREEN),
    ]

    for j, (txt, col) in enumerate(body_metrics):
        add_text_box(slide, left_x + Inches(0.4), col_y + Inches(0.8 + j * 0.42),
                     Inches(4.4), Inches(0.4), txt, font_size=13, color=BODY_GRAY)

    # RIGHT: Mind — Unmeasured
    box_r = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, right_x, col_y, col_w, col_h)
    box_r.fill.solid()
    box_r.fill.fore_color.rgb = CARD_BG
    box_r.line.color.rgb = ACCENT_CORAL
    box_r.line.width = Pt(1.5)
    box_r.adjustments[0] = 0.04

    add_text_box(slide, right_x + Inches(0.4), col_y + Inches(0.25), Inches(4.0), Inches(0.4),
                 "Mind: Unmeasured", font_size=20, color=ACCENT_CORAL, bold=True,
                 font_name=HEADLINE_FONT)

    mind_metrics = [
        ("?   Stress \u2014 subjective 1-10 scale", ACCENT_CORAL),
        ("?   Mood \u2014 self-report questionnaire", ACCENT_CORAL),
        ("?   Anxiety \u2014 PHQ / GAD forms", ACCENT_CORAL),
        ("?   Depression \u2014 episodic screening", ACCENT_CORAL),
        ("?   Cognitive load \u2014 not measured", ACCENT_CORAL),
        ("?   Mental fatigue \u2014 not measured", ACCENT_CORAL),
    ]

    for j, (txt, col) in enumerate(mind_metrics):
        add_text_box(slide, right_x + Inches(0.4), col_y + Inches(0.8 + j * 0.42),
                     Inches(4.4), Inches(0.4), txt, font_size=13, color=BODY_GRAY)

    # Bottom callout
    callout_y = Inches(6.1)
    add_divider_line(slide, Inches(0.8), callout_y, Inches(11.7), RGBColor(0xE8, 0xE8, 0xF0), Pt(1))

    add_text_box(slide, Inches(0.8), callout_y + Inches(0.15), Inches(11.5), Inches(0.5),
                 "This is the biggest blind spot in the quantified-self stack.",
                 font_size=17, color=CHARCOAL, bold=True)

    add_slide_number(slide, 4)


def slide_05_options_failure(prs):
    """Slide 5 — Why Current Mental Health Solutions Fail (LIGHT background). NEW in v3."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, LIGHT_BG)

    add_text_box(slide, Inches(0.8), Inches(0.5), Inches(3.0), Inches(0.35),
                 "THE PROBLEM", font_size=11, color=STEEL, bold=True)
    add_divider_line(slide, Inches(0.8), Inches(0.9), Inches(2.0), STEEL, Pt(2))

    add_text_box(slide, Inches(0.8), Inches(1.1), Inches(11.5), Inches(0.8),
                 "Why Current Mental Health Solutions Fail",
                 font_size=36, color=CHARCOAL, bold=True, font_name=HEADLINE_FONT)

    add_text_box(slide, Inches(0.8), Inches(1.85), Inches(11.5), Inches(0.5),
                 "The QS audience has tried everything. Nothing sticks.",
                 font_size=17, color=BODY_GRAY, line_spacing=1.4)

    # Three columns showing failure modes
    col_w = Inches(3.6)
    col_h = Inches(3.8)
    start_x = Inches(0.8)
    col_y = Inches(2.6)
    col_gap = Inches(0.4)

    columns = [
        ("Therapy", ACCENT_CORAL, [
            ("50%", "drop out before treatment ends"),
            ("2.4", "sessions before quitting"),
            ("48 days", "wait for first appointment"),
            ("$150\u2013300", "per session, often out-of-network"),
        ]),
        ("Mental Health Apps", RGBColor(0x7B, 0x68, 0xD4), [
            ("3.3%", "median 30-day retention"),
            ("7.65%", "Headspace retention"),
            ("8.34%", "Calm retention"),
            ("Woebot", "shut down consumer app (2025)"),
        ]),
        ("Wearables", STEEL, [
            ("170K", "WHOOP mental health surveys"),
            ("Oura", "launched stress tracking"),
            ("Garmin", "added stress via HRV"),
            ("Demand", "is there \u2014 tools aren't"),
        ]),
    ]

    for i, (title, accent, stats) in enumerate(columns):
        x = start_x + i * (col_w + col_gap)

        # Card
        box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, col_y, col_w, col_h)
        box.fill.solid()
        box.fill.fore_color.rgb = CARD_BG
        box.line.color.rgb = CARD_BORDER
        box.line.width = Pt(1)
        box.adjustments[0] = 0.05

        # Accent bar
        add_divider_line(slide, x + Inches(0.3), col_y + Inches(0.2), Inches(0.8), accent, Pt(3))

        # Column title
        add_text_box(slide, x + Inches(0.3), col_y + Inches(0.4), col_w - Inches(0.6), Inches(0.4),
                     title, font_size=18, color=CHARCOAL, bold=True, font_name=HEADLINE_FONT)

        # Stats
        for j, (stat, desc) in enumerate(stats):
            stat_y = col_y + Inches(1.0 + j * 0.68)
            add_text_box(slide, x + Inches(0.3), stat_y, col_w - Inches(0.6), Inches(0.3),
                         stat, font_size=18, color=accent, bold=True, font_name=HEADLINE_FONT)
            add_text_box(slide, x + Inches(0.3), stat_y + Inches(0.28), col_w - Inches(0.6), Inches(0.35),
                         desc, font_size=11, color=BODY_GRAY, line_spacing=1.2)

    # Bottom callout
    add_divider_line(slide, Inches(0.8), Inches(6.65), Inches(11.7), RGBColor(0xE8, 0xE8, 0xF0), Pt(1))
    add_text_box(slide, Inches(0.8), Inches(6.8), Inches(11.5), Inches(0.4),
                 "Quantified-self users want objective mental health data. Nothing on the market delivers it.",
                 font_size=16, color=CHARCOAL, bold=True)

    add_slide_number(slide, 5)


def slide_06_solution(prs):
    """Slide 6 — The Solution (LIGHT background)."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, LIGHT_BG)

    add_text_box(slide, Inches(0.8), Inches(0.5), Inches(3.0), Inches(0.35),
                 "THE SOLUTION", font_size=11, color=STEEL, bold=True)
    add_divider_line(slide, Inches(0.8), Inches(0.9), Inches(2.0), STEEL, Pt(2))

    # Left side content
    add_text_box(slide, Inches(0.8), Inches(1.2), Inches(6.0), Inches(0.8),
                 "Your Voice Is a Biomarker", font_size=36, color=CHARCOAL,
                 bold=True, font_name=HEADLINE_FONT)

    add_text_box(slide, Inches(0.8), Inches(2.05), Inches(5.5), Inches(0.6),
                 "25 seconds of natural speech \u2192 clinical-grade mental health signals",
                 font_size=17, color=STEEL, bold=True, line_spacing=1.4)

    # Key bullets
    bullets = [
        ("Passive", "No questionnaires, no self-reporting"),
        ("Continuous", "Track trends over days, weeks, months"),
        ("Objective", "Clinical-grade biomarker analysis, not subjective mood logs"),
        ("Private", "100% on-device processing \u2014 your voice never leaves your phone"),
    ]

    for i, (title, desc) in enumerate(bullets):
        y = Inches(2.9 + i * 0.75)
        add_accent_dot(slide, Inches(0.9), y + Inches(0.05), Inches(0.1), STEEL)
        add_text_box(slide, Inches(1.15), y, Inches(5.0), Inches(0.35),
                     title, font_size=15, color=CHARCOAL, bold=True)
        add_text_box(slide, Inches(1.15), y + Inches(0.3), Inches(5.0), Inches(0.35),
                     desc, font_size=12, color=BODY_GRAY)

    # Right side — dashboard screenshot, proper aspect ratio
    dash_w, dash_h = fit_image_dims("dashboard", 5.5, 4.0)
    dash_left = Inches(7.2)
    dash_top = Inches(1.8)
    if os.path.exists(IMG_DASHBOARD):
        # Subtle shadow frame
        shadow = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                        dash_left - Inches(0.05), dash_top - Inches(0.05),
                                        dash_w + Inches(0.1), dash_h + Inches(0.1))
        shadow.fill.solid()
        shadow.fill.fore_color.rgb = RGBColor(0xE0, 0xE0, 0xE8)
        shadow.line.fill.background()
        shadow.adjustments[0] = 0.03

        slide.shapes.add_picture(IMG_DASHBOARD, dash_left, dash_top, dash_w, dash_h)

    add_slide_number(slide, 6)


def slide_07_how_it_works(prs):
    """Slide 7 — How It Works (LIGHT background)."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, LIGHT_BG)

    add_text_box(slide, Inches(0.8), Inches(0.5), Inches(3.0), Inches(0.35),
                 "HOW IT WORKS", font_size=11, color=STEEL, bold=True)
    add_divider_line(slide, Inches(0.8), Inches(0.9), Inches(2.0), STEEL, Pt(2))

    add_text_box(slide, Inches(0.8), Inches(1.1), Inches(11.5), Inches(0.8),
                 "Three Steps. 25 Seconds. Zero Effort.", font_size=36,
                 color=CHARCOAL, bold=True, font_name=HEADLINE_FONT)

    steps = [
        ("1", "Install", "Download the app.\nGrant microphone access.\nThat's it.",
         STEEL, "\u2B07"),
        ("2", "Speak", "Talk naturally for 25 seconds.\nRead a prompt or speak freely.\nThe AI listens to how you sound, not what you say.",
         ACCENT_GREEN, "\U0001F399"),
        ("3", "Discover", "Get your Wellness Score instantly.\nTrack depression risk, stress, mood, energy.\nSee trends over time on your personal dashboard.",
         ACCENT_GOLD, "\U0001F4CA"),
    ]

    card_w = Inches(3.5)
    card_h = Inches(3.5)
    start_x = Inches(0.8)
    card_y = Inches(2.3)
    gap = Inches(0.6)

    for i, (num, title, body, accent, icon) in enumerate(steps):
        x = start_x + i * (card_w + gap)

        # Card
        box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, card_y, card_w, card_h)
        box.fill.solid()
        box.fill.fore_color.rgb = CARD_BG
        box.line.color.rgb = CARD_BORDER
        box.line.width = Pt(1)
        box.adjustments[0] = 0.06

        # Step number circle
        circle = slide.shapes.add_shape(MSO_SHAPE.OVAL,
                                        x + Inches(0.3), card_y + Inches(0.3),
                                        Inches(0.5), Inches(0.5))
        circle.fill.solid()
        circle.fill.fore_color.rgb = accent
        circle.line.fill.background()

        add_text_box(slide, x + Inches(0.3), card_y + Inches(0.32),
                     Inches(0.5), Inches(0.5),
                     num, font_size=20, color=WHITE, bold=True,
                     alignment=PP_ALIGN.CENTER, font_name=HEADLINE_FONT)

        # Title
        add_text_box(slide, x + Inches(0.3), card_y + Inches(1.0),
                     card_w - Inches(0.6), Inches(0.4),
                     title, font_size=20, color=CHARCOAL, bold=True,
                     font_name=HEADLINE_FONT)

        # Body
        add_text_box(slide, x + Inches(0.3), card_y + Inches(1.5),
                     card_w - Inches(0.6), Inches(1.8),
                     body, font_size=13, color=BODY_GRAY, line_spacing=1.5)

    # Connector arrows between cards
    for i in range(2):
        arrow_x = start_x + (i + 1) * card_w + (i * gap) + gap * 0.3
        add_text_box(slide, arrow_x, card_y + Inches(1.5), Inches(0.4), Inches(0.5),
                     "\u2192", font_size=28, color=MEDIUM_GRAY,
                     alignment=PP_ALIGN.CENTER, bold=True)

    add_slide_number(slide, 7)


def slide_08_clinical_validation(prs):
    """Slide 8 — Clinical Validation (LIGHT background). Fixed overlap."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, LIGHT_BG)

    add_text_box(slide, Inches(0.8), Inches(0.5), Inches(5.0), Inches(0.35),
                 "CLINICAL VALIDATION", font_size=11, color=STEEL, bold=True)
    add_divider_line(slide, Inches(0.8), Inches(0.9), Inches(2.5), STEEL, Pt(2))

    add_text_box(slide, Inches(0.8), Inches(1.1), Inches(7.0), Inches(0.8),
                 "The Largest Voice Biomarker\nStudy Ever Published", font_size=32,
                 color=CHARCOAL, bold=True, font_name=HEADLINE_FONT, line_spacing=1.15)

    # Three stat boxes — LEFT column, stacked vertically
    stats = [
        ("14,898", "Adult Participants", STEEL),
        ("71.3%", "Sensitivity", ACCENT_GREEN),
        ("73.5%", "Specificity", ACCENT_GOLD),
    ]

    stat_x = Inches(0.8)
    stat_w = Inches(2.8)
    stat_h = Inches(1.15)
    stat_start_y = Inches(2.5)

    for i, (num, label, accent) in enumerate(stats):
        y = stat_start_y + i * (stat_h + Inches(0.25))
        box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, stat_x, y, stat_w, stat_h)
        box.fill.solid()
        box.fill.fore_color.rgb = CARD_BG
        box.line.color.rgb = CARD_BORDER
        box.line.width = Pt(1)
        box.adjustments[0] = 0.06

        add_text_box(slide, stat_x, y + Inches(0.1), stat_w, Inches(0.5),
                     num, font_size=36, color=accent, bold=True,
                     alignment=PP_ALIGN.CENTER, font_name=HEADLINE_FONT)
        add_text_box(slide, stat_x, y + Inches(0.65), stat_w, Inches(0.35),
                     label, font_size=13, color=BODY_GRAY,
                     alignment=PP_ALIGN.CENTER)

    # Additional context below stats — FIXED: moved down to Inches(6.3) to avoid overlap
    add_text_box(slide, Inches(0.8), Inches(6.3), Inches(4.5), Inches(0.5),
                 "Peer-reviewed \u2022 PubMed-indexed \u2022 Multi-site validated",
                 font_size=12, color=MEDIUM_GRAY, line_spacing=1.3)

    # RIGHT side — PubMed image, portrait aspect ratio
    pub_max_w = 3.8
    pub_max_h = 4.0
    pub_w, pub_h = fit_image_dims("pubmed", pub_max_w, pub_max_h)
    pub_left = Inches(8.5)
    pub_top = Inches(1.5)

    if os.path.exists(IMG_PUBMED):
        # Frame
        frame = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                       pub_left - Inches(0.08), pub_top - Inches(0.08),
                                       pub_w + Inches(0.16), pub_h + Inches(0.16))
        frame.fill.solid()
        frame.fill.fore_color.rgb = RGBColor(0xE0, 0xE0, 0xE8)
        frame.line.fill.background()
        frame.adjustments[0] = 0.02

        slide.shapes.add_picture(IMG_PUBMED, pub_left, pub_top, pub_w, pub_h)

    # Citation below image — moved down proportionally
    cite_top = pub_top + pub_h + Inches(0.2)
    add_text_box(slide, pub_left - Inches(0.5), cite_top, Inches(4.8), Inches(0.6),
                 "Mazur et al. (2025) Annals of Family Medicine\nDOI: 10.1370/afm.240091",
                 font_size=10, color=MEDIUM_GRAY, alignment=PP_ALIGN.CENTER,
                 line_spacing=1.3)

    add_slide_number(slide, 8)


def slide_09_technology(prs):
    """Slide 9 — The Technology (NAVY background). Updated Kintsugi narrative."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, NAVY)

    add_text_box(slide, Inches(0.8), Inches(0.5), Inches(3.0), Inches(0.35),
                 "THE TECHNOLOGY", font_size=11, color=STEEL, bold=True)
    add_divider_line(slide, Inches(0.8), Inches(0.9), Inches(2.0), STEEL, Pt(2))

    add_text_box(slide, Inches(0.8), Inches(1.2), Inches(11.5), Inches(0.8),
                 "$30M of Clinical R&D. In Your Pocket.",
                 font_size=36, color=WHITE, bold=True, font_name=HEADLINE_FONT)

    add_text_box(slide, Inches(0.8), Inches(2.1), Inches(11.5), Inches(0.6),
                 "Built on the Depression Assessment Model (DAM) \u2014 the most validated voice biomarker AI in existence.",
                 font_size=16, color=RGBColor(0xAA, 0xBB, 0xCC), line_spacing=1.4)

    # Four stat boxes
    tech_stats = [
        ("$30M", "Clinical R&D\nInvestment"),
        ("244M", "Trainable\nParameters"),
        ("863 hrs", "Clinical Voice\nData"),
        ("35,000", "Patient\nStudy Pool"),
    ]

    box_w = Inches(2.5)
    box_h = Inches(2.2)
    start_x = Inches(0.8)
    box_y = Inches(3.2)
    gap = Inches(0.5)

    for i, (num, label) in enumerate(tech_stats):
        x = start_x + i * (box_w + gap)
        box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, box_y, box_w, box_h)
        box.fill.solid()
        box.fill.fore_color.rgb = RGBColor(0x22, 0x44, 0x5C)
        box.line.color.rgb = RGBColor(0x33, 0x55, 0x6E)
        box.line.width = Pt(1)
        box.adjustments[0] = 0.06

        add_text_box(slide, x, box_y + Inches(0.35), box_w, Inches(0.7),
                     num, font_size=40, color=STEEL, bold=True,
                     alignment=PP_ALIGN.CENTER, font_name=HEADLINE_FONT)
        add_text_box(slide, x, box_y + Inches(1.2), box_w, Inches(0.8),
                     label, font_size=14, color=RGBColor(0xAA, 0xBB, 0xCC),
                     alignment=PP_ALIGN.CENTER, line_spacing=1.3)

    # Bottom note — UPDATED: Kintsugi opportunity narrative
    add_text_box(slide, Inches(0.8), Inches(5.7), Inches(11.5), Inches(0.8),
                 "Kintsugi Health (2020\u20132026) invested $30M building this technology, then shut down \u2014 "
                 "killed by FDA regulatory economics, not bad science. They open-sourced everything. "
                 "Lucid inherits validated, peer-reviewed voice biomarker AI without the regulatory trap.",
                 font_size=12, color=RGBColor(0xAA, 0xBB, 0xCC), alignment=PP_ALIGN.LEFT,
                 line_spacing=1.4)

    add_slide_number(slide, 9, color=RGBColor(0x66, 0x77, 0x88))


def slide_10_product_deep_dive(prs):
    """Slide 10 — Product Deep Dive (LIGHT background). Fixed spacing."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, LIGHT_BG)

    add_text_box(slide, Inches(0.8), Inches(0.5), Inches(3.0), Inches(0.35),
                 "PRODUCT", font_size=11, color=STEEL, bold=True)
    add_divider_line(slide, Inches(0.8), Inches(0.9), Inches(2.0), STEEL, Pt(2))

    add_text_box(slide, Inches(0.8), Inches(1.1), Inches(11.5), Inches(0.8),
                 "A Complete Mental Wellness Platform", font_size=32,
                 color=CHARCOAL, bold=True, font_name=HEADLINE_FONT)

    # 6 feature cards in 2x3 grid — LEFT side
    features = [
        ("Wellness Score", "Composite 0-100 score from voice biomarkers"),
        ("Depression Screening", "PHQ-8 validated clinical screening"),
        ("Stress Detection", "Real-time stress level from vocal patterns"),
        ("Mood Tracking", "Longitudinal mood trends with AI analysis"),
        ("Energy Assessment", "Vocal energy and fatigue indicators"),
        ("Trend Dashboard", "Historical trends with daily/weekly views"),
    ]

    card_w = Inches(2.4)
    card_h = Inches(1.35)
    grid_x = Inches(0.8)
    grid_y = Inches(2.0)
    h_gap = Inches(0.25)
    v_gap = Inches(0.25)

    for i, (title, body) in enumerate(features):
        row = i // 2
        col = i % 2
        x = grid_x + col * (card_w + h_gap)
        y = grid_y + row * (card_h + v_gap)
        add_card(slide, x, y, card_w, card_h, title, body, body_size=11)

    # Dashboard image — RIGHT side — FIXED: moved right to Inches(6.5) for more gap
    dash_w, dash_h = fit_image_dims("dashboard", 5.5, 4.5)
    dash_left = Inches(6.5)
    dash_top = Inches(2.0)

    if os.path.exists(IMG_DASHBOARD):
        frame = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                       dash_left - Inches(0.05), dash_top - Inches(0.05),
                                       dash_w + Inches(0.1), dash_h + Inches(0.1))
        frame.fill.solid()
        frame.fill.fore_color.rgb = RGBColor(0xE0, 0xE0, 0xE8)
        frame.line.fill.background()
        frame.adjustments[0] = 0.03

        slide.shapes.add_picture(IMG_DASHBOARD, dash_left, dash_top, dash_w, dash_h)

    add_slide_number(slide, 10)


def slide_11_market(prs):
    """Slide 11 — Market Opportunity (LIGHT background). Redesigned as horizontal bars."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, LIGHT_BG)

    add_text_box(slide, Inches(0.8), Inches(0.5), Inches(3.0), Inches(0.35),
                 "MARKET", font_size=11, color=STEEL, bold=True)
    add_divider_line(slide, Inches(0.8), Inches(0.9), Inches(2.0), STEEL, Pt(2))

    add_text_box(slide, Inches(0.8), Inches(1.1), Inches(11.5), Inches(0.8),
                 "A $250B Market with Zero Voice-First Players",
                 font_size=32, color=CHARCOAL, bold=True, font_name=HEADLINE_FONT)

    # Horizontal stacked bars — TAM / SAM / SOM
    bar_x = Inches(0.8)
    bar_start_y = Inches(2.4)
    max_bar_w = Inches(11.0)
    bar_h = Inches(1.1)
    bar_gap = Inches(0.3)

    bars = [
        ("TAM", "$250B", "Global Mental Health Market", max_bar_w,
         RGBColor(0xD8, 0xE8, 0xF4), CHARCOAL,
         "Therapy, pharmaceuticals, digital therapeutics, workplace wellness, telehealth. 1B+ people globally."),
        ("SAM", "$38B", "Digital Mental Health", Inches(7.0),
         RGBColor(0xB8, 0xD4, 0xE8), CHARCOAL,
         "Apps, platforms, and digital tools. Growing ~18% CAGR."),
        ("SOM", "$4.2B", "Consumer Voice Wellness", Inches(3.5),
         STEEL, WHITE,
         "Voice-based wellness for QS users: Whoop, Oura, Apple Watch owners."),
    ]

    for i, (tier, amount, name, bar_w, bg_color, text_color, desc) in enumerate(bars):
        y = bar_start_y + i * (bar_h + bar_gap)

        # Bar shape
        bar = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, bar_x, y, bar_w, bar_h)
        bar.fill.solid()
        bar.fill.fore_color.rgb = bg_color
        bar.line.fill.background()
        bar.adjustments[0] = 0.1

        # Tier label + amount — constrained to bar width
        add_text_box(slide, bar_x + Inches(0.3), y + Inches(0.08), Inches(0.8), Inches(0.3),
                     tier, font_size=11, color=text_color, bold=True)
        add_text_box(slide, bar_x + Inches(1.1), y + Inches(0.03), Inches(1.8), Inches(0.4),
                     amount, font_size=26, color=text_color, bold=True, font_name=HEADLINE_FONT)
        add_text_box(slide, bar_x + Inches(0.3), y + Inches(0.48), Inches(2.5), Inches(0.3),
                     name, font_size=12, color=text_color, bold=True)

        # Description — inside for wide bars, to the right for narrow bars
        if bar_w > Inches(5.0):
            add_text_box(slide, bar_x + Inches(3.5), y + Inches(0.12), bar_w - Inches(4.0), Inches(0.85),
                         desc, font_size=11, color=BODY_GRAY, line_spacing=1.3)
        else:
            # Narrow bar — description goes to the right of labels
            add_text_box(slide, bar_x + Inches(3.5), y + Inches(0.12), Inches(8.0), Inches(0.85),
                         desc, font_size=11, color=BODY_GRAY, line_spacing=1.3)

    # Bottom callout
    callout_y = Inches(6.3)
    add_divider_line(slide, Inches(0.8), callout_y, Inches(11.7), RGBColor(0xE8, 0xE8, 0xF0), Pt(1))
    add_text_box(slide, Inches(0.8), callout_y + Inches(0.15), Inches(11.5), Inches(0.5),
                 "Lucid is the only voice-first, clinically validated player targeting the consumer QS segment.",
                 font_size=16, color=CHARCOAL, bold=True)

    add_slide_number(slide, 11)


def slide_12_business_model(prs):
    """Slide 12 — Business Model (LIGHT background). Updated pricing to $14.99/mo."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, LIGHT_BG)

    add_text_box(slide, Inches(0.8), Inches(0.5), Inches(3.0), Inches(0.35),
                 "BUSINESS MODEL", font_size=11, color=STEEL, bold=True)
    add_divider_line(slide, Inches(0.8), Inches(0.9), Inches(2.5), STEEL, Pt(2))

    add_text_box(slide, Inches(0.8), Inches(1.1), Inches(11.5), Inches(0.8),
                 "Premium Subscription with Enterprise Expansion",
                 font_size=32, color=CHARCOAL, bold=True, font_name=HEADLINE_FONT)

    # LEFT — Pricing card (navy card as contained element)
    price_x = Inches(0.8)
    price_y = Inches(2.2)
    price_w = Inches(4.5)
    price_h = Inches(4.3)

    box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, price_x, price_y, price_w, price_h)
    box.fill.solid()
    box.fill.fore_color.rgb = NAVY
    box.line.fill.background()
    box.adjustments[0] = 0.05

    add_text_box(slide, price_x + Inches(0.4), price_y + Inches(0.3),
                 Inches(3.5), Inches(0.4),
                 "Lucid Premium", font_size=14, color=STEEL, bold=True)

    # UPDATED: $14.99/mo
    add_text_box(slide, price_x + Inches(0.4), price_y + Inches(0.7),
                 Inches(3.5), Inches(0.6),
                 "$14.99/mo", font_size=40, color=WHITE, bold=True,
                 font_name=HEADLINE_FONT)

    # UPDATED: $119.99/year
    add_text_box(slide, price_x + Inches(0.4), price_y + Inches(1.3),
                 Inches(3.5), Inches(0.3),
                 "or $119.99/year (save 33%)", font_size=13,
                 color=RGBColor(0xAA, 0xBB, 0xCC))

    price_features = [
        "\u2713  Unlimited voice analyses",
        "\u2713  Clinical wellness scoring",
        "\u2713  Depression risk screening",
        "\u2713  Stress & mood tracking",
        "\u2713  Trend analytics dashboard",
        "\u2713  Export health reports",
        "\u2713  100% on-device processing",
    ]

    for j, feat in enumerate(price_features):
        add_text_box(slide, price_x + Inches(0.4), price_y + Inches(1.85 + j * 0.32),
                     Inches(3.5), Inches(0.3),
                     feat, font_size=12, color=RGBColor(0xCC, 0xDD, 0xEE))

    # RIGHT — Expansion revenue
    exp_x = Inches(6.2)

    add_text_box(slide, exp_x, Inches(2.2), Inches(6.0), Inches(0.4),
                 "Revenue Expansion Paths", font_size=20, color=CHARCOAL,
                 bold=True, font_name=HEADLINE_FONT)

    expansions = [
        ("B2B2C: Employer Wellness", "White-label integration with corporate wellness platforms. Per-employee licensing ($5-15/employee/mo)."),
        ("Health System Partnerships", "EHR integration for mental health screening at scale. Revenue-share model with health systems."),
        ("Insurance / Payer Integration", "Risk-adjusted premium discounts for continuous mental health monitoring. Data licensing (anonymized, aggregated)."),
        ("Platform Licensing", "License the voice biomarker engine to telehealth platforms, EAPs, and digital therapeutics companies."),
    ]

    for i, (title, desc) in enumerate(expansions):
        y = Inches(2.8 + i * 1.0)
        add_accent_dot(slide, exp_x, y + Inches(0.05), Inches(0.1), STEEL)
        add_text_box(slide, exp_x + Inches(0.25), y, Inches(6.0), Inches(0.3),
                     title, font_size=14, color=CHARCOAL, bold=True)
        add_text_box(slide, exp_x + Inches(0.25), y + Inches(0.3), Inches(6.0), Inches(0.55),
                     desc, font_size=11, color=BODY_GRAY, line_spacing=1.3)

    add_slide_number(slide, 12)


def slide_13_competitive_landscape(prs):
    """Slide 13 — Competitive Landscape (LIGHT background). Punchier with dead competitors."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, LIGHT_BG)

    add_text_box(slide, Inches(0.8), Inches(0.5), Inches(5.0), Inches(0.35),
                 "COMPETITIVE LANDSCAPE", font_size=11, color=STEEL, bold=True)
    add_divider_line(slide, Inches(0.8), Inches(0.9), Inches(2.5), STEEL, Pt(2))

    add_text_box(slide, Inches(0.8), Inches(1.1), Inches(11.5), Inches(0.8),
                 "Clinical Depth \u00D7 Consumer Accessibility",
                 font_size=32, color=CHARCOAL, bold=True, font_name=HEADLINE_FONT)

    # 2x2 Matrix
    matrix_x = Inches(1.5)
    matrix_y = Inches(2.2)
    matrix_w = Inches(7.0)
    matrix_h = Inches(4.5)

    # Background quadrants
    quads = [
        (matrix_x, matrix_y, Inches(3.5), Inches(2.25), RGBColor(0xF0, 0xF4, 0xF8)),
        (matrix_x + Inches(3.5), matrix_y, Inches(3.5), Inches(2.25), RGBColor(0xE8, 0xF0, 0xF8)),
        (matrix_x, matrix_y + Inches(2.25), Inches(3.5), Inches(2.25), RGBColor(0xF4, 0xF4, 0xF4)),
        (matrix_x + Inches(3.5), matrix_y + Inches(2.25), Inches(3.5), Inches(2.25), RGBColor(0xF0, 0xF0, 0xF4)),
    ]
    for qx, qy, qw, qh, qc in quads:
        q = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, qx, qy, qw, qh)
        q.fill.solid()
        q.fill.fore_color.rgb = qc
        q.line.color.rgb = CARD_BORDER
        q.line.width = Pt(0.5)

    # Axis labels
    add_text_box(slide, matrix_x + Inches(1.0), matrix_y + matrix_h + Inches(0.1),
                 Inches(5.0), Inches(0.3),
                 "Consumer Accessibility \u2192", font_size=12, color=BODY_GRAY,
                 alignment=PP_ALIGN.CENTER, bold=True)

    add_text_box(slide, matrix_x - Inches(1.4), matrix_y + Inches(1.5),
                 Inches(1.2), Inches(1.5),
                 "Clinical\nDepth\n\u2191", font_size=12, color=BODY_GRAY,
                 alignment=PP_ALIGN.CENTER, bold=True, line_spacing=1.2)

    # Competitor dots — UPDATED: bigger Lucid, dead competitor labels, Woebot shut down
    competitors = [
        ("Headspace", matrix_x + Inches(4.8), matrix_y + Inches(3.5), MEDIUM_GRAY, 0.18, ""),
        ("Calm", matrix_x + Inches(5.5), matrix_y + Inches(3.8), MEDIUM_GRAY, 0.18, ""),
        ("Woebot (shut down)", matrix_x + Inches(3.8), matrix_y + Inches(2.8), ACCENT_CORAL, 0.18, ""),
        ("Kintsugi (shut down)", matrix_x + Inches(1.2), matrix_y + Inches(0.6), ACCENT_CORAL, 0.18, ""),
        ("Sonde Health", matrix_x + Inches(0.8), matrix_y + Inches(1.2), MEDIUM_GRAY, 0.18, ""),
    ]

    for name, cx, cy, color, dot_sz, note in competitors:
        dot_size = Inches(dot_sz)
        dot = slide.shapes.add_shape(MSO_SHAPE.OVAL, cx, cy, dot_size, dot_size)
        dot.fill.solid()
        dot.fill.fore_color.rgb = color
        dot.line.fill.background()

        label_color = ACCENT_CORAL if "shut down" in name else BODY_GRAY
        add_text_box(slide, cx + dot_size + Inches(0.08), cy - Inches(0.02),
                     Inches(2.0), Inches(0.3),
                     name, font_size=11, color=label_color)

    # LUCID — much bigger dot with highlight ring
    lucid_cx = matrix_x + Inches(5.0)
    lucid_cy = matrix_y + Inches(0.3)
    lucid_dot_size = Inches(0.4)

    # Highlight ring (glow effect)
    ring_size = Inches(0.6)
    ring = slide.shapes.add_shape(MSO_SHAPE.OVAL,
                                   lucid_cx - Inches(0.1), lucid_cy - Inches(0.1),
                                   ring_size, ring_size)
    ring.fill.solid()
    ring.fill.fore_color.rgb = RGBColor(0xD8, 0xE8, 0xF4)
    ring.line.fill.background()

    # Lucid dot
    lucid_dot = slide.shapes.add_shape(MSO_SHAPE.OVAL, lucid_cx, lucid_cy,
                                         lucid_dot_size, lucid_dot_size)
    lucid_dot.fill.solid()
    lucid_dot.fill.fore_color.rgb = STEEL
    lucid_dot.line.fill.background()

    add_text_box(slide, lucid_cx + lucid_dot_size + Inches(0.1), lucid_cy - Inches(0.02),
                 Inches(1.5), Inches(0.35),
                 "Lucid", font_size=14, color=CHARCOAL, bold=True)

    # RIGHT side — legend with updated descriptions
    legend_x = Inches(9.2)
    legend_items = [
        ("Kintsugi Health", "Clinical-grade voice AI.\nShut down Feb 2026.", ACCENT_CORAL),
        ("Woebot", "AI chatbot therapy.\nShut down consumer app Jun 2025.", ACCENT_CORAL),
        ("Sonde Health", "Voice biomarkers for enterprises.\nNo consumer product.", MEDIUM_GRAY),
        ("Headspace / Calm", "Consumer meditation apps.\nNo clinical validation.", MEDIUM_GRAY),
    ]

    for i, (name, desc, color) in enumerate(legend_items):
        y = Inches(2.4 + i * 0.95)
        add_accent_dot(slide, legend_x, y + Inches(0.05), Inches(0.1), color)
        add_text_box(slide, legend_x + Inches(0.2), y, Inches(3.5), Inches(0.25),
                     name, font_size=12, color=CHARCOAL, bold=True)
        add_text_box(slide, legend_x + Inches(0.2), y + Inches(0.25), Inches(3.5), Inches(0.5),
                     desc, font_size=10, color=BODY_GRAY, line_spacing=1.2)

    # Bold callout box at bottom right
    callout_x = Inches(9.0)
    callout_y = Inches(6.3)
    callout_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                          callout_x, callout_y, Inches(3.8), Inches(0.65))
    callout_box.fill.solid()
    callout_box.fill.fore_color.rgb = NAVY
    callout_box.line.fill.background()
    callout_box.adjustments[0] = 0.1

    add_text_box(slide, callout_x + Inches(0.15), callout_y + Inches(0.08), Inches(3.5), Inches(0.5),
                 "Only player with clinical validation AND consumer distribution",
                 font_size=11, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER,
                 line_spacing=1.2)

    add_slide_number(slide, 13)


def slide_14_competitive_moats(prs):
    """Slide 14 — Competitive Moats (LIGHT background). All circles STEEL blue."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, LIGHT_BG)

    add_text_box(slide, Inches(0.8), Inches(0.5), Inches(3.0), Inches(0.35),
                 "DEFENSIBILITY", font_size=11, color=STEEL, bold=True)
    add_divider_line(slide, Inches(0.8), Inches(0.9), Inches(2.0), STEEL, Pt(2))

    add_text_box(slide, Inches(0.8), Inches(1.1), Inches(11.5), Inches(0.8),
                 "Five Moats That Get Deeper Over Time",
                 font_size=32, color=CHARCOAL, bold=True, font_name=HEADLINE_FONT)

    # UPDATED: All circles STEEL blue, Kintsugi reference updated in moat #5
    moats = [
        ("1", "Clinical Validation", "14,898-participant peer-reviewed study. FDA Breakthrough Device Designation. No competitor has this depth of validation."),
        ("2", "Proprietary Model", "244M-parameter DAM trained on 863 hours of clinical voice data. $30M+ of R&D invested. Can't be replicated cheaply."),
        ("3", "On-Device Processing", "100% local inference. No cloud dependency. Strongest privacy story in the category. Regulatory advantage for health data."),
        ("4", "Longitudinal Data", "Every user generates a unique mental health time series. More data = better personalization = higher retention = more data."),
        ("5", "Consumer Distribution", "Kintsugi and Woebot are dead. Sonde is B2B-only. Lucid owns the consumer relationship. Direct-to-consumer = brand loyalty + higher LTV."),
    ]

    for i, (num, title, body) in enumerate(moats):
        y = Inches(2.1 + i * 1.0)

        # Number circle — ALL STEEL BLUE
        circle = slide.shapes.add_shape(MSO_SHAPE.OVAL,
                                        Inches(0.8), y + Inches(0.05),
                                        Inches(0.4), Inches(0.4))
        circle.fill.solid()
        circle.fill.fore_color.rgb = STEEL
        circle.line.fill.background()

        add_text_box(slide, Inches(0.8), y + Inches(0.07),
                     Inches(0.4), Inches(0.4),
                     num, font_size=16, color=WHITE, bold=True,
                     alignment=PP_ALIGN.CENTER, font_name=HEADLINE_FONT)

        add_text_box(slide, Inches(1.4), y, Inches(3.0), Inches(0.35),
                     title, font_size=16, color=CHARCOAL, bold=True)

        add_text_box(slide, Inches(1.4), y + Inches(0.35), Inches(10.5), Inches(0.55),
                     body, font_size=12, color=BODY_GRAY, line_spacing=1.3)

    add_slide_number(slide, 14)


def slide_15_comparable_valuations(prs):
    """Slide 15 — Comparable Valuations (LIGHT background). Updated Kintsugi."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, LIGHT_BG)

    add_text_box(slide, Inches(0.8), Inches(0.5), Inches(3.0), Inches(0.35),
                 "COMPARABLES", font_size=11, color=STEEL, bold=True)
    add_divider_line(slide, Inches(0.8), Inches(0.9), Inches(2.0), STEEL, Pt(2))

    add_text_box(slide, Inches(0.8), Inches(1.1), Inches(11.5), Inches(0.8),
                 "Where Lucid Fits in the Landscape",
                 font_size=32, color=CHARCOAL, bold=True, font_name=HEADLINE_FONT)

    # Bar chart — UPDATED: Kintsugi label
    chart_data = CategoryChartData()
    chart_data.categories = ['Oura\n($11B)', 'Whoop\n($3.6B)', 'Calm\n($2.0B)',
                             'Headspace\n($3.0B)', 'Kintsugi\n($72M, shut down)']
    chart_data.add_series('Valuation ($M)', (11000, 3600, 2000, 3000, 72))

    chart_left = Inches(0.8)
    chart_top = Inches(2.0)
    chart_w = Inches(7.5)
    chart_h = Inches(4.0)

    chart_frame = slide.shapes.add_chart(
        XL_CHART_TYPE.COLUMN_CLUSTERED, chart_left, chart_top, chart_w, chart_h,
        chart_data
    )
    chart = chart_frame.chart
    chart.has_legend = False

    # Style the chart
    plot = chart.plots[0]
    plot.gap_width = 100

    series = chart.series[0]
    series.format.fill.solid()
    series.format.fill.fore_color.rgb = STEEL

    # Category axis
    cat_axis = chart.category_axis
    cat_axis.tick_labels.font.size = Pt(10)
    cat_axis.tick_labels.font.color.rgb = BODY_GRAY
    cat_axis.has_major_gridlines = False
    cat_axis.major_tick_mark = 2  # XL_TICK_MARK.NONE

    # Value axis
    val_axis = chart.value_axis
    val_axis.has_major_gridlines = True
    val_axis.major_gridlines.format.line.color.rgb = RGBColor(0xE0, 0xE0, 0xE8)
    val_axis.tick_labels.font.size = Pt(10)
    val_axis.tick_labels.font.color.rgb = BODY_GRAY

    # Data labels
    plot.has_data_labels = True
    data_labels = plot.data_labels
    data_labels.font.size = Pt(10)
    data_labels.font.color.rgb = CHARCOAL
    data_labels.font.bold = True
    data_labels.number_format = '$#,##0"M"'
    data_labels.label_position = XL_LABEL_POSITION.OUTSIDE_END

    # RIGHT side — comparison notes — UPDATED Kintsugi
    notes_x = Inches(8.8)

    comparisons = [
        {"text": "Oura ($11B)", "size": 14, "bold": True, "color": CHARCOAL, "font": HEADLINE_FONT, "space_after": 1},
        {"text": "Consumer wearable ring. Sleep + recovery. 5.5M+ users. Premium subscription model.", "size": 11, "color": BODY_GRAY, "space_after": 12},
        {"text": "Whoop ($3.6B)", "size": 14, "bold": True, "color": CHARCOAL, "font": HEADLINE_FONT, "space_after": 1},
        {"text": "Performance wearable. HRV + strain + recovery. Strong QS brand loyalty.", "size": 11, "color": BODY_GRAY, "space_after": 12},
        {"text": "Kintsugi ($72M before shutdown)", "size": 14, "bold": True, "color": CHARCOAL, "font": HEADLINE_FONT, "space_after": 1},
        {"text": "Voice biomarker AI. B2B only. Same underlying tech \u2014 killed by FDA regulatory economics, not bad science.", "size": 11, "color": BODY_GRAY, "space_after": 12},
    ]

    add_multiline_box(slide, notes_x, Inches(2.2), Inches(4.0), Inches(4.0),
                      comparisons, line_spacing=1.3)

    # Bottom callout — UPDATED
    add_divider_line(slide, Inches(0.8), Inches(6.3), Inches(11.7), RGBColor(0xE8, 0xE8, 0xF0), Pt(1))

    add_text_box(slide, Inches(0.8), Inches(6.45), Inches(11.5), Inches(0.5),
                 "Lucid = clinical rigor of Kintsugi (RIP) + consumer accessibility of Oura",
                 font_size=16, color=CHARCOAL, bold=True)

    add_slide_number(slide, 15)


def slide_16_traction(prs):
    """Slide 16 — Traction & Milestones (LIGHT background)."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, LIGHT_BG)

    add_text_box(slide, Inches(0.8), Inches(0.5), Inches(3.0), Inches(0.35),
                 "TRACTION", font_size=11, color=STEEL, bold=True)
    add_divider_line(slide, Inches(0.8), Inches(0.9), Inches(2.0), STEEL, Pt(2))

    add_text_box(slide, Inches(0.8), Inches(1.1), Inches(11.5), Inches(0.8),
                 "Milestones Completed & Next Steps",
                 font_size=32, color=CHARCOAL, bold=True, font_name=HEADLINE_FONT)

    # LEFT — Completed
    left_x = Inches(0.8)
    col_w = Inches(5.3)

    add_text_box(slide, left_x, Inches(2.1), col_w, Inches(0.4),
                 "Completed", font_size=20, color=ACCENT_GREEN, bold=True,
                 font_name=HEADLINE_FONT)

    completed = [
        "Clinical-grade voice biomarker model integrated",
        "On-device inference pipeline (zero cloud dependency)",
        "macOS desktop app \u2014 fully functional",
        "Wellness scoring with 6 clinical dimensions",
        "Speaker verification & voice activity detection",
        "Longitudinal trend tracking dashboard",
        "Privacy-first architecture (all data local)",
    ]

    for i, item in enumerate(completed):
        y = Inches(2.6 + i * 0.48)
        add_text_box(slide, left_x, y, Inches(0.3), Inches(0.3),
                     "\u2713", font_size=14, color=ACCENT_GREEN, bold=True)
        add_text_box(slide, left_x + Inches(0.35), y, col_w - Inches(0.35), Inches(0.4),
                     item, font_size=13, color=BODY_GRAY)

    # RIGHT — Next milestones
    right_x = Inches(7.0)

    add_text_box(slide, right_x, Inches(2.1), col_w, Inches(0.4),
                 "Next", font_size=20, color=STEEL, bold=True,
                 font_name=HEADLINE_FONT)

    next_items = [
        "iOS app launch (Q2 2026)",
        "Beta testing with 500 quantified-self users",
        "App Store launch (Q3 2026)",
        "B2B pilot with 2 employer wellness programs",
        "Android app (Q4 2026)",
        "1,000 paying subscribers target",
        "Series A readiness (Q1 2027)",
    ]

    for i, item in enumerate(next_items):
        y = Inches(2.6 + i * 0.48)
        circle = slide.shapes.add_shape(MSO_SHAPE.OVAL,
                                        right_x + Inches(0.03), y + Inches(0.05),
                                        Inches(0.15), Inches(0.15))
        circle.fill.background()
        circle.line.color.rgb = STEEL
        circle.line.width = Pt(1.5)

        add_text_box(slide, right_x + Inches(0.35), y, col_w - Inches(0.35), Inches(0.4),
                     item, font_size=13, color=BODY_GRAY)

    add_slide_number(slide, 16)


def slide_17_growth_projections(prs):
    """Slide 17 — Growth Projections (LIGHT background). NEW in v3."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, LIGHT_BG)

    add_text_box(slide, Inches(0.8), Inches(0.5), Inches(3.0), Inches(0.35),
                 "GROWTH", font_size=11, color=STEEL, bold=True)
    add_divider_line(slide, Inches(0.8), Inches(0.9), Inches(2.0), STEEL, Pt(2))

    add_text_box(slide, Inches(0.8), Inches(1.1), Inches(11.5), Inches(0.8),
                 "Growth Projections: Benchmarked Against Category Leaders",
                 font_size=32, color=CHARCOAL, bold=True, font_name=HEADLINE_FONT)

    # LEFT side — Growth milestone timeline
    left_x = Inches(0.8)
    timeline_y = Inches(2.3)

    add_text_box(slide, left_x, timeline_y, Inches(5.5), Inches(0.35),
                 "Lucid Growth Trajectory", font_size=18, color=CHARCOAL,
                 bold=True, font_name=HEADLINE_FONT)

    milestones = [
        ("Month 3", "2,000\u20135,000 users", "QS community launch", STEEL),
        ("Month 6", "15,000\u201325,000 users", "Product Hunt + influencer partnerships", STEEL),
        ("Month 12", "40K\u201375K users \u2022 3K\u20135K paid", "$540K\u2013$900K ARR", ACCENT_GREEN),
        ("Month 18", "80K\u2013150K users \u2022 8K\u201312K paid", "$1.4M\u2013$2.2M ARR", ACCENT_GREEN),
        ("Month 24", "200K+ users \u2022 20K+ paid", "$3.6M+ ARR", ACCENT_GOLD),
    ]

    for i, (period, users, detail, accent) in enumerate(milestones):
        y = timeline_y + Inches(0.55 + i * 0.72)

        # Timeline dot
        dot = slide.shapes.add_shape(MSO_SHAPE.OVAL,
                                      left_x + Inches(0.05), y + Inches(0.08),
                                      Inches(0.18), Inches(0.18))
        dot.fill.solid()
        dot.fill.fore_color.rgb = accent
        dot.line.fill.background()

        # Vertical connector line (except last)
        if i < len(milestones) - 1:
            line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                           left_x + Inches(0.12), y + Inches(0.28),
                                           Pt(2), Inches(0.52))
            line.fill.solid()
            line.fill.fore_color.rgb = RGBColor(0xD0, 0xD8, 0xE0)
            line.line.fill.background()

        # Period
        add_text_box(slide, left_x + Inches(0.35), y, Inches(1.2), Inches(0.3),
                     period, font_size=13, color=accent, bold=True)

        # Users
        add_text_box(slide, left_x + Inches(1.6), y, Inches(3.5), Inches(0.3),
                     users, font_size=13, color=CHARCOAL, bold=True)

        # Detail
        add_text_box(slide, left_x + Inches(1.6), y + Inches(0.28), Inches(3.5), Inches(0.3),
                     detail, font_size=11, color=BODY_GRAY)

    # RIGHT side — Benchmark comparisons
    right_x = Inches(6.8)
    bench_y = Inches(2.3)

    add_text_box(slide, right_x, bench_y, Inches(5.8), Inches(0.35),
                 "Category Benchmarks", font_size=18, color=CHARCOAL,
                 bold=True, font_name=HEADLINE_FONT)

    benchmarks = [
        ("Oura Ring", "4 years to 100K users, then 10x'd in 3 years.\n$11B valuation.", RGBColor(0x7B, 0x68, 0xD4)),
        ("WHOOP", "6 years building with athletes, then consumer breakout.\n$3.6B valuation.", ACCENT_GREEN),
        ("Health App Retention", "Annual plan 12-month retention: 44\u201360%.\nHealth & fitness = #1 category for payer LTV.", STEEL),
    ]

    for i, (name, desc, accent) in enumerate(benchmarks):
        y = bench_y + Inches(0.55 + i * 1.3)

        # Card
        box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                      right_x, y, Inches(5.8), Inches(1.1))
        box.fill.solid()
        box.fill.fore_color.rgb = CARD_BG
        box.line.color.rgb = CARD_BORDER
        box.line.width = Pt(1)
        box.adjustments[0] = 0.08

        add_divider_line(slide, right_x + Inches(0.2), y + Inches(0.15), Inches(0.5), accent, Pt(3))

        add_text_box(slide, right_x + Inches(0.2), y + Inches(0.3), Inches(5.2), Inches(0.3),
                     name, font_size=14, color=CHARCOAL, bold=True)
        add_text_box(slide, right_x + Inches(0.2), y + Inches(0.58), Inches(5.2), Inches(0.45),
                     desc, font_size=11, color=BODY_GRAY, line_spacing=1.3)

    # Bottom callout — unit economics
    add_divider_line(slide, Inches(0.8), Inches(6.3), Inches(11.7), RGBColor(0xE8, 0xE8, 0xF0), Pt(1))

    add_text_box(slide, Inches(0.8), Inches(6.45), Inches(11.5), Inches(0.5),
                 "At $14.99/mo with 44% annual retention, LTV = $80+. CAC in QS community: $15\u201325.",
                 font_size=16, color=CHARCOAL, bold=True)

    add_slide_number(slide, 17)


def slide_18_go_to_market(prs):
    """Slide 18 — Go-to-Market (LIGHT background)."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, LIGHT_BG)

    add_text_box(slide, Inches(0.8), Inches(0.5), Inches(3.0), Inches(0.35),
                 "GO-TO-MARKET", font_size=11, color=STEEL, bold=True)
    add_divider_line(slide, Inches(0.8), Inches(0.9), Inches(2.0), STEEL, Pt(2))

    add_text_box(slide, Inches(0.8), Inches(1.1), Inches(11.5), Inches(0.8),
                 "Land with Consumers. Expand into Enterprise.",
                 font_size=32, color=CHARCOAL, bold=True, font_name=HEADLINE_FONT)

    # Three phase cards
    phases = [
        ("Phase 1: Seed", "Q2-Q3 2026",
         "Launch on iOS & macOS App Stores\n"
         "Target: quantified-self communities\n"
         "Channels: Reddit, Twitter/X, Product Hunt\n"
         "Partnership: Whoop/Oura influencers\n"
         "Goal: 5,000 users, 500 paid",
         STEEL),
        ("Phase 2: Grow", "Q4 2026 - Q1 2027",
         "Expand to Android\n"
         "Launch referral program (1 month free)\n"
         "Content marketing: voice biomarker education\n"
         "First 2 B2B employer wellness pilots\n"
         "Goal: 25,000 users, 3,000 paid",
         ACCENT_GREEN),
        ("Phase 3: Scale", "2027+",
         "Insurance / payer partnerships\n"
         "White-label for telehealth platforms\n"
         "International expansion (UK, EU, Japan)\n"
         "Platform API for third-party integrations\n"
         "Goal: 100,000+ users, ARR $5M+",
         ACCENT_GOLD),
    ]

    card_w = Inches(3.6)
    card_h = Inches(4.2)
    start_x = Inches(0.8)
    card_y = Inches(2.1)
    gap = Inches(0.4)

    for i, (title, timeline, body, accent) in enumerate(phases):
        x = start_x + i * (card_w + gap)

        # Card
        box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, card_y, card_w, card_h)
        box.fill.solid()
        box.fill.fore_color.rgb = CARD_BG
        box.line.color.rgb = CARD_BORDER
        box.line.width = Pt(1)
        box.adjustments[0] = 0.05

        # Accent bar at top
        add_divider_line(slide, x + Inches(0.3), card_y + Inches(0.2), Inches(0.8), accent, Pt(3))

        # Phase title
        add_text_box(slide, x + Inches(0.3), card_y + Inches(0.4), card_w - Inches(0.6), Inches(0.4),
                     title, font_size=17, color=CHARCOAL, bold=True, font_name=HEADLINE_FONT)

        # Timeline
        add_text_box(slide, x + Inches(0.3), card_y + Inches(0.8), card_w - Inches(0.6), Inches(0.3),
                     timeline, font_size=12, color=accent, bold=True)

        # Body
        add_text_box(slide, x + Inches(0.3), card_y + Inches(1.2), card_w - Inches(0.6), Inches(2.8),
                     body, font_size=12, color=BODY_GRAY, line_spacing=1.5)

    add_slide_number(slide, 18)


def slide_19_closing(prs):
    """Slide 19 — Closing (NAVY background — bookend)."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, NAVY)

    add_divider_line(slide, Inches(4.5), Inches(1.8), Inches(4.3), STEEL, Pt(2))

    if os.path.exists(IMG_ICON):
        slide.shapes.add_picture(IMG_ICON, Inches(5.9), Inches(2.1), Inches(1.5), Inches(1.5))

    add_text_box(slide, Inches(2.0), Inches(3.8), Inches(9.3), Inches(0.8),
                 "Clinical-Grade Mental Health Monitoring.",
                 font_size=30, color=WHITE, bold=True,
                 alignment=PP_ALIGN.CENTER, font_name=HEADLINE_FONT)

    add_text_box(slide, Inches(2.0), Inches(4.6), Inches(9.3), Inches(0.5),
                 "No Clinic Required.", font_size=26, color=STEEL,
                 alignment=PP_ALIGN.CENTER, font_name=HEADLINE_FONT)

    add_divider_line(slide, Inches(5.5), Inches(5.4), Inches(2.3), RGBColor(0x33, 0x55, 0x6E), Pt(1))

    add_text_box(slide, Inches(2.0), Inches(5.7), Inches(9.3), Inches(0.4),
                 "hello@lucidvoice.com", font_size=16,
                 color=RGBColor(0xAA, 0xBB, 0xCC), alignment=PP_ALIGN.CENTER)

    add_text_box(slide, Inches(2.0), Inches(6.1), Inches(9.3), Inches(0.4),
                 "lucidvoice.com", font_size=14,
                 color=RGBColor(0x88, 0x99, 0xAA), alignment=PP_ALIGN.CENTER)

    add_slide_number(slide, 19, color=RGBColor(0x66, 0x77, 0x88))


# ── Main ────────────────────────────────────────────────────────

def main():
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    slide_01_title(prs)
    slide_02_problem_blind_spot(prs)
    slide_03_qs_market(prs)                # NEW
    slide_04_measurement_gap(prs)
    slide_05_options_failure(prs)           # NEW
    slide_06_solution(prs)
    slide_07_how_it_works(prs)
    slide_08_clinical_validation(prs)
    slide_09_technology(prs)
    slide_10_product_deep_dive(prs)
    slide_11_market(prs)
    slide_12_business_model(prs)
    slide_13_competitive_landscape(prs)
    slide_14_competitive_moats(prs)
    slide_15_comparable_valuations(prs)
    slide_16_traction(prs)
    slide_17_growth_projections(prs)        # NEW
    slide_18_go_to_market(prs)
    slide_19_closing(prs)

    prs.save(OUTPUT)
    print(f"Saved {TOTAL_SLIDES}-slide deck to: {OUTPUT}")


if __name__ == "__main__":
    main()
