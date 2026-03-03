#!/usr/bin/env python3
"""
Add 6 appendix slides (20-25) to the Lucid Investment Deck.
Updates all slide numbers from /19 to /25.

Usage:
    python3 add_appendix_slides.py
"""

import os
import re
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from lxml import etree

# ── Design System (matching generate_deck.py) ─────────────────
LIGHT_BG = RGBColor(0xF8, 0xF9, 0xFA)
STEEL = RGBColor(0x5B, 0x8D, 0xB8)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
CHARCOAL = RGBColor(0x2D, 0x2D, 0x3A)
BODY_GRAY = RGBColor(0x55, 0x55, 0x66)
MEDIUM_GRAY = RGBColor(0x88, 0x88, 0x99)
CARD_BG = RGBColor(0xFF, 0xFF, 0xFF)
CARD_BORDER = RGBColor(0xE0, 0xE0, 0xE8)

ACCENT_CORAL = RGBColor(0xE8, 0x6B, 0x6B)
ACCENT_GREEN = RGBColor(0x4E, 0xC9, 0xB0)
ACCENT_PURPLE = RGBColor(0x7B, 0x68, 0xD4)
ACCENT_GOLD = RGBColor(0xD4, 0xA5, 0x37)

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)
NEW_TOTAL = 25

HEADLINE_FONT = "Georgia"
BODY_FONT = "Helvetica Neue"

BASE = os.path.dirname(os.path.abspath(__file__))
INPUT = os.path.join(BASE, "Lucid Investment Deck.pptx")
OUTPUT = INPUT


# ── Helpers (same patterns as generate_deck.py) ───────────────

def add_bg(slide, color):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Emu(0), Emu(0), SLIDE_W, SLIDE_H)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    sp = shape._element
    sp.getparent().remove(sp)
    slide.shapes._spTree.insert(2, sp)


def add_text_box(slide, left, top, width, height, text, font_size=18,
                 color=CHARCOAL, bold=False, alignment=PP_ALIGN.LEFT,
                 font_name=BODY_FONT, line_spacing=1.2, anchor=MSO_ANCHOR.TOP):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    txBox.text_frame.word_wrap = True
    txBox.text_frame.auto_size = None
    p = txBox.text_frame.paragraphs[0]
    p.alignment = alignment
    p.space_after = Pt(0)
    p.space_before = Pt(0)
    pPr = p._p.get_or_add_pPr()
    lnSpc = etree.SubElement(pPr, '{http://schemas.openxmlformats.org/drawingml/2006/main}lnSpc')
    spcPct = etree.SubElement(lnSpc, '{http://schemas.openxmlformats.org/drawingml/2006/main}spcPct')
    spcPct.set('val', str(int(line_spacing * 100000)))
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.color.rgb = color
    run.font.bold = bold
    run.font.name = font_name
    txBox.text_frame._txBody.bodyPr.set('anchor', {
        MSO_ANCHOR.TOP: 't', MSO_ANCHOR.MIDDLE: 'ctr', MSO_ANCHOR.BOTTOM: 'b'
    }.get(anchor, 't'))
    return txBox


def add_multiline_box(slide, left, top, width, height, lines, default_size=13,
                      default_color=BODY_GRAY, default_bold=False, alignment=PP_ALIGN.LEFT,
                      default_font=BODY_FONT, line_spacing=1.2):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    txBox.text_frame.word_wrap = True
    txBox.text_frame.auto_size = None
    for i, line in enumerate(lines):
        if isinstance(line, str):
            line = {"text": line}
        if i == 0:
            p = txBox.text_frame.paragraphs[0]
        else:
            p = txBox.text_frame.add_paragraph()
        p.alignment = line.get("alignment", alignment)
        p.space_after = Pt(line.get("space_after", 4))
        p.space_before = Pt(line.get("space_before", 0))
        pPr = p._p.get_or_add_pPr()
        lnSpc = etree.SubElement(pPr, '{http://schemas.openxmlformats.org/drawingml/2006/main}lnSpc')
        spcPct = etree.SubElement(lnSpc, '{http://schemas.openxmlformats.org/drawingml/2006/main}spcPct')
        spcPct.set('val', str(int(line.get("line_spacing", line_spacing) * 100000)))
        run = p.add_run()
        run.text = line.get("text", "")
        run.font.size = Pt(line.get("size", default_size))
        run.font.color.rgb = line.get("color", default_color)
        run.font.bold = line.get("bold", default_bold)
        run.font.name = line.get("font", default_font)
        if line.get("italic"):
            run.font.italic = True
    return txBox


def add_slide_number(slide, num):
    add_text_box(slide, Inches(11.8), Inches(6.9), Inches(1.2), Inches(0.4),
                 f"{num} / {NEW_TOTAL}", font_size=10, color=MEDIUM_GRAY,
                 alignment=PP_ALIGN.RIGHT)


def add_section_label(slide, text):
    add_text_box(slide, Inches(0.8), Inches(0.5), Inches(5), Inches(0.3),
                 text, font_size=11, color=STEEL, bold=True)


def add_title(slide, text, top=Inches(0.85)):
    add_text_box(slide, Inches(0.8), top, Inches(11), Inches(0.7),
                 text, font_size=32, color=CHARCOAL, bold=True, font_name=HEADLINE_FONT)


def add_stat_block(slide, left, top, number, label, num_color=STEEL):
    """Add a stat number + label below it, no background card."""
    add_text_box(slide, left, top, Inches(3.5), Inches(0.5),
                 number, font_size=18, color=num_color, bold=True, font_name=HEADLINE_FONT)
    add_text_box(slide, left, top + Inches(0.4), Inches(3.5), Inches(0.4),
                 label, font_size=11, color=BODY_GRAY)


def add_card_box(slide, left, top, width, height):
    """Add a white rounded card background."""
    box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    box.fill.solid()
    box.fill.fore_color.rgb = CARD_BG
    box.line.color.rgb = CARD_BORDER
    box.line.width = Pt(1)
    box.adjustments[0] = 0.04
    return box


def add_divider_line(slide, left, top, width, color=STEEL, thickness=Pt(1.5)):
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, thickness)
    line.fill.solid()
    line.fill.fore_color.rgb = color
    line.line.fill.background()


def set_notes(slide, notes_text):
    if not slide.has_notes_slide:
        slide.notes_slide
    tf = slide.notes_slide.notes_text_frame
    tf.text = notes_text.strip()


# ── Speaker Notes Content ─────────────────────────────────────

NOTES_20 = """=== APPENDIX: PAIN QUANTIFICATION — The Economic Cost ===

THE ECONOMIC CRATER OF MENTAL HEALTH (Section 1A):
- US annual economic cost of mental illness: $282B/year, 1.7% of GDP (Yale/Columbia/UW-Madison, April 2024)
- Global productivity loss from depression + anxiety: $1T/year, 12B working days lost (WHO "Mental Health at Work" factsheet)
- US employer cost from untreated mental health: $105B/year (Center for Prevention & Health Services)
- Presenteeism cost per employee: $12,000-$19,875/year — 7.5x the cost of absenteeism (CDC/Global Corporate Challenge)
- Depressed employee productivity reduction: 35% (American Psychiatric Association)
- Cost per missed workday: $340/day (Kaiser Permanente)
- Projected global cost by 2030: $6T/year (World Economic Forum, 2024)

The takeaway: Mental illness costs more than cancer, diabetes, and respiratory disease combined in lost productivity. The economic case for better monitoring is enormous.

THE PHYSICAL VS MENTAL HEALTH DATA GAP (Section 1D):
- Rich physical data: Wearables provide continuous, objective data on HR, HRV, sleep, SpO2, skin temp, activity
- Impoverished mental data: "There is no available standard biomarker for the detection of mental health conditions" (PMC survey on wearable sensors, 2023)
- QS user frustration: Users are "virtually drowning in data" but lacking tools for mental/emotional insight (Choe et al. 2014, CHI)
- The gap Lucid fills: Voice is the only passive, hardware-free biomarker with direct neurological links to mental state

TECH WORKER BURNOUT — TARGET DEMOGRAPHIC (Section 1E):
- 82% of tech employees feeling close to burnout (Talkspace Business)
- 52% of tech workers with depression/anxiety (Talkspace Business)
- 65% of engineers experienced burnout past year (Interview Guys 2025)
- 45% of tech founders rate mental health "bad/very bad" (CEREVITY 2025)
- >50% of Millennials/Gen Z have been in therapy; 39% plan to go (Thriving Center of Psych)

Our early adopter demographic — tech-forward, QS-oriented professionals — is also the demographic most affected by burnout and mental health challenges."""

NOTES_21 = """=== APPENDIX: PAIN QUANTIFICATION — Spending & Failures ===

WHAT QS USERS ALREADY SPEND ON WELLNESS (Section 1B):
- Average American wellness spending: $5,321/year (Fortune Well / Global Wellness Institute, Feb 2024)
- US wellness market: $480B, growing ~10%/year (McKinsey 2024)
- Gen Z + Millennials: 41% of wellness spend despite being 36% of adult population (McKinsey Future of Wellness 2024)
- Oura average ticket: >$400 (hardware + subscription) with 52% YoY spending growth (Earnest Analytics 2024)
- WHOOP annual cost: $199-$399/year (WHOOP pricing page)
- Willingness to pay for health apps: $6.50/month median, 58.9% willing to pay (Liu, Xie & Or 2024, SAGE Digital Health)
- 40% of consumers willing to pay specifically for mental health services (research2guidance)
- H&F payer LTV is highest of all app categories: Median $16.44, upper quartile $31.12 (RevenueCat 2025)

$14.99/month for the only mental health data stream is well within their established spending envelope.

WHAT THEY'VE TRIED FOR MENTAL HEALTH — AND FAILED (Section 1C):
- Therapy: $100-$250/session (avg ~$140); 58% cite cost as #1 barrier; 34.8% dropout rate; only 4.3% of primary care visits include screening (Swift & Greenberg 2012; Samples et al. 2020)
- Mental health apps (Calm $596M revenue, Headspace $348M): 3.3% median 30-day retention; requires active engagement (Baumel et al. 2019)
- Telehealth therapy: $134/session without insurance; same engagement problems as in-person; scheduling friction (GoodRx 2024)
- Corporate EAPs: $12-$40/employee/year, $68.4B market; only 4-6% utilization rate (EASNA; Fortune Business Insights)
- Psychiatric medication: $5-$489/month; does not provide the "data" QS users crave; no tracking loop (K Health pricing data)
- AI chatbots (Woebot): $114M raised; consumer product shut down June 2025 (MobiHealthNews)

Every existing solution requires active effort — scheduling, showing up, typing, meditating. QS users don't want another chore; they want another data stream. Lucid is the only passive mental health monitoring tool."""

NOTES_22 = """=== APPENDIX: BEHAVIORAL SCIENCE — Engagement Architecture ===

THE ENGAGEMENT PROBLEM MENTAL HEALTH APPS NEVER SOLVED (Section 2A):
- Mental health apps: 3.3% 30-day retention (Baumel et al. 2019)
- Monitoring apps (Oura, Apple Watch): 88-89%+ 12-month retention
- Root cause: Mental health apps require active effort (meditation, journaling, CBT exercises). Monitoring apps are passive and create compounding data value.
- Key insight: Lucid is a MONITORING TOOL (like Oura), not a TREATMENT APP (like Calm). This is the single most important architectural decision we made.

BEHAVIORAL SCIENCE FRAMEWORKS (Section 2B):
1. Nir Eyal's Hook Model: Trigger (morning notification) -> Action (25-sec recording) -> Variable Reward (Canopy Score reveal with particle animation) -> Investment (historical data accumulates)
2. BJ Fogg's Behavior Model: Behavior = Motivation x Ability x Prompt. High motivation (health data) x ultra-low ability barrier (25 seconds) x system prompt (morning notification)
3. Loss Aversion (Kahneman & Tversky): Losses feel 2x as painful as equivalent gains. Grove feature: missed days create wilted trees; users feel compelled to maintain streaks.
4. Endowed Progress (Nunes & Dreze 2006): Pre-giving progress nearly doubles completion (19% -> 34%). First 2 waypoints auto-complete on signup; Canopy Score starts building immediately.
5. Variable Ratio Reinforcement (Skinner): Unpredictable rewards are most engaging. Echoes (pattern discovery) appear at unpredictable intervals; Sanctuary celebrations have varied triggers.
6. Zeigarnik Effect: Incomplete tasks are remembered 1.9x more than completed ones. Rhythm Rings show partial progress; Grove trees in mid-growth stages.

WHY PREVIOUS SOLUTIONS LACK THESE MECHANICS (Section 2D):
- Therapy: No data, no score, no streak, no progress visualization — session-based, not continuous
- Calm/Headspace: Same experience every time (no variable reward); no personal data; 3.3% 30-day retention
- BetterHelp/Talkspace: Scheduling friction; therapist dependency; no passive monitoring
- Woebot: Chatbot = active engagement required; no biometric data; shut down June 2025
- EAPs: Stigma; no data; no habit loop; requires self-referral; 4-6% utilization despite being free

WHY QS USERS ARE UNIQUELY SUSCEPTIBLE (Section 2E):
- QS identity = "self-knowledge through numbers" — intrinsically motivated by data
- Already habituated to daily health rituals (checking Oura score, closing Apple Watch rings)
- High conscientiousness and emotional stability personality correlates
- 2-3x higher willingness to pay for health subscriptions vs general population
- Social proof is powerful in QS communities (r/quantifiedself 155K, r/Biohackers 350K+)"""

NOTES_23 = """=== APPENDIX: BEHAVIORAL SCIENCE — 10 Features Mapped ===

LUCID'S 10 ENGAGEMENT FEATURES MAPPED TO DOPAMINE LOOPS (Section 2C):
1. Canopy Score (0-100 daily wellness score): Progressive disclosure + count-up animation + leaf particles. Comparable: Oura Readiness Score. Dopamine: Anticipation -> Relief -> Baseline comparison.
2. Grove (streak forest with wilting): Loss aversion + collection + recovery ritual. Comparable: Duolingo streaks (8x retention for streak users). Dopamine: Status + loss pain + recovery satisfaction.
3. Sanctuary Overlay (micro-celebrations): Variable reward + sensory delight; 5-min cooldown prevents fatigue. Comparable: Instagram like notification. Dopamine: Immediate positive feedback at milestone moments.
4. Rhythm Rings (3 daily goal rings): Progress visualization + adaptive goals (+5%/week). Comparable: Apple Watch Activity Rings. Dopamine: Completion -> Celebration -> Next-day reset.
5. Waypoints (30-tier achievement trail): Endowed progress + tiered unlocking across 6 stages over 90 days. Comparable: Strava/Xbox achievements. Dopamine: Progression gates + long-term goal orientation.
6. Morning Briefing: Ritual creation + narrative framing ("Yesterday you..."). Comparable: Oura morning readiness ritual. Dopamine: Curiosity -> Validation -> Motivation.
7. Weekly Wrapped: Retrospective validation + Spotify Wrapped-style social design. Comparable: Spotify Wrapped (200M users engaged in 24 hours). Dopamine: Self-narrative + shareable format.
8. Echoes (pattern discovery): Variable reward + discovery delight; appear after 7+ days. Comparable: Strava segment discoveries. Dopamine: Curiosity -> Insight -> Secondary waypoint unlock.
9. Notifications ("The Pulse"): Contextual alerts + positive reinforcement; rate-limited to 4/hour. Comparable: Apple Watch stand reminders. Dopamine: Awareness -> Personalized comparison -> Action.
10. Onboarding (6-step flow): Micro-commitments + endowed progress + baseline calibration. Comparable: Duolingo onboarding (placement test -> immediate lesson). Dopamine: Low friction entry -> Early wins -> Commitment.

KEY ENGAGEMENT STATS FROM COMPARABLE APPS (Section 2F):
- Duolingo streaks: Streak users retain at 8x (40% vs 5% at Day 30) — Duolingo S-1/earnings
- Apple Watch ring closing: Users check wrists ~80 times/day; 89%+ 12-month retention — Apple investor materials
- Oura morning readiness: Users open app 3+ times/day — Oura press materials
- Spotify Wrapped: 200M users engaged in 24 hours; 500M social shares — Spotify year-end reports
- Strava Segments + Kudos: 120M users; social features drive 2x retention vs solo — Strava press releases"""

NOTES_24 = """=== APPENDIX: FINANCIAL MODEL — Unit Economics & Margins ===

PRICING & REVENUE MODEL (Section 3A):
- Consumer Monthly: $14.99/month (V10 website pricing)
- Consumer Annual: $149/year (~$12.42/mo, 17% discount vs monthly)
- Blended Monthly ARPU: ~$12.42 (assuming 60% annual / 40% monthly mix)
- Annual ARPU: ~$149
- Free Tier: $0, basic check-ins, conversion funnel
- Corporate (Phase 2): $2-$5 PEPM (per-employee-per-month)

GROSS MARGIN ANALYSIS (Section 3C):
App Store Path (Year 1 & 2+):
- Revenue per user/month: $14.99
- Payment processing: -$2.25 (15% small biz rate)
- Server/infra costs: ~$0 (on-device processing)
- Support allocation: -$0.50
- COGS total: $2.75 | Gross Margin: 81.7%

Direct Web/Stripe Path:
- Payment processing: -$0.73 (Stripe 2.9% + $0.30)
- COGS total: $1.23 | Gross Margin: 91.8%

Key advantage: On-device processing = zero compute COGS. Margin profile closer to pure software than AI company.
Benchmarks: Typical SaaS 75-85% (Benchmarkit 2024), Top-quartile software 85%+, Oura ~60% (HW blend), Peloton ~40% (HW drag), Lucid 82-92%.

UNIT ECONOMICS — THREE SCENARIOS (Section 3D):
Conservative: Blended ARPU $10/mo, 35% 12-mo retention, 8-mo avg lifespan, LTV $80, CAC $80, LTV/CAC 1.0x, 8-mo payback, 82% margin
Target: ARPU $12.42/mo, 44% 12-mo retention (RevenueCat median), 12-mo lifespan, LTV $149, CAC $60, LTV/CAC 2.5x, 5-mo payback, 85% margin
Optimistic: ARPU $14.99/mo, 60% 12-mo retention (top quartile), 18-mo lifespan, LTV $270, CAC $40, LTV/CAC 6.8x, 3-mo payback, 92% margin

Path to 3:1+: Organic/referral channels must comprise 40%+ of acquisition. Consistent with WHOOP and Oura growth patterns — both grew primarily through organic/community channels before scaling paid."""

NOTES_25 = """=== APPENDIX: FINANCIAL MODEL — Channels, Breakeven & Comps ===

CUSTOMER ACQUISITION COST BY CHANNEL (Section 3B):
1. Apple Search Ads: CPI $3-$8, 9.4% install-to-paid, CAC $32-$85 [HIGH confidence] (Watsspace 2025; SplitMetrics 2025)
2. Meta/Facebook: CPM $15.77, CPC $1.10, CPI $10.42, CAC $111 [HIGH] (Mesha 2025; SuperAds 2025)
3. Instagram: CPC $1.83-$3.35, CPI ~$3.50, CAC $37 [MEDIUM] (Quimby Digital 2025; Birch 2025)
4. Google Search: CPC $4.22 (mental health, +42% YoY), CAC $160-$265 [MEDIUM] (WordStream 2025)
5. TikTok: CPC $0.40-$1.00, CPM $4-$7, CAC $21-$53 [LOW] (Quimby Digital 2025)
6. YouTube: CPV $0.071 (healthcare), CAC $160-$300 [LOW] (Awisee 2025; AdBacklog 2025)
7. Reddit: CPC $0.10-$0.80, CPM $2-$6, CAC $32-$85 [MEDIUM] (AdBacklog 2025)
8. Podcast Sponsorships: CPM $25-$50 (host-read), CAC $160-$320 [LOW] (Ad Results Media 2025)
9. LinkedIn: CPC $5.58-$10.00, CAC $200-$500 for B2B leads [MEDIUM] (The B2B House 2025)
10. Organic/SEO: $1,500-$3K/mo budget, CAC $15-$50 at scale after 12+ months [MEDIUM] (SearchAtlas 2025)

Blended CAC Strategy: Weight toward Apple Search Ads + Reddit + Organic (lowest CAC). Target blended CAC of $50-$80.

BREAKEVEN ANALYSIS (Section 3E):
- Lean (3 FTE, $25K/mo burn): 2,034 subs needed, $364K ARR
- Seed (5-7 FTE, $50K/mo burn): 4,068 subs needed, $729K ARR
- Growth (10 FTE, $100K/mo burn): 7,843 subs needed, $1.4M ARR
Formula: Subs needed = Monthly burn / (ARPU x Gross margin) = $25K / ($14.99 x 0.82) = 2,034
Key insight: Software-only model = breakeven at MUCH lower subscriber count than hardware companies. Oura needed millions of ring sales; Lucid needs ~2,000-4,000 subscribers.

COMPARABLE COMPANY BENCHMARKS (Section 3F):
- Oura: ~$500M rev (2024), 2M paying, ~$55 ARPU*, $11B valuation, 10.4x revenue multiple, ~60% GM (HW blend), expanding profitability (Sacra; CNBC Oct 2025)
- WHOOP: $260M+ rev (2025), ~1M est. subs, ~$260 ARPU, $3.7B valuation, ~14x multiple (GetLatka; Sacra)
- Calm: $596M rev (2024), 4M+ subs, ~$75 ARPU, $2B valuation, 3.4x multiple, profitable since 2018 (GetLatka; CNBC)
- Headspace: $348M rev (2024), 2.8M subs (declining), ~$124 ARPU, $3B valuation, 8.6x multiple (GetLatka; Business of Apps)
- Noom: ~$1B rev (2023), 1.5M subs, ~$420 ARPU (Sacra)
*Oura ARPU appears low because hardware revenue is separate; subscription is $5.99/mo = $72/yr

Revenue multiple implications: At $1M ARR -> 10x = $10M; At $5M ARR -> $50M; At $50M ARR -> $500M.

DATA QUALITY DISCLOSURE (Section 3G):
- HIGH confidence: Apple/Google commission rates, Stripe fees, RevenueCat churn/retention/ARPU, SaaS gross margin benchmarks, Calm revenue trajectory
- MEDIUM confidence: Meta/Google/Reddit CPM/CPC, startup burn rates, LTV/CAC thresholds
- DERIVED: CAC per paid subscriber (CPI / conversion rate), Lucid gross margin projections, breakeven subscriber counts
- NOT FOUND (flagged): TikTok health-specific CPI, YouTube health conversion rate, podcast-to-install rate, WHOOP exact subscriber count"""


# ── Slide Builders ─────────────────────────────────────────────

def slide_20_economic_cost(prs):
    """Slide 20 — The Economic Cost of Ignoring Mental Health"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, LIGHT_BG)
    add_section_label(slide, "APPENDIX: PAIN QUANTIFICATION")
    add_title(slide, "The Economic Cost of Ignoring Mental Health")

    # Divider under title
    add_divider_line(slide, Inches(0.8), Inches(1.55), Inches(4), STEEL)

    # 6 stat blocks in 3x2 grid
    col_starts = [Inches(0.8), Inches(4.8), Inches(8.8)]
    row_starts = [Inches(2.0), Inches(4.4)]

    stats = [
        ("$282B/yr", "US economic cost of mental illness\n(Yale/Columbia/UW-Madison 2024)", ACCENT_CORAL),
        ("$1T/yr", "Global productivity loss\n12B working days (WHO)", ACCENT_CORAL),
        ("35%", "Productivity reduction in\ndepressed employees (APA)", STEEL),
        ("$6T/yr", "Projected global cost\nby 2030 (WEF)", ACCENT_CORAL),
        ("4.3%", "Primary care visits with\nmental health screening", ACCENT_PURPLE),
        ("82%", "Tech employees feeling\nclose to burnout (Talkspace)", ACCENT_PURPLE),
    ]

    for idx, (number, label, color) in enumerate(stats):
        col = idx % 3
        row = idx // 3
        left = col_starts[col]
        top = row_starts[row]

        # Card background
        add_card_box(slide, left, top, Inches(3.5), Inches(1.9))
        # Stat number
        add_text_box(slide, left + Inches(0.3), top + Inches(0.25), Inches(2.9), Inches(0.6),
                     number, font_size=28, color=color, bold=True, font_name=HEADLINE_FONT)
        # Label
        add_text_box(slide, left + Inches(0.3), top + Inches(0.85), Inches(2.9), Inches(0.8),
                     label, font_size=11, color=BODY_GRAY, line_spacing=1.3)

    add_slide_number(slide, 20)
    set_notes(slide, NOTES_20)


def slide_21_spending_failures(prs):
    """Slide 21 — What QS Users Spend — And What's Failed"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, LIGHT_BG)
    add_section_label(slide, "APPENDIX: PAIN QUANTIFICATION")
    add_title(slide, "What QS Users Spend \u2014 And What\u2019s Failed")

    add_divider_line(slide, Inches(0.8), Inches(1.55), Inches(4), STEEL)

    # Left column — What They Spend
    left_x = Inches(0.8)
    add_card_box(slide, left_x, Inches(1.9), Inches(5.6), Inches(5.0))

    add_text_box(slide, left_x + Inches(0.3), Inches(2.1), Inches(5.0), Inches(0.4),
                 "What They Spend", font_size=16, color=ACCENT_GREEN, bold=True, font_name=HEADLINE_FONT)

    spend_items = [
        {"text": "$5,321/yr", "size": 14, "bold": True, "color": ACCENT_GREEN, "font": HEADLINE_FONT, "space_after": 1},
        {"text": "Average American wellness spending", "size": 11, "color": BODY_GRAY, "space_after": 10},
        {"text": ">$400", "size": 14, "bold": True, "color": ACCENT_GREEN, "font": HEADLINE_FONT, "space_after": 1},
        {"text": "Oura average ticket (52% YoY growth)", "size": 11, "color": BODY_GRAY, "space_after": 10},
        {"text": "$199\u2013$399/yr", "size": 14, "bold": True, "color": ACCENT_GREEN, "font": HEADLINE_FONT, "space_after": 1},
        {"text": "WHOOP annual cost", "size": 11, "color": BODY_GRAY, "space_after": 10},
        {"text": "58.9%", "size": 14, "bold": True, "color": ACCENT_GREEN, "font": HEADLINE_FONT, "space_after": 1},
        {"text": "Willing to pay for health apps", "size": 11, "color": BODY_GRAY, "space_after": 10},
        {"text": "$14.99/mo", "size": 14, "bold": True, "color": STEEL, "font": HEADLINE_FONT, "space_after": 1},
        {"text": "Lucid\u2019s price fits the envelope", "size": 11, "color": BODY_GRAY},
    ]
    add_multiline_box(slide, left_x + Inches(0.3), Inches(2.55), Inches(5.0), Inches(4.0), spend_items)

    # Right column — What's Failed
    right_x = Inches(6.9)
    add_card_box(slide, right_x, Inches(1.9), Inches(5.6), Inches(5.0))

    add_text_box(slide, right_x + Inches(0.3), Inches(2.1), Inches(5.0), Inches(0.4),
                 "What\u2019s Failed", font_size=16, color=ACCENT_CORAL, bold=True, font_name=HEADLINE_FONT)

    fail_items = [
        {"text": "Therapy", "size": 14, "bold": True, "color": ACCENT_CORAL, "font": HEADLINE_FONT, "space_after": 1},
        {"text": "34.8% dropout, $140/session avg", "size": 11, "color": BODY_GRAY, "space_after": 10},
        {"text": "Mental Health Apps", "size": 14, "bold": True, "color": ACCENT_CORAL, "font": HEADLINE_FONT, "space_after": 1},
        {"text": "3.3% median 30-day retention", "size": 11, "color": BODY_GRAY, "space_after": 10},
        {"text": "Telehealth", "size": 14, "bold": True, "color": ACCENT_CORAL, "font": HEADLINE_FONT, "space_after": 1},
        {"text": "$134/session, same friction", "size": 11, "color": BODY_GRAY, "space_after": 10},
        {"text": "EAPs", "size": 14, "bold": True, "color": ACCENT_CORAL, "font": HEADLINE_FONT, "space_after": 1},
        {"text": "4\u20136% utilization despite being free", "size": 11, "color": BODY_GRAY, "space_after": 10},
        {"text": "Woebot", "size": 14, "bold": True, "color": ACCENT_CORAL, "font": HEADLINE_FONT, "space_after": 1},
        {"text": "$114M raised \u2192 shut down 2025", "size": 11, "color": BODY_GRAY},
    ]
    add_multiline_box(slide, right_x + Inches(0.3), Inches(2.55), Inches(5.0), Inches(4.0), fail_items)

    add_slide_number(slide, 21)
    set_notes(slide, NOTES_21)


def slide_22_behavioral_science(prs):
    """Slide 22 — Engagement Architecture: Behavioral Science"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, LIGHT_BG)
    add_section_label(slide, "APPENDIX: BEHAVIORAL SCIENCE")
    add_title(slide, "Engagement Architecture: Behavioral Science")

    add_divider_line(slide, Inches(0.8), Inches(1.55), Inches(4), STEEL)

    # Key insight callout box
    callout_box = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(1.85), Inches(11.7), Inches(0.9))
    callout_box.fill.solid()
    callout_box.fill.fore_color.rgb = RGBColor(0xEE, 0xF4, 0xFA)  # light steel-blue
    callout_box.line.color.rgb = STEEL
    callout_box.line.width = Pt(1)
    callout_box.adjustments[0] = 0.04

    add_text_box(slide, Inches(1.1), Inches(1.95), Inches(11.2), Inches(0.7),
                 "Mental health apps: 3.3% retention at 30 days. Monitoring apps: 88%+ at 12 months. "
                 "Lucid is a monitoring tool, not a treatment app.",
                 font_size=13, color=CHARCOAL, bold=True, font_name=BODY_FONT, line_spacing=1.4)

    # 6 frameworks in 2x3 grid of cards
    frameworks = [
        ("Hook Model", "Eyal", "Notification \u2192 Record \u2192\nScore reveal \u2192 Data accumulates", STEEL),
        ("Behavior Model", "Fogg", "High motivation \u00d7 25-sec\nbarrier \u00d7 System prompt", ACCENT_GREEN),
        ("Loss Aversion", "Kahneman", "Grove wilting on missed\ndays (losses feel 2\u00d7)", ACCENT_CORAL),
        ("Endowed Progress", "Nunes & Dreze", "First 2 waypoints auto-\ncomplete (19% \u2192 34%)", ACCENT_PURPLE),
        ("Variable Reinforcement", "Skinner", "Echoes appear\nunpredictably", ACCENT_GOLD),
        ("Zeigarnik Effect", "", "Rhythm Rings show\npartial progress (1.9\u00d7 recall)", STEEL),
    ]

    col_starts = [Inches(0.8), Inches(4.8), Inches(8.8)]
    row_starts = [Inches(3.05), Inches(5.15)]

    for idx, (name, author, desc, color) in enumerate(frameworks):
        col = idx % 3
        row = idx // 3
        left = col_starts[col]
        top = row_starts[row]

        add_card_box(slide, left, top, Inches(3.5), Inches(1.8))

        # Framework name
        label = f"{name}" if not author else f"{name} ({author})"
        add_text_box(slide, left + Inches(0.25), top + Inches(0.15), Inches(3.0), Inches(0.45),
                     label, font_size=13, color=color, bold=True, font_name=BODY_FONT)
        # Description
        add_text_box(slide, left + Inches(0.25), top + Inches(0.65), Inches(3.0), Inches(0.9),
                     desc, font_size=11, color=BODY_GRAY, line_spacing=1.3)

    add_slide_number(slide, 22)
    set_notes(slide, NOTES_22)


def slide_23_ten_features(prs):
    """Slide 23 — 10 Features x Engagement Loops"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, LIGHT_BG)
    add_section_label(slide, "APPENDIX: BEHAVIORAL SCIENCE")
    add_title(slide, "10 Features \u00d7 Engagement Loops")

    add_divider_line(slide, Inches(0.8), Inches(1.55), Inches(4), STEEL)

    # Two-column layout for 10 features (5 per column)
    features_left = [
        ("1. Canopy Score (0\u2013100)", "cf. Oura Readiness Score", STEEL),
        ("2. Grove (streak forest)", "cf. Duolingo streaks (8\u00d7 retention)", ACCENT_GREEN),
        ("3. Sanctuary (micro-celebrations)", "Variable reward + cooldown", ACCENT_GOLD),
        ("4. Rhythm Rings (3 goals)", "cf. Apple Watch Activity Rings", ACCENT_CORAL),
        ("5. Waypoints (30-tier trail)", "Endowed progress over 90 days", ACCENT_PURPLE),
    ]

    features_right = [
        ("6. Morning Briefing", "cf. Oura morning ritual", STEEL),
        ("7. Weekly Wrapped", "cf. Spotify Wrapped (200M users)", ACCENT_GREEN),
        ("8. Echoes (pattern discovery)", "Variable reward after 7+ days", ACCENT_GOLD),
        ("9. The Pulse (notifications)", "Rate-limited contextual alerts", ACCENT_CORAL),
        ("10. Onboarding (6-step)", "Micro-commitments + calibration", ACCENT_PURPLE),
    ]

    # Left column card
    add_card_box(slide, Inches(0.8), Inches(1.9), Inches(5.6), Inches(4.2))
    y = Inches(2.05)
    for name, desc, color in features_left:
        add_text_box(slide, Inches(1.1), y, Inches(5.0), Inches(0.35),
                     name, font_size=13, color=color, bold=True, font_name=BODY_FONT)
        add_text_box(slide, Inches(1.1), y + Inches(0.3), Inches(5.0), Inches(0.3),
                     desc, font_size=11, color=BODY_GRAY)
        y += Inches(0.78)

    # Right column card
    add_card_box(slide, Inches(6.9), Inches(1.9), Inches(5.6), Inches(4.2))
    y = Inches(2.05)
    for name, desc, color in features_right:
        add_text_box(slide, Inches(7.2), y, Inches(5.0), Inches(0.35),
                     name, font_size=13, color=color, bold=True, font_name=BODY_FONT)
        add_text_box(slide, Inches(7.2), y + Inches(0.3), Inches(5.0), Inches(0.3),
                     desc, font_size=11, color=BODY_GRAY)
        y += Inches(0.78)

    # Bottom engagement stats bar
    add_card_box(slide, Inches(0.8), Inches(6.3), Inches(11.7), Inches(0.7))
    stats_text = ("Duolingo: 8\u00d7 retention with streaks  |  Apple Watch: 89%+ at 12 mo  |  "
                  "Oura: 3\u00d7/day opens  |  Spotify Wrapped: 200M in 24hrs")
    add_text_box(slide, Inches(1.1), Inches(6.35), Inches(11.2), Inches(0.5),
                 stats_text, font_size=11, color=STEEL, bold=True,
                 alignment=PP_ALIGN.CENTER)

    add_slide_number(slide, 23)
    set_notes(slide, NOTES_23)


def slide_24_unit_economics(prs):
    """Slide 24 — Unit Economics & Gross Margins"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, LIGHT_BG)
    add_section_label(slide, "APPENDIX: FINANCIAL MODEL")
    add_title(slide, "Unit Economics & Gross Margins")

    add_divider_line(slide, Inches(0.8), Inches(1.55), Inches(4), STEEL)

    # Pricing line
    add_text_box(slide, Inches(0.8), Inches(1.8), Inches(11.7), Inches(0.4),
                 "$14.99/mo  |  $149/yr  |  Blended ARPU ~$12.42",
                 font_size=14, color=CHARCOAL, bold=True, font_name=HEADLINE_FONT,
                 alignment=PP_ALIGN.CENTER)

    # Gross margin comparison — left card
    add_card_box(slide, Inches(0.8), Inches(2.4), Inches(5.6), Inches(2.5))

    add_text_box(slide, Inches(1.1), Inches(2.55), Inches(5.0), Inches(0.35),
                 "Gross Margin Comparison", font_size=15, color=CHARCOAL, bold=True, font_name=HEADLINE_FONT)

    margin_lines = [
        {"text": "Lucid (App Store): 81.7%", "size": 13, "bold": True, "color": ACCENT_GREEN, "space_after": 4},
        {"text": "Lucid (Direct): 91.8%", "size": 13, "bold": True, "color": ACCENT_GREEN, "space_after": 10},
        {"text": "vs Oura ~60%  |  Peloton ~40%  |  Typical SaaS 75\u201385%", "size": 11, "color": BODY_GRAY, "space_after": 10},
        {"text": "Key: Zero compute COGS \u2014 all on-device", "size": 11, "color": STEEL, "bold": True},
    ]
    add_multiline_box(slide, Inches(1.1), Inches(2.95), Inches(5.0), Inches(1.7), margin_lines)

    # Unit economics — right card
    add_card_box(slide, Inches(6.9), Inches(2.4), Inches(5.6), Inches(2.5))

    add_text_box(slide, Inches(7.2), Inches(2.55), Inches(5.0), Inches(0.35),
                 "Three-Scenario Unit Economics", font_size=15, color=CHARCOAL, bold=True, font_name=HEADLINE_FONT)

    # Conservative
    econ_lines = [
        {"text": "Conservative", "size": 12, "bold": True, "color": BODY_GRAY, "space_after": 1},
        {"text": "LTV $80  |  CAC $80  |  LTV/CAC 1.0\u00d7", "size": 12, "color": BODY_GRAY, "space_after": 8},
        {"text": "Target", "size": 12, "bold": True, "color": STEEL, "space_after": 1},
        {"text": "LTV $149  |  CAC $60  |  LTV/CAC 2.5\u00d7", "size": 12, "color": STEEL, "space_after": 8},
        {"text": "Optimistic", "size": 12, "bold": True, "color": ACCENT_GREEN, "space_after": 1},
        {"text": "LTV $270  |  CAC $40  |  LTV/CAC 6.8\u00d7", "size": 12, "color": ACCENT_GREEN},
    ]
    add_multiline_box(slide, Inches(7.2), Inches(2.95), Inches(5.0), Inches(1.7), econ_lines)

    # Bottom insight
    add_card_box(slide, Inches(0.8), Inches(5.2), Inches(11.7), Inches(1.5))
    add_text_box(slide, Inches(1.1), Inches(5.35), Inches(11.2), Inches(0.35),
                 "Path to 3:1+ LTV/CAC", font_size=15, color=CHARCOAL, bold=True, font_name=HEADLINE_FONT)
    add_text_box(slide, Inches(1.1), Inches(5.75), Inches(11.2), Inches(0.7),
                 "Organic/referral channels must comprise 40%+ of acquisition. Consistent with WHOOP and Oura "
                 "growth patterns \u2014 both grew primarily through organic/community channels before scaling paid. "
                 "Conservative case (1.0\u00d7) models worst-case paid-only; Target (2.5\u00d7) is viable with organic supplement.",
                 font_size=12, color=BODY_GRAY, line_spacing=1.4)

    add_slide_number(slide, 24)
    set_notes(slide, NOTES_24)


def slide_25_channels_comps(prs):
    """Slide 25 — Acquisition Channels, Breakeven & Comps"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, LIGHT_BG)
    add_section_label(slide, "APPENDIX: FINANCIAL MODEL")
    add_title(slide, "Acquisition Channels, Breakeven & Comps")

    add_divider_line(slide, Inches(0.8), Inches(1.55), Inches(4), STEEL)

    # Top channels — left card
    add_card_box(slide, Inches(0.8), Inches(1.85), Inches(5.6), Inches(2.2))
    add_text_box(slide, Inches(1.1), Inches(1.95), Inches(5.0), Inches(0.35),
                 "Top Acquisition Channels (CAC)", font_size=15, color=CHARCOAL, bold=True, font_name=HEADLINE_FONT)

    channel_lines = [
        {"text": "Apple Search Ads: $32\u2013$85", "size": 12, "color": ACCENT_GREEN, "bold": True, "space_after": 4},
        {"text": "Reddit: $32\u2013$85", "size": 12, "color": ACCENT_GREEN, "bold": True, "space_after": 4},
        {"text": "Organic/SEO: $15\u2013$50 (at scale)", "size": 12, "color": ACCENT_GREEN, "bold": True, "space_after": 4},
        {"text": "Instagram: $37", "size": 12, "color": STEEL, "bold": True, "space_after": 8},
        {"text": "Blended target: $50\u2013$80", "size": 12, "color": CHARCOAL, "bold": True},
    ]
    add_multiline_box(slide, Inches(1.1), Inches(2.35), Inches(5.0), Inches(1.5), channel_lines)

    # Breakeven — right card
    add_card_box(slide, Inches(6.9), Inches(1.85), Inches(5.6), Inches(2.2))
    add_text_box(slide, Inches(7.2), Inches(1.95), Inches(5.0), Inches(0.35),
                 "Breakeven Scenarios", font_size=15, color=CHARCOAL, bold=True, font_name=HEADLINE_FONT)

    break_lines = [
        {"text": "Lean (3 FTE)", "size": 12, "bold": True, "color": STEEL, "space_after": 1},
        {"text": "2,034 subs  |  $364K ARR", "size": 12, "color": BODY_GRAY, "space_after": 8},
        {"text": "Seed (5\u20137 FTE)", "size": 12, "bold": True, "color": STEEL, "space_after": 1},
        {"text": "4,068 subs  |  $729K ARR", "size": 12, "color": BODY_GRAY, "space_after": 8},
        {"text": "Growth (10 FTE)", "size": 12, "bold": True, "color": STEEL, "space_after": 1},
        {"text": "7,843 subs  |  $1.4M ARR", "size": 12, "color": BODY_GRAY},
    ]
    add_multiline_box(slide, Inches(7.2), Inches(2.35), Inches(5.0), Inches(1.5), break_lines)

    # Comparable companies — full-width bottom card
    add_card_box(slide, Inches(0.8), Inches(4.3), Inches(11.7), Inches(2.6))
    add_text_box(slide, Inches(1.1), Inches(4.4), Inches(11.2), Inches(0.35),
                 "Comparable Companies", font_size=15, color=CHARCOAL, bold=True, font_name=HEADLINE_FONT)

    # 4 comps in a row
    comps = [
        ("Oura", "$500M rev", "$11B", "10.4\u00d7"),
        ("WHOOP", "$260M rev", "$3.7B", "~14\u00d7"),
        ("Calm", "$596M rev", "$2B", "3.4\u00d7"),
        ("Headspace", "$348M rev", "$3B", "8.6\u00d7"),
    ]

    comp_x_starts = [Inches(1.1), Inches(3.95), Inches(6.8), Inches(9.65)]
    for i, (name, rev, val, mult) in enumerate(comps):
        x = comp_x_starts[i]
        add_text_box(slide, x, Inches(4.85), Inches(2.5), Inches(0.3),
                     name, font_size=14, color=STEEL, bold=True, font_name=HEADLINE_FONT)

        comp_detail_lines = [
            {"text": rev, "size": 12, "color": BODY_GRAY, "space_after": 2},
            {"text": f"Valuation: {val}", "size": 12, "color": BODY_GRAY, "space_after": 2},
            {"text": f"Rev Multiple: {mult}", "size": 12, "color": ACCENT_GREEN, "bold": True},
        ]
        add_multiline_box(slide, x, Inches(5.15), Inches(2.5), Inches(1.2), comp_detail_lines)

    add_slide_number(slide, 25)
    set_notes(slide, NOTES_25)


# ── Update Existing Slide Numbers ──────────────────────────────

def update_slide_numbers(prs):
    """Update all slide number text from 'N / 19' to 'N / 25' across all slides."""
    ns = '{http://schemas.openxmlformats.org/drawingml/2006/main}'
    count = 0
    for slide in prs.slides:
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            for para in shape.text_frame.paragraphs:
                full_text = para.text.strip()
                # Match patterns like "1 / 19", "19 / 19"
                if re.match(r'^\d{1,2}\s*/\s*19$', full_text):
                    for run in para.runs:
                        old = run.text
                        new = old.replace("/ 19", f"/ {NEW_TOTAL}")
                        if new != old:
                            run.text = new
                            count += 1
    print(f"  Updated {count} slide number(s) from /19 to /{NEW_TOTAL}")


# ── Main ───────────────────────────────────────────────────────

def main():
    if not os.path.exists(INPUT):
        print(f"ERROR: Input file not found: {INPUT}")
        return

    prs = Presentation(INPUT)
    slides_before = len(prs.slides)
    print(f"Loaded deck with {slides_before} slides")

    if slides_before != 19:
        print(f"WARNING: Expected 19 slides, found {slides_before}")

    # Add 6 appendix slides
    print("Adding appendix slides...")
    slide_20_economic_cost(prs)
    print("  Added slide 20: The Economic Cost of Ignoring Mental Health")
    slide_21_spending_failures(prs)
    print("  Added slide 21: What QS Users Spend — And What's Failed")
    slide_22_behavioral_science(prs)
    print("  Added slide 22: Engagement Architecture: Behavioral Science")
    slide_23_ten_features(prs)
    print("  Added slide 23: 10 Features × Engagement Loops")
    slide_24_unit_economics(prs)
    print("  Added slide 24: Unit Economics & Gross Margins")
    slide_25_channels_comps(prs)
    print("  Added slide 25: Acquisition Channels, Breakeven & Comps")

    # Update all slide numbers
    print("Updating slide numbers...")
    update_slide_numbers(prs)

    prs.save(OUTPUT)
    slides_after = len(prs.slides)
    print(f"\nSaved deck: {OUTPUT}")
    print(f"Total slides: {slides_after}")


if __name__ == "__main__":
    main()
